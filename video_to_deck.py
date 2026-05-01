#!/usr/bin/env python3
"""
video_to_deck.py — Turn a screen-recording MP4 + SRT transcript into a PowerPoint deck.

Pipeline:
  1. FFmpeg scene detection -> candidate frames + timestamps
  2. Perceptual-hash dedup -> drop near-identical frames
  3. Parse SRT transcript -> timestamped segments
  4. Align transcript chunks to each frame's time window
  5. Build .pptx with frame as visual + transcript as speaker notes

Usage:
    python video_to_deck.py input.mp4 transcript.srt -o output.pptx
    python video_to_deck.py input.mp4 transcript.srt --scene-threshold 0.25

Requirements:
    System:  ffmpeg, ffprobe
    Python:  pip install imagehash Pillow python-pptx
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple


# ---------- Step 1: Frame extraction (scene detection + periodic sampling) ----------

def extract_scene_frames(video: Path, out_dir: Path, threshold: float = 0.2) -> List[Tuple[Path, float]]:
    """Run ffmpeg scene detection. Returns list of (frame_path, timestamp_seconds)."""
    out_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        "ffmpeg", "-hide_banner", "-y",
        "-i", str(video),
        "-vf", f"select='eq(n,0)+gt(scene,{threshold})',showinfo",
        "-vsync", "vfr",
        "-q:v", "3",
        str(out_dir / "scene_%04d.jpg"),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg failed:\n{proc.stderr}")

    timestamps = [float(m) for m in re.findall(r"pts_time:([\d.]+)", proc.stderr)]
    frame_files = sorted(out_dir.glob("scene_*.jpg"))
    if len(frame_files) != len(timestamps):
        n = min(len(frame_files), len(timestamps))
        frame_files, timestamps = frame_files[:n], timestamps[:n]

    return list(zip(frame_files, timestamps))


def extract_interval_frames(video: Path, out_dir: Path, interval: float, duration: float) -> List[Tuple[Path, float]]:
    """Sample one frame every `interval` seconds. Returns (frame_path, timestamp) list."""
    out_dir.mkdir(parents=True, exist_ok=True)
    if interval <= 0:
        return []

    fps = 1.0 / interval
    cmd = [
        "ffmpeg", "-hide_banner", "-y",
        "-i", str(video),
        "-vf", f"fps={fps}",
        "-q:v", "3",
        str(out_dir / "interval_%04d.jpg"),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg interval sampling failed:\n{proc.stderr}")

    frame_files = sorted(out_dir.glob("interval_*.jpg"))
    # ffmpeg's fps filter places the first frame at t≈interval/2; reconstruct timestamps.
    timestamps = [min(duration, interval * (i + 0.5)) for i in range(len(frame_files))]
    return list(zip(frame_files, timestamps))


def merge_frames(*frame_lists: List[Tuple[Path, float]]) -> List[Tuple[Path, float]]:
    """Merge multiple frame lists, sort by timestamp."""
    combined: List[Tuple[Path, float]] = []
    for fl in frame_lists:
        combined.extend(fl)
    combined.sort(key=lambda x: x[1])
    return combined


# ---------- Step 2: Perceptual-hash dedup (consecutive only) ----------

def dedupe_frames(frames: List[Tuple[Path, float]], hash_distance: int = 4) -> List[Tuple[Path, float]]:
    """
    Drop a frame only if it's near-identical to the *previous kept* frame.
    Comparing only consecutively preserves recurring views (e.g., returning to an editor).
    """
    import imagehash
    from PIL import Image

    kept: List[Tuple[Path, float]] = []
    prev_hash = None
    for path, ts in frames:
        h = imagehash.phash(Image.open(path))
        if prev_hash is not None and (h - prev_hash) <= hash_distance:
            path.unlink(missing_ok=True)
            continue
        kept.append((path, ts))
        prev_hash = h
    return kept


# ---------- Step 3: SRT parser ----------

@dataclass
class TranscriptSegment:
    start: float
    end: float
    text: str


def _srt_time_to_seconds(t: str) -> float:
    """'00:01:23,456' -> 83.456"""
    h, m, rest = t.split(":")
    s, ms = rest.split(",")
    return int(h) * 3600 + int(m) * 60 + int(s) + int(ms) / 1000


def parse_srt(srt_path: Path) -> List[TranscriptSegment]:
    """Parse a .srt file. Tolerant of BOM, CRLF, dot-vs-comma ms separator, simple markup."""
    text = srt_path.read_text(encoding="utf-8-sig")
    blocks = re.split(r"\r?\n\r?\n", text.strip())

    segments: List[TranscriptSegment] = []
    time_re = re.compile(r"(\d{2}:\d{2}:\d{2}[,.]\d{3})\s*-->\s*(\d{2}:\d{2}:\d{2}[,.]\d{3})")

    for block in blocks:
        lines = [ln for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        time_line_idx = next((i for i, ln in enumerate(lines) if time_re.search(ln)), None)
        if time_line_idx is None:
            continue
        m = time_re.search(lines[time_line_idx])
        start = _srt_time_to_seconds(m.group(1).replace(".", ","))
        end = _srt_time_to_seconds(m.group(2).replace(".", ","))
        body = " ".join(lines[time_line_idx + 1:])
        body = re.sub(r"<[^>]+>", "", body)         # <i>, <b>, <font>
        body = re.sub(r"\{[^}]+\}", "", body).strip()  # {\an8}, etc.
        if body:
            segments.append(TranscriptSegment(start=start, end=end, text=body))

    return segments


# ---------- Step 4: Align transcript to frames ----------

def assign_transcript_to_frames(
    frames: List[Tuple[Path, float]],
    segments: List[TranscriptSegment],
    video_duration: float,
) -> List[Tuple[Path, float, float, str]]:
    """
    For each frame at time t_i, the slide covers [t_i, t_{i+1}).
    Collect transcript segments whose midpoint falls in that window.
    Returns: list of (frame_path, slide_start, slide_end, transcript_text).
    """
    results = []
    for i, (path, start) in enumerate(frames):
        end = frames[i + 1][1] if i + 1 < len(frames) else video_duration
        chunk = " ".join(
            seg.text for seg in segments
            if start <= (seg.start + seg.end) / 2 < end
        )
        results.append((path, start, end, chunk))
    return results


def get_video_duration(video: Path) -> float:
    out = subprocess.check_output([
        "ffprobe", "-v", "error", "-show_entries", "format=duration",
        "-of", "json", str(video),
    ])
    return float(json.loads(out)["format"]["duration"])


# ---------- Step 5: Build the .pptx ----------

def format_timestamp(seconds: float) -> str:
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    return f"{h:d}:{m:02d}:{s:02d}" if h else f"{m:d}:{s:02d}"


def _shorten(text: str, max_chars: int = 280) -> str:
    """Trim a transcript chunk to a slide-friendly length on a sentence boundary."""
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    # Prefer to end on sentence punctuation, then on a word boundary
    for sep in (". ", "! ", "? "):
        idx = cut.rfind(sep)
        if idx >= max_chars * 0.5:
            return cut[: idx + 1].strip()
    idx = cut.rfind(" ")
    return (cut[:idx] if idx > 0 else cut).strip() + "…"


def build_deck(slides_data, output: Path, video_name: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = video_name
    p.font.size = Pt(40)
    p.font.bold = True
    sub = title_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.5))
    sub.text_frame.paragraphs[0].text = f"{len(slides_data)} slides extracted from video"
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    # Layout: image on top (~5"), transcript snippet below (~1.6"), footer at bottom.
    img_top_in = 0.3
    img_max_h_in = 5.0
    img_max_w_in = 12.33
    text_top_in = 5.5
    text_h_in = 1.6
    footer_top_in = 7.15

    for idx, (frame_path, start, end, transcript) in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        with Image.open(frame_path) as im:
            iw, ih = im.size
        scale = min(img_max_w_in / (iw / 96), img_max_h_in / (ih / 96))
        disp_w = (iw / 96) * scale
        disp_h = (ih / 96) * scale
        left = (13.33 - disp_w) / 2

        slide.shapes.add_picture(
            str(frame_path),
            Inches(left), Inches(img_top_in),
            width=Inches(disp_w), height=Inches(disp_h),
        )

        # On-slide transcript snippet
        snippet = _shorten(transcript, 280) if transcript else ""
        if snippet:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(text_top_in),
                                          Inches(12.33), Inches(text_h_in))
            tf = tx.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = snippet
            para.font.size = Pt(14)

        # Footer with time range + slide number
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(footer_top_in), Inches(12.7), Inches(0.3))
        fp = footer.text_frame.paragraphs[0]
        fp.text = f"{format_timestamp(start)} – {format_timestamp(end)}   |   Slide {idx}"
        fp.font.size = Pt(10)

        if transcript:
            slide.notes_slide.notes_text_frame.text = transcript

    prs.save(str(output))


# ---------- Config / API key handling ----------

CONFIG_DIR = Path.home() / ".config" / "video-digest"
CONFIG_ENV_FILE = CONFIG_DIR / ".env"


def load_env_files() -> None:
    """Populate os.environ from .env files. Real env vars always win.

    Order (lowest priority first; later loads do NOT override earlier-set keys):
      1. Real env vars (from the shell)
      2. CWD/.env (project-local)
      3. ~/.config/video-digest/.env (global fallback)
    """
    from dotenv import load_dotenv

    load_dotenv()  # CWD/.env, only fills in missing
    if CONFIG_ENV_FILE.exists():
        load_dotenv(CONFIG_ENV_FILE)


def ensure_api_key() -> None:
    """Make sure ANTHROPIC_API_KEY is set, prompting + saving on first run if interactive."""
    if os.environ.get("ANTHROPIC_API_KEY"):
        return

    msg = (
        "ANTHROPIC_API_KEY is not set.\n"
        "Get a key from: https://console.anthropic.com/settings/keys"
    )
    if not sys.stdin.isatty():
        sys.exit(
            f"{msg}\n"
            "Then either export it (`export ANTHROPIC_API_KEY=...`), put it in a .env "
            "file in the current directory, or save it globally via:\n"
            "  mkdir -p ~/.config/video-digest && echo 'ANTHROPIC_API_KEY=sk-ant-...' "
            "> ~/.config/video-digest/.env"
        )

    print(msg)
    key = input("Paste your API key (or press Enter to abort): ").strip()
    if not key:
        sys.exit("Aborted.")

    save = input(
        f"Save it to {CONFIG_ENV_FILE} so future runs find it automatically? [Y/n] "
    ).strip().lower()
    if save in ("", "y", "yes"):
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)
        CONFIG_ENV_FILE.write_text(f"ANTHROPIC_API_KEY={key}\n")
        try:
            os.chmod(CONFIG_ENV_FILE, 0o600)
        except OSError:
            pass
        print(f"      saved to {CONFIG_ENV_FILE}")
    os.environ["ANTHROPIC_API_KEY"] = key


# ---------- YouTube fetch ----------

URL_RE = re.compile(r"^https?://", re.IGNORECASE)


def is_url(s: str) -> bool:
    return bool(URL_RE.match(s))


def fetch_youtube(url: str, cache_root: Path) -> Tuple[Path, Path]:
    """Download mp4 + English SRT from YouTube. Cached by video ID under cache_root.

    Returns (mp4_path, srt_path).
    """
    import yt_dlp

    cache_root.mkdir(parents=True, exist_ok=True)

    # Probe first to get the video ID for stable cache layout
    with yt_dlp.YoutubeDL({"quiet": True, "no_warnings": True}) as ydl:
        info = ydl.extract_info(url, download=False)
    video_id = info["id"]
    out_dir = cache_root / video_id
    out_dir.mkdir(parents=True, exist_ok=True)

    mp4_path = out_dir / f"{video_id}.mp4"
    srt_path = out_dir / f"{video_id}.en.srt"

    if mp4_path.exists() and srt_path.exists():
        print(f"      using cached {out_dir}/")
        return mp4_path, srt_path

    ydl_opts = {
        "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
        "merge_output_format": "mp4",
        "outtmpl": str(out_dir / f"{video_id}.%(ext)s"),
        "writesubtitles": True,
        "writeautomaticsub": True,
        "subtitleslangs": ["en", "en-US", "en-GB"],
        "subtitlesformat": "srt/vtt/best",
        "postprocessors": [{"key": "FFmpegSubtitlesConvertor", "format": "srt"}],
        "quiet": True,
        "no_warnings": True,
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])

    # Resolve actual filenames yt-dlp produced (lang suffix can vary)
    candidates_mp4 = list(out_dir.glob(f"{video_id}.*"))
    mp4_found = next((p for p in candidates_mp4 if p.suffix in (".mp4", ".mkv", ".webm")), None)
    if mp4_found and mp4_found != mp4_path:
        mp4_found.rename(mp4_path)
    elif not mp4_path.exists():
        raise RuntimeError(f"yt-dlp finished but no video file found in {out_dir}")

    srt_found = next((p for p in out_dir.glob(f"{video_id}.*.srt")), None)
    if srt_found and srt_found != srt_path:
        srt_found.rename(srt_path)
    elif not srt_path.exists():
        raise RuntimeError(
            f"No English subtitles available for {url} (neither manual nor auto-captions)."
        )

    return mp4_path, srt_path


# ---------- Markdown digest (transcript-primary, LLM-summarized) ----------

DIGEST_SYSTEM_PROMPT = (
    "You distill video transcripts into readable digests so a reader can grasp "
    "the essence of a video without watching it.\n\n"
    "Given a timestamped transcript, identify the natural topic segments. "
    "Return 5–12 sections depending on length and content density (favor fewer "
    "for short videos, more for long talks).\n\n"
    "For each topic:\n"
    "- title: a short, descriptive heading (not a question, not a teaser)\n"
    "- start_time: the timestamp in SECONDS where the topic begins, taken from "
    "the transcript's bracketed timestamps\n"
    "- summary: 2–4 sentences of informative prose distilling what is actually "
    "said. Concrete claims, names, numbers, and reasoning — not vague paraphrases.\n"
    "- key_points: 3–6 short bullet points capturing the most useful takeaways. "
    "Bullets should add detail beyond the summary, not restate it.\n\n"
    "Also produce a 2–3 sentence overview of the entire video.\n\n"
    "Write for a reader, not a viewer. Skip filler ('in this video I'll show you'); "
    "go straight to the substance."
)


def generate_digest(
    segments: List[TranscriptSegment],
    video_title: str,
    model: str,
):
    """Call Claude to segment the transcript into topics. Returns a parsed VideoDigest."""
    import anthropic
    from pydantic import BaseModel
    from typing import List as TList

    class Topic(BaseModel):
        title: str
        start_time: float
        summary: str
        key_points: TList[str]

    class VideoDigest(BaseModel):
        title: str
        overview: str
        topics: TList[Topic]

    transcript = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    user_text = (
        f"Video title: {video_title}\n\n"
        f"Total duration: {format_timestamp(segments[-1].end if segments else 0)}\n\n"
        f"Timestamped transcript:\n\n{transcript}"
    )

    client = anthropic.Anthropic()
    response = client.messages.parse(
        model=model,
        max_tokens=16000,
        system=DIGEST_SYSTEM_PROMPT,
        messages=[{
            "role": "user",
            "content": [{
                "type": "text",
                "text": user_text,
                "cache_control": {"type": "ephemeral"},
            }],
        }],
        output_format=VideoDigest,
    )
    return response.parsed_output, response.usage


def _candidates_for_topic(
    topic_start: float,
    topic_end: float,
    frames: List[Tuple[Path, float]],
    max_per_topic: int = 5,
    overlap_pre: float = 5.0,
) -> List[Tuple[Path, float]]:
    """Frames whose timestamp falls in [start - overlap, end). Downsample to max_per_topic."""
    in_window = [
        (p, t) for p, t in frames if (topic_start - overlap_pre) <= t < topic_end
    ]
    if len(in_window) <= max_per_topic:
        return in_window
    # Even spacing across the window
    step = len(in_window) / max_per_topic
    return [in_window[int(i * step)] for i in range(max_per_topic)]


def _encode_frame_for_vision(path: Path, max_long_edge: int = 1024) -> str:
    """Resize a frame to <= max_long_edge on the long side and return base64-encoded JPEG bytes."""
    import base64
    import io
    from PIL import Image

    with Image.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        scale = max_long_edge / max(w, h)
        if scale < 1:
            im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=82)
    return base64.standard_b64encode(buf.getvalue()).decode("ascii")


def vision_pick_frames(
    digest,
    frames: List[Tuple[Path, float]],
    video_duration: float,
    model: str,
):
    """Use Claude's vision to pick the best frame per topic from in-window candidates.

    Returns a dict {topic_index -> chosen_frame_path}, plus the API usage object.
    Topics with no in-window candidates are omitted (caller falls back to timestamp-based pick).
    """
    import anthropic
    from pydantic import BaseModel
    from typing import List as TList

    class TopicChoice(BaseModel):
        topic_index: int
        candidate_index: int
        rationale: str

    class FrameChoices(BaseModel):
        choices: TList[TopicChoice]

    topics = digest.topics

    # Build per-topic candidate lists and a flat list of (topic_idx, cand_idx, path, ts)
    per_topic: List[List[Tuple[Path, float]]] = []
    for i, topic in enumerate(topics):
        end = topics[i + 1].start_time if i + 1 < len(topics) else video_duration
        per_topic.append(_candidates_for_topic(topic.start_time, end, frames))

    # Build the message: text intro -> for each topic, label + summary + numbered candidate images
    content: list = []
    intro = (
        "For each topic below, pick the candidate frame that best illustrates what the "
        "narrator is discussing. Prefer frames showing the most informative visual content "
        "(diagrams, code, distinctive UI) over generic framing or talking-head shots. "
        "Return one choice per topic that has candidates.\n\n"
        f"Total topics: {len(topics)}\n"
    )
    content.append({"type": "text", "text": intro})

    for ti, topic in enumerate(topics):
        cands = per_topic[ti]
        if not cands:
            content.append({
                "type": "text",
                "text": f"\n--- Topic {ti} (no candidates available — skip) ---\n"
                        f"Title: {topic.title}\n",
            })
            continue
        header = (
            f"\n--- Topic {ti} ---\n"
            f"Title: {topic.title}\n"
            f"Summary: {topic.summary}\n"
            f"Candidates ({len(cands)} frames):\n"
        )
        content.append({"type": "text", "text": header})
        for ci, (path, ts) in enumerate(cands):
            content.append({
                "type": "text",
                "text": f"Candidate {ci} (at {format_timestamp(ts)}):",
            })
            content.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/jpeg",
                    "data": _encode_frame_for_vision(path),
                },
            })

    system = (
        "You select illustrative frames for a video digest. You will be shown a list of "
        "topics, each with a small set of candidate frames. For each topic with candidates, "
        "return the (topic_index, candidate_index) pair that best illustrates the topic, "
        "with a one-sentence rationale. Skip topics that say 'no candidates available'."
    )

    client = anthropic.Anthropic()
    response = client.messages.parse(
        model=model,
        max_tokens=4000,
        system=system,
        messages=[{"role": "user", "content": content}],
        output_format=FrameChoices,
    )

    chosen: dict = {}
    for choice in response.parsed_output.choices:
        if 0 <= choice.topic_index < len(topics):
            cands = per_topic[choice.topic_index]
            if 0 <= choice.candidate_index < len(cands):
                chosen[choice.topic_index] = cands[choice.candidate_index][0]
    return chosen, response.usage


def _pick_topic_frame(
    topic_start: float,
    topic_end: float,
    candidates: List[Tuple[Path, float]],
    used: set,
) -> Optional[Tuple[Path, float]]:
    """Pick the best frame for a topic: prefer a frame inside [start, end), else closest."""
    in_window = [(p, t) for p, t in candidates if topic_start <= t < topic_end and p not in used]
    if in_window:
        midpoint = (topic_start + min(topic_end, in_window[-1][1] + 1)) / 2
        return min(in_window, key=lambda x: abs(x[1] - midpoint))
    available = [(p, t) for p, t in candidates if p not in used]
    if not available:
        return None
    return min(available, key=lambda x: abs(x[1] - topic_start))


def write_markdown_digest(
    digest,
    candidate_frames: List[Tuple[Path, float]],
    video_duration: float,
    output_md: Path,
    images_dir: Path,
    vision_picks: Optional[dict] = None,
) -> None:
    """Render the digest as Markdown with <img> tags. Copies frames into images_dir.

    If vision_picks is provided, prefer those mappings; fall back to timestamp-based picks
    for any topic not covered.
    """
    images_dir.mkdir(parents=True, exist_ok=True)

    used: set = set()
    topic_images: List[Optional[Path]] = []
    topics = digest.topics
    for i, topic in enumerate(topics):
        # Vision pick takes priority
        if vision_picks and i in vision_picks:
            src_path = vision_picks[i]
            used.add(src_path)
        else:
            end = topics[i + 1].start_time if i + 1 < len(topics) else video_duration
            pick = _pick_topic_frame(topic.start_time, end, candidate_frames, used)
            if pick is None:
                topic_images.append(None)
                continue
            src_path, _ = pick
            used.add(src_path)
        dest = images_dir / f"topic_{i + 1:02d}.jpg"
        shutil.copy(src_path, dest)
        topic_images.append(dest)

    rel_dir = images_dir.name  # render as "images/topic_NN.jpg"
    lines: List[str] = []
    lines.append(f"# {digest.title}")
    lines.append("")
    lines.append(digest.overview)
    lines.append("")
    for i, topic in enumerate(topics):
        ts = format_timestamp(topic.start_time)
        lines.append(f"## {topic.title}  <sub>*{ts}*</sub>")
        lines.append("")
        if topic_images[i] is not None:
            lines.append(f'<img src="{rel_dir}/{topic_images[i].name}" width="800">')
            lines.append("")
        lines.append(topic.summary)
        lines.append("")
        for kp in topic.key_points:
            lines.append(f"- {kp}")
        if topic.key_points:
            lines.append("")

    output_md.write_text("\n".join(lines).rstrip() + "\n")


# ---------- Main ----------

def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("video", help="Input MP4 file path OR a YouTube URL (auto-downloads mp4 + SRT)")
    ap.add_argument("srt", nargs="?", default=None,
                    help="SRT transcript file (required if 'video' is a local path; "
                         "ignored for URLs — fetched automatically)")
    ap.add_argument("-o", "--output", type=Path, default=None, metavar="PATH",
                    help="Digest output path (default: <video-name>_digest.md)")
    ap.add_argument("--deck", nargs="?", const="__default__", default=None, metavar="PATH",
                    help="Also write a PowerPoint deck. With no path, defaults to "
                         "<video-name>_deck.pptx.")
    ap.add_argument("--deck-only", action="store_true",
                    help="Skip the digest entirely — only build the deck. No API key needed.")
    ap.add_argument("--no-vision", action="store_true",
                    help="Disable vision-based frame picking for the digest. Cheaper but the "
                         "frames may be less illustrative.")
    ap.add_argument("--digest-model", default="claude-sonnet-4-6",
                    help="Claude model for the digest (default: claude-sonnet-4-6). "
                         "Use claude-opus-4-7 for the highest-quality summarization.")
    ap.add_argument("--scene-threshold", type=float, default=0.2,
                    help="Scene-detection sensitivity, 0.1=lots of frames, 0.5=only major changes (default: 0.2)")
    ap.add_argument("--interval", type=float, default=20.0,
                    help="Also sample one frame every N seconds (0 to disable). Useful for "
                         "screen recordings where gradual changes don't trip scene detection. (default: 20)")
    ap.add_argument("--hash-distance", type=int, default=4,
                    help="Perceptual hash dedup threshold; lower = stricter. Compared only against "
                         "the previous kept frame, so recurring views are preserved. (default: 4)")
    ap.add_argument("--keep-frames", action="store_true",
                    help="Keep extracted frames in ./frames_<videoname>/ instead of cleaning up")
    ap.add_argument("--downloads-dir", type=Path, default=Path("downloads"),
                    help="Where to cache YouTube downloads (default: ./downloads)")
    args = ap.parse_args()

    load_env_files()

    for tool in ("ffmpeg", "ffprobe"):
        if not shutil.which(tool):
            sys.exit(f"{tool} not found on PATH. Install with `brew install ffmpeg` or your package manager.")

    do_digest = not args.deck_only
    if do_digest:
        ensure_api_key()

    if is_url(args.video):
        print(f"[0/5] Fetching YouTube video: {args.video}")
        video_path, srt_path = fetch_youtube(args.video, args.downloads_dir)
        print(f"      mp4: {video_path}")
        print(f"      srt: {srt_path}")
    else:
        video_path = Path(args.video)
        if not video_path.exists():
            sys.exit(f"Video not found: {video_path}")
        if args.srt is None:
            sys.exit("SRT path is required when 'video' is a local file.")
        srt_path = Path(args.srt)
        if not srt_path.exists():
            sys.exit(f"SRT not found: {srt_path}")

    base = video_path.stem
    digest_path = args.output if args.output is not None else Path(f"{base}_digest.md")
    if args.deck == "__default__":
        deck_path = Path(f"{base}_deck.pptx")
    elif args.deck is not None:
        deck_path = Path(args.deck)
    else:
        deck_path = None
    if args.deck_only and deck_path is None:
        deck_path = Path(f"{base}_deck.pptx")

    workdir = Path(tempfile.mkdtemp(prefix="v2d_"))
    scene_dir = workdir / "scene"
    interval_dir = workdir / "interval"

    try:
        duration = get_video_duration(video_path)

        print(f"[1/5] Extracting frames "
              f"(scene threshold={args.scene_threshold}, interval={args.interval}s)...")
        scene_frames = extract_scene_frames(video_path, scene_dir, args.scene_threshold)
        interval_frames = extract_interval_frames(video_path, interval_dir, args.interval, duration)
        frames = merge_frames(scene_frames, interval_frames)
        print(f"      {len(scene_frames)} scene + {len(interval_frames)} interval = "
              f"{len(frames)} candidate frames")

        print(f"[2/5] Deduping consecutive near-identical frames (hash distance <= {args.hash_distance})...")
        frames = dedupe_frames(frames, args.hash_distance)
        print(f"      {len(frames)} unique frames")

        print(f"[3/5] Parsing SRT: {srt_path.name}")
        segments = parse_srt(srt_path)
        print(f"      {len(segments)} transcript segments")

        print("[4/5] Aligning transcript to frames...")
        slides_data = assign_transcript_to_frames(frames, segments, duration)

        if deck_path is not None:
            print(f"[5/5] Building deck -> {deck_path}")
            build_deck(slides_data, deck_path, video_path.stem)
        else:
            print("[5/5] Skipping deck (use --deck to also write one)")

        if do_digest:
            digest_path.parent.mkdir(parents=True, exist_ok=True)
            images_dir = digest_path.parent / f"{digest_path.stem}_images"
            print(f"[+] Generating digest with {args.digest_model} -> {digest_path}")
            digest, usage = generate_digest(segments, video_path.stem, args.digest_model)
            print(f"      {len(digest.topics)} topics  |  "
                  f"input: {usage.input_tokens} tokens "
                  f"(cache read: {getattr(usage, 'cache_read_input_tokens', 0)}, "
                  f"cache write: {getattr(usage, 'cache_creation_input_tokens', 0)})  |  "
                  f"output: {usage.output_tokens} tokens")

            vision_picks = None
            if not args.no_vision:
                print(f"[+] Vision-picking frames with {args.digest_model}...")
                vision_picks, v_usage = vision_pick_frames(
                    digest, frames, duration, args.digest_model
                )
                print(f"      vision-selected {len(vision_picks)}/{len(digest.topics)} topics  |  "
                      f"input: {v_usage.input_tokens} tokens  |  "
                      f"output: {v_usage.output_tokens} tokens")

            write_markdown_digest(digest, frames, duration, digest_path, images_dir, vision_picks)
            print(f"      Digest written. Images in {images_dir}/")

        if args.keep_frames:
            dest = Path.cwd() / f"frames_{video_path.stem}"
            dest.mkdir(parents=True, exist_ok=True)
            for d in (scene_dir, interval_dir):
                if d.exists():
                    shutil.copytree(d, dest, dirs_exist_ok=True)
            print(f"      frames saved to {dest}")

        outputs = []
        if do_digest:
            outputs.append(str(digest_path))
        if deck_path is not None:
            outputs.append(str(deck_path))
        print(f"\nDone. Wrote: {', '.join(outputs)}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


if __name__ == "__main__":
    main()
