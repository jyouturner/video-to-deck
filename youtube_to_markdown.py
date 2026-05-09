#!/usr/bin/env python3
"""
youtube-to-markdown — Turn a YouTube video (or local MP4 + SRT) into a readable
Markdown digest with embedded frame images. Optional PowerPoint export.

Default flow (digest):
  yt2md "https://youtu.be/..."
  yt2md input.mp4 transcript.srt

Also build a deck:
  yt2md "https://youtu.be/..." --deck

Just the deck (no API key needed):
  yt2md input.mp4 transcript.srt --deck-only

Requirements:
  System:  ffmpeg, ffprobe
  API:     ANTHROPIC_API_KEY (prompted on first run unless --deck-only)
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple


# ---------- Step 1: Frame extraction (scene detection + periodic sampling) ----------

# Hard cap on scene frames a single video can yield. Slide-heavy talks
# with subtle reveal animations can blow past 500-1000 candidates at the
# default threshold; beyond a few hundred there's negligible additional
# visual signal but real hashing/IO overhead. Truncating keeps the dedupe
# and vision-pick stages bounded.
SCENE_FRAME_HARD_CAP = 500


def extract_scene_frames(video: Path, out_dir: Path, threshold: float = 0.2) -> List[Tuple[Path, float]]:
    """Run ffmpeg scene detection. Returns list of (frame_path, timestamp_seconds)."""
    out_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        "ffmpeg", "-hide_banner", "-y",
        "-i", str(video),
        "-vf", f"select='eq(n,0)+gt(scene,{threshold})',showinfo",
        "-vsync", "vfr",
        "-q:v", "3",
        str(out_dir / "scene_%04d.jpg"),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg failed:\n{proc.stderr}")

    timestamps = [float(m) for m in re.findall(r"pts_time:([\d.]+)", proc.stderr)]
    frame_files = sorted(out_dir.glob("scene_*.jpg"))
    if len(frame_files) != len(timestamps):
        n = min(len(frame_files), len(timestamps))
        frame_files, timestamps = frame_files[:n], timestamps[:n]

    if len(frame_files) > SCENE_FRAME_HARD_CAP:
        # Evenly-spaced subsample so we keep visual diversity across the
        # whole video, not just the first chunk.
        step = len(frame_files) / SCENE_FRAME_HARD_CAP
        keep = [int(i * step) for i in range(SCENE_FRAME_HARD_CAP)]
        frame_files = [frame_files[i] for i in keep]
        timestamps = [timestamps[i] for i in keep]

    return list(zip(frame_files, timestamps))


def extract_interval_frames(video: Path, out_dir: Path, interval: float, duration: float) -> List[Tuple[Path, float]]:
    """Sample one frame every `interval` seconds. Returns (frame_path, timestamp) list."""
    out_dir.mkdir(parents=True, exist_ok=True)
    if interval <= 0:
        return []

    fps = 1.0 / interval
    cmd = [
        "ffmpeg", "-hide_banner", "-y",
        "-i", str(video),
        "-vf", f"fps={fps}",
        "-q:v", "3",
        str(out_dir / "interval_%04d.jpg"),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg interval sampling failed:\n{proc.stderr}")

    frame_files = sorted(out_dir.glob("interval_*.jpg"))
    # ffmpeg's fps filter places the first frame at t≈interval/2; reconstruct timestamps.
    timestamps = [min(duration, interval * (i + 0.5)) for i in range(len(frame_files))]
    return list(zip(frame_files, timestamps))


def merge_frames(*frame_lists: List[Tuple[Path, float]]) -> List[Tuple[Path, float]]:
    """Merge multiple frame lists, sort by timestamp."""
    combined: List[Tuple[Path, float]] = []
    for fl in frame_lists:
        combined.extend(fl)
    combined.sort(key=lambda x: x[1])
    return combined


# ---------- Step 2: Perceptual-hash dedup (consecutive only) ----------

def dedupe_frames(frames: List[Tuple[Path, float]], hash_distance: int = 4) -> List[Tuple[Path, float]]:
    """Cluster runs of near-identical consecutive frames; keep the LAST of each cluster.

    Animated slide reveals (bullets / diagram elements appearing over time) and slow
    pans both hit this case: each intermediate frame is similar to its neighbor, but
    only the final frame shows the fully-revealed / settled state. Keeping the last
    frame of the cluster preserves that state rather than the partial opening one.

    Discrete scene changes (dist > threshold between neighbors) break the cluster,
    so recurring views — e.g. switching to an editor and back to slides — still
    survive as separate kept frames.
    """
    if not frames:
        return []

    import imagehash
    from PIL import Image

    hashed: List[Tuple[Path, float, "imagehash.ImageHash"]] = []
    for path, ts in frames:
        with Image.open(path) as im:
            hashed.append((path, ts, imagehash.phash(im)))

    clusters: List[List[Tuple[Path, float, "imagehash.ImageHash"]]] = [[hashed[0]]]
    for i in range(1, len(hashed)):
        if (hashed[i][2] - hashed[i - 1][2]) <= hash_distance:
            clusters[-1].append(hashed[i])
        else:
            clusters.append([hashed[i]])

    kept: List[Tuple[Path, float]] = []
    for cluster in clusters:
        for path, _, _ in cluster[:-1]:
            path.unlink(missing_ok=True)
        kept.append((cluster[-1][0], cluster[-1][1]))
    return kept


def global_phash_cluster(
    frames: List[Tuple[Path, float]],
    distance: int = 6,
) -> List[Tuple[Path, float]]:
    """Drop frames whose perceptual hash is close to ANY earlier kept frame.

    Consecutive dedup (dedupe_frames) merges runs of similar adjacent frames
    but doesn't catch the talk-deck pattern where the camera cuts speaker→
    slide→speaker→same-slide. Each return to the slide creates a new
    consecutive cluster, so we end up with N copies of the same slide.

    This is an O(N²) pass that compares each new frame against every
    previously-kept hash. N is small after consecutive dedup (typically
    <200), so the cost is negligible. The chronologically-first occurrence
    of a slide is kept; later returns are dropped.

    The default distance (6) is looser than dedupe_frames' default (4) —
    we want to merge "same slide with speaker overlay in the corner" or
    "same chart at 480p vs 720p", which can drift past 4 in pHash space.
    """
    if not frames:
        return []
    import imagehash
    from PIL import Image

    kept: List[Tuple[Path, float]] = []
    kept_hashes: List["imagehash.ImageHash"] = []
    for path, ts in frames:
        with Image.open(path) as im:
            h = imagehash.phash(im)
        if any((h - kh) <= distance for kh in kept_hashes):
            continue
        kept.append((path, ts))
        kept_hashes.append(h)
    return kept


# ---------- Vision-grid slide classifier ----------
#
# Tile candidate frames into 3×3 grid images and ask Claude (Haiku, by
# default) to classify each cell as NEW_SLIDE / SAME_AS_PREVIOUS_CELL /
# TALKING_HEAD / TRANSITION. Sending grids instead of single frames is
# ~10× cheaper in input tokens and lets the LLM compare adjacent cells
# directly — easier than reasoning across separate API calls.

_GRID_COLS = 3
_GRID_ROWS = 3
_GRID_CELLS = _GRID_COLS * _GRID_ROWS  # 9
_GRID_CELL_W = 400  # 16:9 thumbnails
_GRID_CELL_H = 225


def _render_classification_grids(
    frames: List[Tuple[Path, float]],
    out_dir: Path,
) -> List[Tuple[Path, List[int]]]:
    """Render frames into one or more numbered 3×3 grid images.

    Returns a list of (grid_path, [original_frame_index_for_each_cell]). The
    last grid may be padded with blank cells; the index list reflects only
    real frames (length ≤ 9 per grid). Cells are numbered 1..9 with a small
    overlay so the LLM can address them unambiguously.

    Each grid carries a 1-cell overlap from the previous grid (the last real
    frame of grid N-1 is duplicated as cell 1 of grid N) so the LLM has
    context for the SAME-AS-PREVIOUS judgement at boundaries.
    """
    from PIL import Image, ImageDraw, ImageFont

    out_dir.mkdir(parents=True, exist_ok=True)
    if not frames:
        return []

    # Try a system font for legibility; fall back to PIL default if not available.
    try:
        font = ImageFont.truetype(
            "/System/Library/Fonts/Helvetica.ttc", 32,
        )
    except (OSError, IOError):
        font = ImageFont.load_default()

    grids: List[Tuple[Path, List[int]]] = []
    grid_w = _GRID_CELL_W * _GRID_COLS
    grid_h = _GRID_CELL_H * _GRID_ROWS

    # Walk the frame list in chunks of 8 (leaving cell 0 of each grid for
    # the overlap from the previous grid). The first grid has no overlap so
    # it gets all 9 cells.
    i = 0
    grid_idx = 0
    prev_last_idx: Optional[int] = None
    while i < len(frames):
        canvas = Image.new("RGB", (grid_w, grid_h), color="black")
        draw = ImageDraw.Draw(canvas)
        cell_indices: List[int] = []

        # Position 0: overlap cell from previous grid (None on the first grid).
        # We embed the overlap so the model can judge "is cell 1 the same as
        # cell 0?" right at the grid boundary.
        if prev_last_idx is not None and grid_idx > 0:
            cells_to_fill = [prev_last_idx]
            slots_for_new = _GRID_CELLS - 1
        else:
            cells_to_fill = []
            slots_for_new = _GRID_CELLS

        # Fill the remaining cells with new frames.
        end = min(i + slots_for_new, len(frames))
        cells_to_fill.extend(range(i, end))

        for cell_pos, frame_idx in enumerate(cells_to_fill):
            row, col = divmod(cell_pos, _GRID_COLS)
            x0 = col * _GRID_CELL_W
            y0 = row * _GRID_CELL_H
            with Image.open(frames[frame_idx][0]) as im:
                im.thumbnail((_GRID_CELL_W, _GRID_CELL_H))
                # Center the thumbnail in its cell (frames may not be exact 16:9).
                tw, th = im.size
                px = x0 + (_GRID_CELL_W - tw) // 2
                py = y0 + (_GRID_CELL_H - th) // 2
                canvas.paste(im, (px, py))
            # Big numbered label, white-on-black with a small offset, so it
            # stands out against either light or dark cell backgrounds.
            label = str(cell_pos + 1)
            draw.rectangle(
                [(x0 + 4, y0 + 4), (x0 + 44, y0 + 44)],
                fill="black", outline="white",
            )
            draw.text((x0 + 12, y0 + 4), label, fill="white", font=font)
            cell_indices.append(frame_idx)

        grid_path = out_dir / f"grid_{grid_idx:03d}.jpg"
        canvas.save(grid_path, "JPEG", quality=85)
        # Subtract 1 from cell_indices to get the "real new frames" range
        # (cell 0 may be the overlap, which the caller should skip when
        # mapping LLM responses back).
        grids.append((grid_path, cell_indices))

        prev_last_idx = cells_to_fill[-1]
        i = end
        grid_idx += 1

    return grids


def classify_slides_via_grids(
    frames: List[Tuple[Path, float]],
    *,
    backend,
    model: str,
    workdir: Path,
    log_video_id: Optional[str] = None,
) -> List[Tuple[Path, float]]:
    """Use vision LLM (typically Haiku via grids) to filter frames down to
    distinct deck slides. Returns the kept frames in chronological order.

    On any failure (LLM error, malformed response, or implausibly small
    output), returns the input frames unchanged so the caller still gets a
    deck — just a less-pruned one.

    log_video_id (optional): when provided, each per-grid LLM call is
    recorded to the cost-audit log under kind='slide_classifier'.
    """
    from pydantic import BaseModel
    from typing import List as TList, Literal

    class CellLabel(BaseModel):
        cell: int  # 1..9
        label: Literal["NEW_SLIDE", "SAME_AS_PREVIOUS_CELL",
                       "TALKING_HEAD", "TRANSITION"]

    class GridLabels(BaseModel):
        labels: TList[CellLabel]

    if not frames:
        return frames
    if not getattr(backend, "vision_supported", False):
        return frames  # backend without vision (e.g. Claude Code w/o opt-in)

    grid_dir = workdir / "slide_classifier_grids"
    grids = _render_classification_grids(frames, grid_dir)
    if not grids:
        return frames

    system_prompt = (
        "You classify frames extracted from a video that contains a slide "
        "deck. Each grid image you receive is a 3×3 layout (cells numbered "
        "1–9, top-left to bottom-right). For each cell, decide what the "
        "frame is and whether it shows a slide we haven't already seen.\n\n"
        "Labels:\n"
        "- NEW_SLIDE — a slide whose content (text, layout, charts) is "
        "different from cell N-1 in the same grid AND from any slide you've "
        "already labeled NEW_SLIDE in earlier grids.\n"
        "- SAME_AS_PREVIOUS_CELL — visually the same slide as cell N-1 (or, "
        "for cell 1 of grids after the first, the same as the overlap cell). "
        "Animations and reveals count as same.\n"
        "- TALKING_HEAD — frame is dominated by the speaker, no slide visible.\n"
        "- TRANSITION — fade, blur, mid-cut, or otherwise not a clean slide.\n\n"
        "Important:\n"
        "- Cell 1 of grids 2+ duplicates the LAST real cell of the previous "
        "grid as boundary context. Use it to compare cell 2 against the "
        "previous grid's last frame.\n"
        "- Output exactly one label per cell present in the grid, in cell-"
        "number order. If the grid has fewer than 9 real frames, only "
        "output labels for the cells that contain images (skip blank ones)."
    )

    kept_global: set = set()  # original frame indices we've decided to keep
    for grid_idx, (grid_path, cell_indices) in enumerate(grids):
        # Build the message: grid image + a short reminder of cell count.
        with open(grid_path, "rb") as f:
            grid_bytes = f.read()
        import base64 as _b64
        b64 = _b64.b64encode(grid_bytes).decode("ascii")
        content_blocks = [
            {"type": "text", "text":
             f"Grid {grid_idx + 1} of {len(grids)} — "
             f"{len(cell_indices)} cell(s) populated."},
            {"type": "image", "source": {
                "type": "base64", "media_type": "image/jpeg", "data": b64,
            }},
            {"type": "text", "text":
             "Classify each cell. Return one label per real cell, in order."},
        ]
        try:
            parsed, grid_usage = backend.vision_parse(
                system=system_prompt,
                content_blocks=content_blocks,
                model=model,
                max_tokens=600,
                schema=GridLabels,
            )
        except Exception:
            # Defensive: any LLM failure leaves the candidate set intact.
            return frames
        if log_video_id is not None:
            record_llm_usage(
                video_id=log_video_id, kind="slide_classifier",
                model=model, backend_name=backend.name, usage=grid_usage,
            )

        for cell_label in parsed.labels:
            cell_pos = cell_label.cell - 1  # back to 0-indexed
            if cell_pos < 0 or cell_pos >= len(cell_indices):
                continue
            frame_idx = cell_indices[cell_pos]
            # Skip the overlap cell on grids 2+ (cell 0) — already considered.
            if grid_idx > 0 and cell_pos == 0:
                continue
            if cell_label.label == "NEW_SLIDE":
                kept_global.add(frame_idx)

    if not kept_global:
        # Implausible result (LLM said no NEW_SLIDE anywhere) — fall back.
        return frames
    return [frames[i] for i in sorted(kept_global)]


# ---------- Step 3: SRT parser ----------

@dataclass
class TranscriptSegment:
    start: float
    end: float
    text: str


def _srt_time_to_seconds(t: str) -> float:
    """'00:01:23,456' -> 83.456"""
    h, m, rest = t.split(":")
    s, ms = rest.split(",")
    return int(h) * 3600 + int(m) * 60 + int(s) + int(ms) / 1000


def parse_srt(srt_path: Path) -> List[TranscriptSegment]:
    """Parse a .srt file. Tolerant of BOM, CRLF, dot-vs-comma ms separator, simple markup."""
    text = srt_path.read_text(encoding="utf-8-sig")
    blocks = re.split(r"\r?\n\r?\n", text.strip())

    segments: List[TranscriptSegment] = []
    time_re = re.compile(r"(\d{2}:\d{2}:\d{2}[,.]\d{3})\s*-->\s*(\d{2}:\d{2}:\d{2}[,.]\d{3})")

    for block in blocks:
        lines = [ln for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        time_line_idx = next((i for i, ln in enumerate(lines) if time_re.search(ln)), None)
        if time_line_idx is None:
            continue
        m = time_re.search(lines[time_line_idx])
        start = _srt_time_to_seconds(m.group(1).replace(".", ","))
        end = _srt_time_to_seconds(m.group(2).replace(".", ","))
        body = " ".join(lines[time_line_idx + 1:])
        body = re.sub(r"<[^>]+>", "", body)         # <i>, <b>, <font>
        body = re.sub(r"\{[^}]+\}", "", body).strip()  # {\an8}, etc.
        if body:
            segments.append(TranscriptSegment(start=start, end=end, text=body))

    return segments


# ---------- Step 4: Align transcript to frames ----------

def assign_transcript_to_frames(
    frames: List[Tuple[Path, float]],
    segments: List[TranscriptSegment],
    video_duration: float,
) -> List[Tuple[Path, float, float, str]]:
    """
    For each frame at time t_i, the slide covers [t_i, t_{i+1}).
    Collect transcript segments whose midpoint falls in that window.
    Returns: list of (frame_path, slide_start, slide_end, transcript_text).
    """
    results = []
    for i, (path, start) in enumerate(frames):
        end = frames[i + 1][1] if i + 1 < len(frames) else video_duration
        chunk = " ".join(
            seg.text for seg in segments
            if start <= (seg.start + seg.end) / 2 < end
        )
        results.append((path, start, end, chunk))
    return results


def get_video_duration(video: Path) -> float:
    out = subprocess.check_output([
        "ffprobe", "-v", "error", "-show_entries", "format=duration",
        "-of", "json", str(video),
    ])
    return float(json.loads(out)["format"]["duration"])


# ---------- Step 5: Build the .pptx ----------

def format_timestamp(seconds: float) -> str:
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    return f"{h:d}:{m:02d}:{s:02d}" if h else f"{m:d}:{s:02d}"


def build_deck(slides_data, output: Path, video_name: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = video_name
    p.font.size = Pt(40)
    p.font.bold = True
    sub = title_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.5))
    sub.text_frame.paragraphs[0].text = f"{len(slides_data)} slides extracted from video"
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    # Layout: image fills most of the slide (image-first, deck-replication
    # framing). The full transcript chunk lives in PowerPoint's speaker
    # notes so anyone who wants the narration can open it via View → Notes
    # Page; the slide itself stays clean.
    img_top_in = 0.3
    img_max_h_in = 6.8
    img_max_w_in = 12.33
    footer_top_in = 7.15

    for idx, (frame_path, start, end, transcript) in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        with Image.open(frame_path) as im:
            iw, ih = im.size
        scale = min(img_max_w_in / (iw / 96), img_max_h_in / (ih / 96))
        disp_w = (iw / 96) * scale
        disp_h = (ih / 96) * scale
        left = (13.33 - disp_w) / 2

        slide.shapes.add_picture(
            str(frame_path),
            Inches(left), Inches(img_top_in),
            width=Inches(disp_w), height=Inches(disp_h),
        )

        # Footer with time range + slide number — useful provenance, small
        # enough not to compete with the image.
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(footer_top_in), Inches(12.7), Inches(0.3))
        fp = footer.text_frame.paragraphs[0]
        fp.text = f"{format_timestamp(start)} – {format_timestamp(end)}   |   Slide {idx}"
        fp.font.size = Pt(10)

        if transcript:
            slide.notes_slide.notes_text_frame.text = transcript

    prs.save(str(output))


# ---------- Config / API key handling ----------

# Everything lives in one visible directory under HOME. Override with YT2MD_DATA.
# Layout under that directory:
#   .env, channels.txt, state.json, digests/, meta/, downloads/, logs/

DEFAULT_DATA_DIR = Path.home() / "yt2md"


def get_data_dir() -> Path:
    return Path(os.environ.get("YT2MD_DATA", str(DEFAULT_DATA_DIR))).expanduser()


def env_file() -> Path:
    return get_data_dir() / ".env"


def load_env_files() -> None:
    """Populate os.environ from .env files. Real env vars always win.

    Order (lowest priority first; later loads do NOT override earlier-set keys):
      1. Real env vars (from the shell)
      2. CWD/.env (project-local override)
      3. <data dir>/.env (default: ~/yt2md/.env)
    """
    from dotenv import load_dotenv

    load_dotenv()  # CWD/.env, only fills in missing
    e = env_file()
    if e.exists():
        load_dotenv(e)


def set_env_var(name: str, value: str) -> Path:
    """Persist <name>=<value> to ~/yt2md/.env, preserving other entries.

    Uses dotenv.set_key for round-trip safe updates (vs. naive overwrite, which
    would clobber co-resident keys). Also updates os.environ so the running
    process sees the new value immediately. Returns the .env path.
    """
    from dotenv import set_key

    e = env_file()
    e.parent.mkdir(parents=True, exist_ok=True)
    if not e.exists():
        e.touch()
    try:
        os.chmod(e, 0o600)
    except OSError:
        pass
    set_key(str(e), name, value, quote_mode="never")
    os.environ[name] = value
    return e


API_KEY_COST_NOTE = (
    "Anthropic bills your API key per request (separate from any Claude.ai "
    "subscription). Rough costs: a 30-min digest is ~$0.03 with "
    "<code>claude-sonnet-4-6</code> (default), ~$0.15 with "
    "<code>claude-opus-4-7</code>. The panel discussion adds one Opus call "
    "(~$0.10). Add a payment method at "
    '<a href="https://console.anthropic.com/settings/billing" target="_blank" '
    'rel="noopener">console.anthropic.com/settings/billing</a>.'
)


def validate_api_key(key: str) -> Optional[str]:
    """Send a 1-token request to the cheapest model to verify auth. Returns
    None on success, otherwise a short human-readable error string.
    """
    import anthropic

    try:
        client = anthropic.Anthropic(api_key=key)
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1,
            messages=[{"role": "user", "content": "ok"}],
        )
        # Even validation pings cost a few tokens — record so the audit
        # log is complete.
        record_llm_usage(
            video_id=None, kind="validation",
            model="claude-haiku-4-5-20251001",
            backend_name="api", usage=resp.usage,
        )
        return None
    except anthropic.AuthenticationError:
        return "key rejected by Anthropic (authentication failed)."
    except anthropic.PermissionDeniedError as e:
        return f"key authenticated but lacks permission: {e}"
    except anthropic.APIConnectionError as e:
        return f"could not reach Anthropic: {e}"
    except Exception as e:
        return f"unexpected error: {type(e).__name__}: {e}"


def ensure_api_key() -> None:
    """Make sure ANTHROPIC_API_KEY is set, prompting + saving on first run if interactive."""
    if os.environ.get("ANTHROPIC_API_KEY"):
        return

    e = env_file()
    msg = (
        "ANTHROPIC_API_KEY is not set.\n"
        "Get a key from: https://console.anthropic.com/settings/keys"
    )
    if not sys.stdin.isatty():
        sys.exit(
            f"{msg}\n"
            "Then either export it (`export ANTHROPIC_API_KEY=...`) or save it via:\n"
            f"  mkdir -p {e.parent} && echo 'ANTHROPIC_API_KEY=sk-ant-...' >> {e}"
        )

    print(msg)
    key = input("Paste your API key (or press Enter to abort): ").strip()
    if not key:
        sys.exit("Aborted.")

    save = input(
        f"Save it to {e} so future runs find it automatically? [Y/n] "
    ).strip().lower()
    if save in ("", "y", "yes"):
        set_env_var("ANTHROPIC_API_KEY", key)
        print(f"      saved to {e}")
    else:
        os.environ["ANTHROPIC_API_KEY"] = key


# ---------- Claude Code sandbox (alternative auth path) ----------
#
# yt2md ships a private Claude Code install under <data dir>/claude-code/.
# Driving the official `claude` binary as a subprocess is the supported way
# for a third-party tool to leverage a user's Claude.ai subscription auth
# without violating ToS (vs. extracting OAuth tokens, which is forbidden).
#
# Sandbox layout:
#   <data>/claude-code/node_modules/.bin/claude   <- the binary we invoke
#   <data>/claude-code/config/                     <- CLAUDE_CONFIG_DIR target
#       settings.json, projects/, plugins/, .credentials.json (Linux/Win)
# On macOS, CLAUDE_CONFIG_DIR isolates settings/projects but credentials still
# go to the system Keychain — this is a documented Claude Code limitation.

CLAUDE_CODE_NPM_PACKAGE = "@anthropic-ai/claude-code"  # always pulls latest
MIN_NODE_MAJOR = 18


def claude_sandbox_dir() -> Path:
    return get_data_dir() / "claude-code"


def claude_config_dir() -> Path:
    return claude_sandbox_dir() / "config"


def claude_binary_path() -> Path:
    return claude_sandbox_dir() / "node_modules" / ".bin" / "claude"


def claude_subprocess_env() -> dict:
    """Env dict for invoking the sandboxed claude binary. Pins CLAUDE_CONFIG_DIR
    into our sandbox so settings/projects/plugins don't collide with any
    system-wide Claude Code install. (macOS Keychain owns credentials regardless.)
    """
    env = os.environ.copy()
    env["CLAUDE_CONFIG_DIR"] = str(claude_config_dir())
    return env


def detect_node() -> Optional[Tuple[int, str]]:
    """Returns (major_version, full_version_string) if Node ≥ 18 is available,
    else None. We need npm to install the sandbox; npm requires Node.
    """
    node = shutil.which("node")
    if not node:
        return None
    try:
        out = subprocess.check_output([node, "-v"], text=True, timeout=5).strip()
    except (subprocess.SubprocessError, OSError):
        return None
    # `node -v` prints e.g. "v18.19.0"
    m = re.match(r"^v(\d+)\.", out)
    if not m:
        return None
    major = int(m.group(1))
    if major < MIN_NODE_MAJOR:
        return None
    return major, out


def claude_code_installed() -> bool:
    """True if our sandbox has a working claude binary."""
    return claude_binary_path().exists()


def install_claude_code(stream_to: Optional[Path] = None) -> Tuple[int, str]:
    """Run `npm install --prefix <sandbox> <package>` to materialize a private
    Claude Code install. Returns (returncode, combined_output). If stream_to
    is given, output is appended there as it's produced (for live UI polling).
    """
    sandbox = claude_sandbox_dir()
    sandbox.mkdir(parents=True, exist_ok=True)
    claude_config_dir().mkdir(parents=True, exist_ok=True)

    npm = shutil.which("npm")
    if npm is None:
        return 127, "npm not found on PATH (install Node.js 18+ first)."

    cmd = [npm, "install", "--prefix", str(sandbox),
           "--no-fund", "--no-audit", "--silent",
           CLAUDE_CODE_NPM_PACKAGE]
    proc = subprocess.Popen(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True,
    )
    chunks: list = []
    assert proc.stdout is not None
    for line in proc.stdout:
        chunks.append(line)
        if stream_to is not None:
            try:
                with open(stream_to, "a") as f:
                    f.write(line)
            except OSError:
                pass
    proc.wait()
    return proc.returncode, "".join(chunks)


def claude_code_logged_in() -> bool:
    """Heuristic: if there's a working binary AND `claude /status` (or a tiny
    -p call) succeeds without auth error, we're logged in. Cheap probe via the
    presence of credentials state."""
    if not claude_code_installed():
        return False
    # The cheapest probe is a 1-token --print call; cache the result for the
    # lifetime of the process (set when we explicitly log in/out).
    return _claude_code_session_state.get("logged_in", False)


# Module-level cache of login probe state. Populated by validate_claude_code()
# after install/login; checked by claude_code_logged_in() to avoid spawning a
# subprocess on every page load.
_claude_code_session_state: dict = {}


def _claude_login_sentinel() -> Path:
    """Touched after a successful validation; lets us assume logged-in across
    server restarts without burning a token-cost probe per boot. Cleared by
    claude_logout(). Real auth failures still reset the session state when
    encountered."""
    return claude_config_dir() / ".yt2md-logged-in"


def claude_probe_login_state() -> None:
    """Cheap startup-time probe: read the sentinel file and populate the
    session-state cache. No subprocess call. Real auth-failure reset happens
    in validate_claude_code() and on first failed LLM call.
    """
    if claude_code_installed() and _claude_login_sentinel().exists():
        _claude_code_session_state["logged_in"] = True
    else:
        _claude_code_session_state["logged_in"] = False


def validate_claude_code() -> Optional[str]:
    """Run a 1-token call through the sandboxed claude binary. Returns None on
    success, otherwise a short error string. Updates the session-state cache
    and the on-disk sentinel.
    """
    if not claude_code_installed():
        return "Claude Code is not installed in the sandbox."
    cmd = [
        str(claude_binary_path()), "-p", "ok",
        "--model", "claude-haiku-4-5-20251001",
        "--output-format", "json",
    ]
    try:
        proc = subprocess.run(
            cmd, env=claude_subprocess_env(),
            capture_output=True, text=True, timeout=60,
        )
    except subprocess.TimeoutExpired:
        _claude_code_session_state["logged_in"] = False
        _claude_login_sentinel().unlink(missing_ok=True)
        return "validation timed out (60s) — the OAuth flow may not have completed."
    except OSError as e:
        _claude_code_session_state["logged_in"] = False
        _claude_login_sentinel().unlink(missing_ok=True)
        return f"could not invoke claude binary: {e}"
    if proc.returncode != 0:
        _claude_code_session_state["logged_in"] = False
        _claude_login_sentinel().unlink(missing_ok=True)
        msg = (proc.stderr or proc.stdout or "").strip().splitlines()
        last = msg[-1] if msg else f"exit code {proc.returncode}"
        return f"claude returned an error: {last}"
    _claude_code_session_state["logged_in"] = True
    try:
        _claude_login_sentinel().parent.mkdir(parents=True, exist_ok=True)
        _claude_login_sentinel().touch()
    except OSError:
        pass
    return None


# In-memory tracker for async setup jobs (install, login). Keyed by job name
# ("install" | "login"). Each entry: {"proc": Popen, "log": Path, "started":
# epoch, "error": Optional[str]}. The web reader is single-process, so a plain
# dict suffices. State resets on server restart, which is fine — finished jobs
# leave their result on disk (sandbox dir exists; credentials persist).

_claude_setup_jobs: dict = {}


def _claude_setup_log(name: str) -> Path:
    p = get_data_dir() / "logs" / f"claude-setup-{name}.log"
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def claude_setup_job_running(name: str) -> bool:
    job = _claude_setup_jobs.get(name)
    if job is None:
        return False
    proc = job.get("proc")
    return proc is not None and proc.poll() is None


def start_install_job() -> Optional[str]:
    """Spawn `npm install` for Claude Code in the sandbox. Idempotent: returns
    None if it's already running, or an error string if Node is missing.
    """
    if claude_setup_job_running("install"):
        return None
    if detect_node() is None:
        return (
            f"Node.js {MIN_NODE_MAJOR}+ is required. Install with "
            "`brew install node` (macOS) or your package manager, then retry."
        )
    sandbox = claude_sandbox_dir()
    sandbox.mkdir(parents=True, exist_ok=True)
    claude_config_dir().mkdir(parents=True, exist_ok=True)
    log_path = _claude_setup_log("install")
    log_path.write_text("")  # truncate prior run
    npm = shutil.which("npm") or "npm"
    cmd = [npm, "install", "--prefix", str(sandbox),
           "--no-fund", "--no-audit",
           CLAUDE_CODE_NPM_PACKAGE]
    log_f = open(log_path, "a")
    proc = subprocess.Popen(
        cmd, stdout=log_f, stderr=subprocess.STDOUT, text=True,
        start_new_session=True,
    )
    import time as _t
    _claude_setup_jobs["install"] = {
        "proc": proc, "log": log_path, "started": _t.time(),
        "log_f": log_f, "error": None,
    }
    return None


def start_login_job() -> Optional[str]:
    """Spawn `claude /login`. The CLI auto-opens the user's default browser
    via `open` for OAuth; the random localhost callback is captured by the
    subprocess itself. We wait for exit-0 then validate.
    """
    if claude_setup_job_running("login"):
        return None
    if not claude_code_installed():
        return "Claude Code is not installed yet. Install it first."
    log_path = _claude_setup_log("login")
    log_path.write_text("")
    cmd = [str(claude_binary_path()), "/login"]
    log_f = open(log_path, "a")
    proc = subprocess.Popen(
        cmd, env=claude_subprocess_env(),
        stdout=log_f, stderr=subprocess.STDOUT,
        stdin=subprocess.DEVNULL,  # no terminal stdin → CLI relies on browser callback
        text=True,
        start_new_session=True,
    )
    import time as _t
    _claude_setup_jobs["login"] = {
        "proc": proc, "log": log_path, "started": _t.time(),
        "log_f": log_f, "error": None, "validated": False,
    }
    return None


def claude_logout() -> Tuple[int, str]:
    """Run `claude /logout` to clear stored credentials. Resets session cache."""
    if not claude_code_installed():
        return 0, "(not installed)"
    cmd = [str(claude_binary_path()), "/logout"]
    proc = subprocess.run(
        cmd, env=claude_subprocess_env(),
        capture_output=True, text=True, timeout=30,
    )
    _claude_code_session_state["logged_in"] = False
    _claude_login_sentinel().unlink(missing_ok=True)
    return proc.returncode, (proc.stdout or "") + (proc.stderr or "")


def claude_setup_snapshot() -> dict:
    """Status blob consumed by the /setup polling JS. Reaps finished jobs and
    runs validation after a successful login.
    """
    install_job = _claude_setup_jobs.get("install")
    login_job = _claude_setup_jobs.get("login")

    # Reap install if it finished — surface the exit code as an error if non-zero.
    if install_job and install_job["proc"].poll() is not None:
        rc = install_job["proc"].returncode
        if rc != 0 and not install_job.get("error"):
            install_job["error"] = f"npm install failed (exit {rc}). See log."
        try:
            install_job["log_f"].close()
        except Exception:
            pass

    # Reap login. On success, validate with a 1-token call (cached).
    if login_job and login_job["proc"].poll() is not None:
        try:
            login_job["log_f"].close()
        except Exception:
            pass
        if not login_job.get("validated"):
            rc = login_job["proc"].returncode
            if rc == 0:
                err = validate_claude_code()
                if err:
                    login_job["error"] = err
            else:
                login_job["error"] = f"login subprocess exited with code {rc}. See log."
            login_job["validated"] = True

    def _tail(path: Path, n: int = 30) -> str:
        try:
            lines = path.read_text(errors="replace").splitlines()
            return "\n".join(lines[-n:])
        except OSError:
            return ""

    return {
        "node_ok": detect_node() is not None,
        "installed": claude_code_installed(),
        "logged_in": _claude_code_session_state.get("logged_in", False),
        "install_running": claude_setup_job_running("install"),
        "login_running": claude_setup_job_running("login"),
        "install_log_tail": _tail(install_job["log"]) if install_job else "",
        "login_log_tail": _tail(login_job["log"]) if login_job else "",
        "install_error": install_job.get("error") if install_job else None,
        "login_error": login_job.get("error") if login_job else None,
    }


# ---------- YouTube fetch ----------

URL_RE = re.compile(r"^https?://", re.IGNORECASE)


def is_url(s: str) -> bool:
    return bool(URL_RE.match(s))


DEFAULT_WHISPER_MODEL = "medium"


def _whisper_secs_to_srt(secs: float) -> str:
    """Convert seconds to SRT-style HH:MM:SS,mmm timestamp."""
    if secs < 0:
        secs = 0.0
    h = int(secs // 3600)
    m = int((secs % 3600) // 60)
    s = int(secs % 60)
    ms = int(round((secs - int(secs)) * 1000))
    if ms == 1000:  # rounding can push us a full ms over
        s += 1
        ms = 0
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"


def _transcribe_with_whisper(
    media_path: Path,
    out_dir: Path,
    video_id: str,
    model_name: str = DEFAULT_WHISPER_MODEL,
) -> Tuple[Path, str]:
    """Transcribe a media file with faster-whisper, write an SRT, return (srt_path, lang).

    Model weights are downloaded on first use to ~/.cache/huggingface and reused
    afterwards. Detected language is used as the SRT filename suffix so the
    existing cache-by-glob logic in fetch_youtube picks it up on re-runs.
    """
    try:
        from faster_whisper import WhisperModel
    except ImportError as e:
        raise RuntimeError(
            "Whisper transcription needed but faster-whisper is not installed. "
            "Run `uv sync` (or `pip install faster-whisper`) and try again."
        ) from e

    print(f"      loading whisper model '{model_name}' (first run downloads weights)")
    # int8 keeps memory low and runs well on CPU; faster-whisper picks Metal/CUDA
    # automatically when device='auto'.
    model = WhisperModel(model_name, device="auto", compute_type="int8")

    print(f"      transcribing audio with whisper ({media_path.name})...")
    segments_iter, info = model.transcribe(
        str(media_path),
        beam_size=5,
        vad_filter=True,  # cuts long silences so the transcript stays tight
    )

    lang = info.language or "und"
    srt_path = out_dir / f"{video_id}.{lang}.srt"

    n = 0
    with srt_path.open("w") as fh:
        for seg in segments_iter:
            n += 1
            text = seg.text.strip().replace("\n", " ")
            if not text:
                continue
            fh.write(
                f"{n}\n"
                f"{_whisper_secs_to_srt(seg.start)} --> {_whisper_secs_to_srt(seg.end)}\n"
                f"{text}\n\n"
            )

    print(
        f"      whisper: {n} segments, language='{lang}' "
        f"(prob={info.language_probability:.2f})"
    )
    return srt_path, lang


def _ensure_js_runtime_available() -> Optional[str]:
    """Find a JS runtime usable for yt-dlp's n-challenge solver.

    yt-dlp accepts deno / node / bun. As of 2026, it marks Node <20 as
    'unsupported' — silently failing the n-challenge and producing only
    storyboard formats. So for Node we collect all candidates (PATH match +
    common version-manager locations: nvm, fnm, asdf, volta), version-rank
    them, and prepend the dir of the highest version to os.environ['PATH']
    so yt-dlp's internal lookups pick the right one.
    """
    # deno / bun: trust the first PATH match (no version-rank needed).
    for rt in ("deno", "bun"):
        if shutil.which(rt):
            return rt

    home = Path.home()
    seen: set = set()
    node_candidates: List[Path] = []

    def _add(p: Optional[Path]):
        if p and p.is_file() and str(p) not in seen:
            seen.add(str(p))
            node_candidates.append(p)

    path_node = shutil.which("node")
    if path_node:
        _add(Path(path_node))
    for p in sorted((home / ".nvm" / "versions" / "node").glob("*/bin/node"), reverse=True):
        _add(p)
    for p in sorted(
        (home / ".local" / "share" / "fnm" / "node-versions").glob("*/installation/bin/node"),
        reverse=True,
    ):
        _add(p)
    for p in sorted((home / ".asdf" / "installs" / "nodejs").glob("*/bin/node"), reverse=True):
        _add(p)
    _add(home / ".volta" / "bin" / "node")

    best_path: Optional[Path] = None
    best_version: Tuple[int, ...] = (0, 0, 0)
    for c in node_candidates:
        try:
            r = subprocess.run(
                [str(c), "--version"], capture_output=True, text=True, timeout=5
            )
            v = tuple(int(p) for p in r.stdout.strip().lstrip("v").split(".")[:3])
        except Exception:
            continue
        if v > best_version:
            best_version = v
            best_path = c

    if best_path is None:
        return None

    # Prepend its dir so yt-dlp's shutil.which lookups pick this one.
    bin_dir = str(best_path.parent)
    current = shutil.which("node")
    if current != str(best_path):
        os.environ["PATH"] = f"{bin_dir}{os.pathsep}{os.environ.get('PATH', '')}"
        print(
            f"[yt2md] using node v{'.'.join(map(str, best_version))} at "
            f"{best_path} for yt-dlp"
        )

    if best_version < (20, 0, 0):
        print(
            f"[yt2md] warning: node v{'.'.join(map(str, best_version))} is below "
            "yt-dlp's required v20+ for YouTube JS challenges. Install a newer "
            "node (`nvm install 20`) or deno (`brew install deno`).",
            file=sys.stderr,
        )

    return "node"


def _pick_caption_lang(info: dict) -> Optional[Tuple[str, bool]]:
    """Choose best caption track from a yt-dlp info dict.

    Returns (lang_code, is_manual) or None if no captions exist at all.

    Priority:
      1. Manual English (any en-* code)
      2. Manual matching the audio language (info['language'])
      3. Any manual track
      4. Auto matching the audio language
      5. Auto English
      6. Any auto track

    Manual is preferred over auto everywhere. Within auto, we prefer the
    original audio language over English: YouTube's auto-EN on a non-English
    video is translation-of-auto-caption (double degradation), whereas
    Claude translating the auto-caption-of-original-audio is single degradation
    and produces better digests.
    """
    manual = info.get("subtitles") or {}
    auto = info.get("automatic_captions") or {}
    audio_lang = (info.get("language") or "").lower()

    def first_starts_with(d, prefix):
        if not prefix:
            return None
        for k in d:
            if k.lower().startswith(prefix):
                return k
        return None

    def first_any(d):
        return next(iter(d), None)

    for picker, source, is_manual in (
        (lambda d: first_starts_with(d, "en"), manual, True),
        (lambda d: first_starts_with(d, audio_lang), manual, True),
        (first_any, manual, True),
        (lambda d: first_starts_with(d, audio_lang), auto, False),
        (lambda d: first_starts_with(d, "en"), auto, False),
        (first_any, auto, False),
    ):
        k = picker(source)
        if k:
            return (k, is_manual)
    return None


def fetch_youtube(
    url: str,
    cache_root: Path,
    whisper_model: str = DEFAULT_WHISPER_MODEL,
    allow_whisper: bool = True,
    cookies_from_browser: Optional[str] = None,
) -> dict:
    """Download mp4 + best-available SRT from YouTube. Cached by video ID under cache_root.

    Returns a dict with:
      mp4 (Path), srt (Path), lang (str),
      title (str), webpage_url (str),
      download_secs (float), whisper_secs (float),
      used_whisper (bool), whisper_model (Optional[str]).

    Falls back to local Whisper transcription when YouTube has no captions.
    Set allow_whisper=False to fail fast instead of falling back.
    """
    import yt_dlp
    import time as _time

    cache_root.mkdir(parents=True, exist_ok=True)

    # YouTube increasingly requires logged-in cookies to bypass the bot challenge.
    # yt-dlp accepts `cookiesfrombrowser` as a tuple — single-item is the simplest
    # form (no profile / domain filter).
    cookie_opt: dict = {}
    if cookies_from_browser:
        cookie_opt["cookiesfrombrowser"] = (cookies_from_browser,)

    # YouTube's "n challenge" obfuscates real format URLs behind a JavaScript
    # function that yt-dlp must execute to deobfuscate. Without a JS runtime +
    # the challenge solver scripts, only thumbnail storyboards come back.
    rt = _ensure_js_runtime_available()
    yt_dlp_runtime_opt: dict = {}
    if rt is not None:
        yt_dlp_runtime_opt["js_runtimes"] = {rt: {}}
        yt_dlp_runtime_opt["remote_components"] = ["ejs:github"]

    base_opts = {**cookie_opt, **yt_dlp_runtime_opt}

    # Probe first to get the video ID for stable cache layout. We use the same
    # format selector as the download below so probe doesn't reject videos whose
    # default yt-dlp selector ('bestvideo*+bestaudio') happens to match nothing
    # (some YouTube videos return a format pool that the default doesn't span).
    probe_opts = {
        "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
        "quiet": True,
        "no_warnings": True,
        **base_opts,
    }
    with yt_dlp.YoutubeDL(probe_opts) as ydl:
        info = ydl.extract_info(url, download=False)
    video_id = info["id"]
    title = info.get("title") or video_id
    webpage_url = info.get("webpage_url") or url
    upload_date = info.get("upload_date")  # YYYYMMDD per yt-dlp; None if missing
    out_dir = cache_root / video_id
    out_dir.mkdir(parents=True, exist_ok=True)

    mp4_path = out_dir / f"{video_id}.mp4"

    # Cache hit: lang is in the filename (works for legacy *.en.srt too).
    existing_srt = next(iter(out_dir.glob(f"{video_id}.*.srt")), None)
    if mp4_path.exists() and existing_srt is not None:
        lang = existing_srt.stem[len(video_id) + 1:]
        print(f"      using cached {out_dir}/ (lang: {lang})")
        return {
            "mp4": mp4_path, "srt": existing_srt, "lang": lang,
            "title": title, "webpage_url": webpage_url, "upload_date": upload_date,
            "download_secs": 0.0, "whisper_secs": 0.0,
            "used_whisper": False, "whisper_model": None,
        }

    picked = _pick_caption_lang(info)

    if picked is not None:
        picked_lang, _is_manual = picked
        ydl_opts = {
            "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
            "merge_output_format": "mp4",
            "outtmpl": str(out_dir / f"{video_id}.%(ext)s"),
            "writesubtitles": True,
            "writeautomaticsub": True,
            "subtitleslangs": [picked_lang],
            "subtitlesformat": "srt/vtt/best",
            "postprocessors": [{"key": "FFmpegSubtitlesConvertor", "format": "srt"}],
            "quiet": True,
            "no_warnings": True,
            **base_opts,
        }
    else:
        # No captions of any kind. Download just the mp4 and transcribe locally.
        if not allow_whisper:
            raise RuntimeError(
                f"No subtitles available for {url} in any language and "
                "Whisper fallback is disabled."
            )
        print("      no captions on YouTube; will transcribe with Whisper after download")
        ydl_opts = {
            "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
            "merge_output_format": "mp4",
            "outtmpl": str(out_dir / f"{video_id}.%(ext)s"),
            "quiet": True,
            "no_warnings": True,
            **base_opts,
        }

    download_t0 = _time.monotonic()
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])
    download_secs = _time.monotonic() - download_t0

    candidates_mp4 = list(out_dir.glob(f"{video_id}.*"))
    mp4_found = next((p for p in candidates_mp4 if p.suffix in (".mp4", ".mkv", ".webm")), None)
    if mp4_found and mp4_found != mp4_path:
        mp4_found.rename(mp4_path)
    elif not mp4_path.exists():
        raise RuntimeError(f"yt-dlp finished but no video file found in {out_dir}")

    if picked is not None:
        # Re-derive lang from the filename yt-dlp actually wrote (it normalizes
        # codes like en-US -> en, so the picked code may differ from the result).
        srt_found = next(iter(out_dir.glob(f"{video_id}.*.srt")), None)
        if srt_found is None:
            raise RuntimeError(
                f"yt-dlp claimed '{picked[0]}' subtitles for {url} but produced no SRT."
            )
        lang = srt_found.stem[len(video_id) + 1:]
        return {
            "mp4": mp4_path, "srt": srt_found, "lang": lang,
            "title": title, "webpage_url": webpage_url, "upload_date": upload_date,
            "download_secs": download_secs, "whisper_secs": 0.0,
            "used_whisper": False, "whisper_model": None,
        }

    whisper_t0 = _time.monotonic()
    srt_path, lang = _transcribe_with_whisper(mp4_path, out_dir, video_id, model_name=whisper_model)
    whisper_secs = _time.monotonic() - whisper_t0
    return {
        "mp4": mp4_path, "srt": srt_path, "lang": lang,
        "title": title, "webpage_url": webpage_url, "upload_date": upload_date,
        "download_secs": download_secs, "whisper_secs": whisper_secs,
        "used_whisper": True, "whisper_model": whisper_model,
    }


# ---------- LLM backend abstraction ----------
#
# Two backends:
#   AnthropicAPIBackend   — direct anthropic.Anthropic() SDK calls; uses
#                           messages.parse for structured output and
#                           cache_control for prompt caching. Requires
#                           ANTHROPIC_API_KEY.
#   ClaudeCodeBackend     — shells out to the sandboxed `claude` binary with
#                           --output-format json (and --json-schema for parse).
#                           Uses the user's Claude.ai subscription auth as
#                           configured in our sandbox. No prompt caching
#                           (each subprocess starts cold) and vision support
#                           is opt-in (off by default per user preference).
#
# Both expose: text(...), parse(schema=...), vision_parse(content_blocks=...).
# Returns (response, usage_namespace) where usage has .input_tokens,
# .output_tokens, .cache_read_input_tokens, .cache_creation_input_tokens.
#
# Backends raise VisionUnsupported when vision is requested but unavailable;
# callers fall back to non-vision paths.

from types import SimpleNamespace as _SN


class VisionUnsupported(Exception):
    """Raised when a backend cannot process image inputs in the current config."""


def _zero_usage() -> _SN:
    return _SN(input_tokens=0, output_tokens=0,
               cache_read_input_tokens=0, cache_creation_input_tokens=0)


# ---------- Cost audit / pricing ----------
#
# Per-million-token rates in USD. Used for the cost-transparency layer
# so the user can see per-call and aggregate spend. Update when Anthropic
# adjusts pricing; users can override via the `model_pricing` setting.
# Ranges as of May 2026; treat as estimates.
#
# Cache-read is the discounted price for tokens the API serves from prompt
# cache; cache-creation is the surcharge for the FIRST time the cached
# block is seen. Anthropic publishes both; both matter for our usage shape
# (digest + panel + takeaway re-cite the transcript).

DEFAULT_MODEL_PRICING: dict = {
    # Sonnet 4.6 — digest, takeaway, vision-pick, on-demand panel re-runs
    "claude-sonnet-4-6": {
        "input": 3.0, "output": 15.0,
        "cache_read": 0.30, "cache_creation": 3.75,
    },
    # Opus 4.7 — panel discussion (highest-quality multi-perspective)
    "claude-opus-4-7": {
        "input": 15.0, "output": 75.0,
        "cache_read": 1.50, "cache_creation": 18.75,
    },
    # Haiku 4.5 — slide classifier, validation pings
    "claude-haiku-4-5-20251001": {
        "input": 0.80, "output": 4.0,
        "cache_read": 0.08, "cache_creation": 1.0,
    },
    # Older / aliases that may show up in saved settings.json:
    "claude-haiku-4-5": {
        "input": 0.80, "output": 4.0,
        "cache_read": 0.08, "cache_creation": 1.0,
    },
}


def _model_pricing(model: str) -> Optional[dict]:
    """Resolve pricing for a model name. Settings can override via
    settings['model_pricing'][model]; falls back to the default table.
    """
    try:
        s = load_settings()
        override = (s.get("model_pricing") or {}).get(model)
        if override:
            return override
    except Exception:
        pass
    return DEFAULT_MODEL_PRICING.get(model)


def estimate_cost_usd(model: str, usage) -> float:
    """Dollar estimate for a single LLM call. Returns 0.0 when pricing is
    unknown for the model — caller can choose to surface that as 'n/a'.
    Token attributes default to 0 if missing (Claude Code backend).
    """
    rates = _model_pricing(model)
    if not rates:
        return 0.0

    def _tok(name: str) -> int:
        return int(getattr(usage, name, 0) or 0)

    return (
        _tok("input_tokens") / 1_000_000 * rates.get("input", 0.0)
        + _tok("output_tokens") / 1_000_000 * rates.get("output", 0.0)
        + _tok("cache_read_input_tokens") / 1_000_000 * rates.get("cache_read", 0.0)
        + _tok("cache_creation_input_tokens") / 1_000_000 * rates.get("cache_creation", 0.0)
    )


def _llm_usage_log_path() -> Path:
    p = get_data_dir() / "logs" / "llm_usage.jsonl"
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def record_llm_usage(
    *,
    video_id: Optional[str],
    kind: str,
    model: str,
    backend_name: str,
    usage,
) -> dict:
    """Append a single usage record to ~/yt2md/logs/llm_usage.jsonl and
    return the recorded dict. Cost is set to 0.0 for the claude-code backend
    (subscription bills the user via their Anthropic plan, not per-call) so
    the audit log stays consistent — token counts still recorded for
    rate-limit awareness.
    """
    import time as _t

    cost = 0.0 if backend_name == "claude-code" else estimate_cost_usd(model, usage)
    entry = {
        "ts": _t.time(),
        "video_id": video_id or "",
        "kind": kind,
        "model": model,
        "backend": backend_name,
        "input_tokens": int(getattr(usage, "input_tokens", 0) or 0),
        "output_tokens": int(getattr(usage, "output_tokens", 0) or 0),
        "cache_read_input_tokens": int(getattr(usage, "cache_read_input_tokens", 0) or 0),
        "cache_creation_input_tokens": int(getattr(usage, "cache_creation_input_tokens", 0) or 0),
        "cost_usd": round(cost, 6),
    }
    try:
        with open(_llm_usage_log_path(), "a") as f:
            f.write(json.dumps(entry) + "\n")
    except OSError:
        pass  # Audit log must never break the digest pipeline.
    return entry


def read_llm_usage_log() -> List[dict]:
    """Load the full usage log into memory. The file stays small in
    practice (a few hundred bytes per entry, one entry per LLM call).
    """
    p = _llm_usage_log_path()
    if not p.exists():
        return []
    rows: List[dict] = []
    try:
        with open(p) as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    rows.append(json.loads(line))
                except json.JSONDecodeError:
                    continue
    except OSError:
        pass
    return rows


class AnthropicAPIBackend:
    name = "api"
    vision_supported = True

    def __init__(self):
        import anthropic
        self._client = anthropic.Anthropic()

    def text(self, *, system: str, user_text: str, model: str,
             max_tokens: int, cache: bool = False):
        block = {"type": "text", "text": user_text}
        if cache:
            block["cache_control"] = {"type": "ephemeral"}
        response = self._client.messages.create(
            model=model, max_tokens=max_tokens, system=system,
            messages=[{"role": "user", "content": [block]}],
        )
        text = next(b.text for b in response.content if b.type == "text")
        return text, response.usage

    def parse(self, *, system: str, user_text: str, model: str,
              max_tokens: int, schema, cache: bool = False):
        block = {"type": "text", "text": user_text}
        if cache:
            block["cache_control"] = {"type": "ephemeral"}
        response = self._client.messages.parse(
            model=model, max_tokens=max_tokens, system=system,
            messages=[{"role": "user", "content": [block]}],
            output_format=schema,
        )
        return response.parsed_output, response.usage

    def vision_parse(self, *, system: str, content_blocks: list,
                     model: str, max_tokens: int, schema):
        response = self._client.messages.parse(
            model=model, max_tokens=max_tokens, system=system,
            messages=[{"role": "user", "content": content_blocks}],
            output_format=schema,
        )
        return response.parsed_output, response.usage


class ClaudeCodeBackend:
    name = "claude-code"

    def __init__(self, *, vision_enabled: bool = False):
        if not claude_code_installed():
            raise RuntimeError("Claude Code is not installed in the sandbox.")
        self._binary = str(claude_binary_path())
        self._env = claude_subprocess_env()
        self._vision_enabled = vision_enabled

    @property
    def vision_supported(self) -> bool:
        return self._vision_enabled

    def _run(self, *, prompt: str, model: str, schema=None,
             timeout: float = 600.0):
        cmd = [self._binary, "-p", prompt, "--model", model,
               "--output-format", "json"]
        if schema is not None:
            cmd += ["--json-schema", json.dumps(schema.model_json_schema())]
        proc = subprocess.run(
            cmd, env=self._env, capture_output=True, text=True, timeout=timeout,
        )
        if proc.returncode != 0:
            raise RuntimeError(
                f"claude -p failed (exit {proc.returncode}): "
                f"{(proc.stderr or proc.stdout or '').strip()[:500]}"
            )
        return self._parse_output(proc.stdout)

    @staticmethod
    def _parse_output(stdout: str):
        """`claude -p --output-format json` emits a JSON envelope. Extract the
        text/result and a usage namespace. Tolerant of envelope shape changes.
        Raises RuntimeError when the envelope's `is_error` flag is set (which
        Claude Code uses for auth + tool failures even when the process exits 0).
        """
        try:
            payload = json.loads(stdout)
        except json.JSONDecodeError:
            # Defensive fallback: treat raw stdout as the result.
            return stdout.strip(), _zero_usage()
        result_text = (
            payload.get("result")
            or payload.get("text")
            or payload.get("response")
            or ""
        )
        if payload.get("is_error"):
            raise RuntimeError(f"claude reported error: {result_text or '(no message)'}")
        usage_dict = payload.get("usage") or {}
        usage = _SN(
            input_tokens=int(usage_dict.get("input_tokens") or 0),
            output_tokens=int(usage_dict.get("output_tokens") or 0),
            cache_read_input_tokens=int(usage_dict.get("cache_read_input_tokens") or 0),
            cache_creation_input_tokens=int(
                usage_dict.get("cache_creation_input_tokens") or 0
            ),
        )
        return result_text, usage

    def text(self, *, system: str, user_text: str, model: str,
             max_tokens: int, cache: bool = False):
        # No system-prompt CLI flag we rely on — combine system + user into
        # one prompt with explicit delineation. cache is ignored (subprocess
        # invocations don't share Anthropic's prompt cache).
        prompt = f"<system>\n{system}\n</system>\n\n{user_text}"
        return self._run(prompt=prompt, model=model)

    def parse(self, *, system: str, user_text: str, model: str,
              max_tokens: int, schema, cache: bool = False):
        prompt = f"<system>\n{system}\n</system>\n\n{user_text}"
        result_text, usage = self._run(prompt=prompt, model=model, schema=schema)
        # --json-schema constrains the output to the schema; the result string
        # IS the JSON we need to parse into the Pydantic model.
        try:
            data = json.loads(result_text) if isinstance(result_text, str) else result_text
        except json.JSONDecodeError as e:
            raise RuntimeError(
                f"claude returned non-JSON for a schema-constrained call: {e}\n"
                f"Output: {result_text[:500]}"
            )
        return schema.model_validate(data), usage

    def vision_parse(self, *, system: str, content_blocks: list,
                     model: str, max_tokens: int, schema):
        if not self._vision_enabled:
            raise VisionUnsupported(
                "Claude Code vision is disabled. Enable claude_code_vision in "
                "Settings to base64-embed images in prompts (token-heavy)."
            )
        # Opt-in path: serialize image blocks as base64 markers in the prompt.
        # This is expensive — each image expands ~33% over its byte size and
        # there's no native multipart in the CLI -p mode.
        text_parts: list = []
        for block in content_blocks:
            if block.get("type") == "text":
                text_parts.append(block["text"])
            elif block.get("type") == "image":
                src = block.get("source") or {}
                if src.get("type") == "base64":
                    text_parts.append(
                        f"[image media_type={src.get('media_type','image/jpeg')} "
                        f"data:base64]\n{src.get('data','')}\n[/image]"
                    )
        prompt = (
            f"<system>\n{system}\n</system>\n\n" + "\n".join(text_parts)
        )
        result_text, usage = self._run(prompt=prompt, model=model, schema=schema)
        try:
            data = json.loads(result_text) if isinstance(result_text, str) else result_text
        except json.JSONDecodeError as e:
            raise RuntimeError(
                f"claude returned non-JSON for vision call: {e}\nOutput: {result_text[:500]}"
            )
        return schema.model_validate(data), usage


def select_backend(*, vision_enabled: Optional[bool] = None):
    """Resolve the active LLM backend from settings + environment.

    Honors settings["llm_backend"] in {"auto", "api", "claude-code"}:
      - "auto":        prefer "api" when ANTHROPIC_API_KEY is set, else
                       "claude-code" when sandboxed claude is installed and
                       logged in, else raises RuntimeError.
      - "api":         requires ANTHROPIC_API_KEY.
      - "claude-code": requires sandbox install + login.
    """
    s = load_settings()
    choice = (s.get("llm_backend") or "auto").lower()
    if vision_enabled is None:
        vision_enabled = bool(s.get("claude_code_vision", False))

    if choice == "api":
        if not os.environ.get("ANTHROPIC_API_KEY"):
            raise RuntimeError(
                "llm_backend=api but ANTHROPIC_API_KEY is not set. "
                "Configure it in Settings or switch backend to claude-code."
            )
        return AnthropicAPIBackend()

    if choice == "claude-code":
        if not claude_code_installed():
            raise RuntimeError(
                "llm_backend=claude-code but Claude Code is not installed. "
                "Run /setup to install it."
            )
        return ClaudeCodeBackend(vision_enabled=vision_enabled)

    # auto: prefer API when key is set (keeps prompt caching + native vision);
    # else fall through to Claude Code only when both installed AND we have a
    # cached login signal (sentinel). Otherwise raise so the caller redirects
    # to /setup rather than burning a doomed call.
    if os.environ.get("ANTHROPIC_API_KEY"):
        return AnthropicAPIBackend()
    if (claude_code_installed()
            and _claude_code_session_state.get("logged_in", False)):
        return ClaudeCodeBackend(vision_enabled=vision_enabled)
    raise RuntimeError(
        "No LLM backend configured. Set ANTHROPIC_API_KEY or install + sign "
        "in to Claude Code via /setup."
    )


# ---------- Markdown digest (transcript-primary, LLM-summarized) ----------

DIGEST_SYSTEM_PROMPT = (
    "You distill video transcripts into readable digests so a reader can grasp "
    "the essence of a video without watching it.\n\n"
    "Given a timestamped transcript, identify the natural topic segments. "
    "Return 5–12 sections depending on length and content density (favor fewer "
    "for short videos, more for long talks).\n\n"
    "For each topic:\n"
    "- title: a short, descriptive heading (not a question, not a teaser)\n"
    "- start_time: the timestamp in SECONDS where the topic begins, taken from "
    "the transcript's bracketed timestamps\n"
    "- summary: 2–4 sentences of informative prose distilling what is actually "
    "said. Concrete claims, names, numbers, and reasoning — not vague paraphrases.\n"
    "- key_points: 3–6 short bullet points capturing the most useful takeaways. "
    "Bullets should add detail beyond the summary, not restate it.\n\n"
    "Also produce a 2–3 sentence overview of the entire video.\n\n"
    "Write for a reader, not a viewer. Skip filler ('in this video I'll show you'); "
    "go straight to the substance."
)


def generate_digest(
    segments: List[TranscriptSegment],
    video_title: str,
    model: str,
    source_lang: str = "en",
    output_language: str = "auto",
    backend=None,
):
    """Call Claude to segment the transcript into topics. Returns a parsed VideoDigest.

    source_lang is the BCP-47 language code of the transcript (e.g. 'en', 'zh-Hans').
    output_language: 'auto' (write in source language) or 'en' (force English).
    backend: an LLMBackend; defaults to select_backend() (auto-resolved).
    """
    from pydantic import BaseModel
    from typing import List as TList

    class Topic(BaseModel):
        title: str
        start_time: float
        summary: str
        key_points: TList[str]

    class VideoDigest(BaseModel):
        title: str
        overview: str
        topics: TList[Topic]

    transcript = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    is_english_source = (source_lang or "").lower().startswith("en")
    lang_note = ""
    if not is_english_source:
        if output_language == "en":
            lang_note = (
                f"NOTE: This transcript is in language code '{source_lang}', not English. "
                "Translate to English while distilling. Title, overview, topic titles, "
                "summaries, and key points must all be written in English regardless of "
                "the source language. Preserve proper nouns (people, places, products) in "
                "their original form when there is no established English rendering.\n\n"
            )
        else:  # "auto" — match the source language
            lang_note = (
                f"NOTE: This transcript is in language code '{source_lang}'. Write the "
                "digest in the SAME language as the transcript. Title, overview, topic "
                "titles, summaries, and key points must all be in the source language. "
                "Preserve proper nouns and established technical terms in their original "
                "form (including English technical jargon when the field uses it that way).\n\n"
            )

    user_text = (
        f"{lang_note}"
        f"Video title: {video_title}\n\n"
        f"Total duration: {format_timestamp(segments[-1].end if segments else 0)}\n\n"
        f"Timestamped transcript:\n\n{transcript}"
    )

    if backend is None:
        backend = select_backend()
    return backend.parse(
        system=DIGEST_SYSTEM_PROMPT, user_text=user_text,
        model=model, max_tokens=16000, schema=VideoDigest, cache=True,
    )


DEFAULT_PANEL_MODEL = "claude-opus-4-7"


PANEL_SYSTEM_PROMPT = (
    "You facilitate a panel of domain experts critically analyzing video content. "
    "Read the digest and transcript carefully, then:\n\n"
    "1. Infer 3–5 experts whose perspectives would best illuminate this material. "
    "Choose them from the actual domain of the video — a neuroscientist for a brain "
    "talk, a hardware engineer + an ML practitioner for a chip-design talk, a historian "
    "of science + a contemporary researcher for a science-history piece. Avoid generic "
    "labels (\"a thoughtful generalist\"); make each expert's specialty concrete enough "
    "that their angle on this material is distinct.\n\n"
    "2. Run a 1500–2500 word panel discussion in markdown. Open with one short paragraph "
    "introducing each panelist (name, role, one credential or claim-to-relevance). Then "
    "the discussion proper, with each turn labeled by the speaker's name.\n\n"
    "Goals for the discussion:\n"
    "- Surface what the speaker glossed over, hand-waved, or assumed without arguing.\n"
    "- Bring contrary readings — where would a competing school of thought disagree?\n"
    "- Connect to adjacent domains the speaker didn't mention.\n"
    "- Examine concrete claims (numbers, names, mechanisms) for how robust they actually are.\n"
    "- Synthesize, but don't paper over disagreements: if two panelists land in "
    "different places, leave them there.\n\n"
    "Style: skip restating the digest — the reader already read it. Open directly with "
    "the moderator framing the first question. No conclusion-summary at the end; let "
    "the discussion close naturally."
)


def generate_panel_discussion(
    digest_md_text: str,
    segments: List["TranscriptSegment"],
    model: str,
    source_lang: str = "en",
    output_language: str = "auto",
    backend=None,
):
    """Call Claude to simulate a panel of domain-relevant experts discussing a video.
    Returns (markdown_text, usage).

    source_lang / output_language follow the same convention as generate_digest:
    'auto' writes the panel in the transcript's language; 'en' forces English.

    Costs ~one Opus call per click (≈ 4–8k input + 2–4k output tokens). Output is
    one markdown document the caller writes to digests/<id>/panel.md.
    """
    transcript_str = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    is_english_source = (source_lang or "").lower().startswith("en")
    lang_directive = ""
    if not is_english_source and output_language == "auto":
        lang_directive = (
            f"\n\nIMPORTANT: The transcript is in language code '{source_lang}'. "
            "Write the entire panel discussion in the SAME language — expert names "
            "(transliterated when appropriate), credentials, the moderator's "
            "questions, every speaker's turns. Preserve proper nouns and technical "
            "terms in their original form when the field uses them that way."
        )
    # output_language == "en" with non-English source: rely on the existing
    # system prompt (no explicit translate directive needed; English is the
    # default Claude output style for this prompt).

    user_text = (
        "## Existing digest (the reader has already seen this)\n\n"
        f"{digest_md_text}\n\n"
        "## Full timestamped transcript\n\n"
        f"{transcript_str}"
        f"{lang_directive}\n\n"
        "Now: introduce the panelists, then run the discussion."
    )

    if backend is None:
        backend = select_backend()
    return backend.text(
        system=PANEL_SYSTEM_PROMPT, user_text=user_text,
        model=model, max_tokens=8000, cache=True,
    )


DEFAULT_TAKEAWAY_MODEL = "claude-sonnet-4-6"


TAKEAWAY_SYSTEM_PROMPT = (
    "You write the audience-facing takeaway for a video the reader has just "
    "finished. They've read the digest and (often) the panel discussion. "
    "Now they want a friend's-eye-view: 'Here's what I got out of this; "
    "here's what to walk away with; here's what's contested.'\n\n"
    "Format: 1–3 short paragraphs of plain prose. Open with the single most "
    "important thing the reader should leave with — the bottom line, "
    "stated directly. Then weave in supporting context: where the speaker's "
    "framing is solid, where the panel pushed back, what's still open. "
    "Close with the implications for the reader (so what / why does this "
    "matter), one or two sentences.\n\n"
    "Genre awareness — first identify (silently, internally) what kind of "
    "video this is, and shape the takeaway accordingly:\n"
    "- Tech talk / explainer → the working position on the frameworks "
    "presented, with critical pushback woven in.\n"
    "- Market / finance → the trade thesis, what's priced in, the catalysts "
    "to watch — with an explicit 'as of <date>' anchor.\n"
    "- News / current events → what changed and what it means.\n"
    "- How-to / tutorial → what they teach you to do, and the gotcha "
    "experienced practitioners flag.\n"
    "- Interview → where the speakers' positions actually differ.\n"
    "- Product launch → what's genuinely new, what's hype, what to "
    "actually use.\n"
    "Do NOT label the genre in the output. Just let it shape the writing.\n\n"
    "Grounding claims in the video — when you state a specific fact, claim, "
    "or quote that lives at a particular moment, mark it with bracketed "
    "timestamps using the original video's M:SS or H:MM:SS format, e.g. "
    "[3:15] or [1:02:48]. Use the bracketed form exactly — no parentheses, "
    "no markdown link syntax. The renderer will turn these into clickable "
    "links to the source video. Use them sparingly (3–8 across the whole "
    "takeaway) and only on substantive points worth verifying — not on "
    "every sentence.\n\n"
    "Time-sensitive content: if the video discusses dated material — market "
    "state, recent product launches, current numbers — mention the publish "
    "date inline ('As of <date>...') so a future reader knows the freshness "
    "window. Use the publish date provided in the user message. For "
    "evergreen content (general knowledge, frameworks, well-established "
    "claims), don't add a date.\n\n"
    "Style:\n"
    "- Conversational, not academic. A friend telling you what they got out "
    "of the talk — not a research abstract.\n"
    "- Concrete > abstract. Use real names, numbers, and frameworks from "
    "the video.\n"
    "- Honest about contested points. If the panel disagreed, say so "
    "(\"though the panel pushed back on the file-system framing as a long-"
    "term abstraction\"). Don't paper it over.\n"
    "- Don't restate the digest. The reader just read it. Synthesize and "
    "go beyond.\n"
    "- No headings, no bullet lists, no preamble like 'Here's my takeaway'. "
    "Just the prose."
)


def generate_takeaway(
    digest_md_text: str,
    panel_md_text: Optional[str],
    segments: List["TranscriptSegment"],
    model: str,
    *,
    publish_date: Optional[str] = None,
    source_lang: str = "en",
    output_language: str = "auto",
    backend=None,
):
    """Final pipeline step: write the audience-facing takeaway as 1-3 short
    paragraphs of prose. Synthesizes digest + panel into a personal
    'what to walk away with' read.

    publish_date: YYYYMMDD or YYYY-MM-DD as returned by yt-dlp's
    `info["upload_date"]`; threaded into the prompt so time-sensitive
    takeaways can anchor with 'as of <date>'.

    Returns (takeaway_text: str, usage). The text contains [M:SS] bracket
    markers that the renderer converts into clickable timestamp links.
    """
    transcript_str = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    # Normalize publish_date to YYYY-MM-DD for the prompt (yt-dlp returns
    # YYYYMMDD by default).
    pub_str = ""
    if publish_date:
        pd = publish_date.replace("-", "")
        if len(pd) == 8 and pd.isdigit():
            pub_str = f"{pd[0:4]}-{pd[4:6]}-{pd[6:8]}"
        else:
            pub_str = publish_date

    is_english_source = (source_lang or "").lower().startswith("en")
    lang_directive = ""
    if not is_english_source and output_language == "auto":
        lang_directive = (
            f"\n\nIMPORTANT: The transcript is in language code '{source_lang}'. "
            "Write the takeaway in the SAME language as the transcript. "
            "Preserve proper nouns and technical terms in their original form."
        )

    panel_section = (
        "## Panel discussion\n\n" + panel_md_text + "\n\n"
        if panel_md_text else
        "## Panel discussion\n\n(none generated)\n\n"
    )

    user_text = (
        f"Video publish date: {pub_str or '(unknown)'}\n\n"
        "## Digest (the reader has seen this)\n\n"
        f"{digest_md_text}\n\n"
        f"{panel_section}"
        "## Full timestamped transcript\n\n"
        f"{transcript_str}"
        f"{lang_directive}\n\n"
        "Now: write the takeaway."
    )

    if backend is None:
        backend = select_backend()
    return backend.text(
        system=TAKEAWAY_SYSTEM_PROMPT, user_text=user_text,
        model=model, max_tokens=2000, cache=True,
    )



def _transcript_slice(
    segments: List["TranscriptSegment"],
    topic_start: float,
    topic_end: float,
) -> str:
    """Render the transcript segments inside [topic_start, topic_end) as one timestamped string.

    Used by vision_pick_frames to ground picks against what the narrator is saying
    at a candidate frame's timestamp. Truncates very long topics to keep token cost
    bounded — first 30 + last 30 segments with a marker between, which captures the
    narrator's framing at start and conclusion at end.
    """
    in_window = [s for s in segments if topic_start <= s.start < topic_end]
    if len(in_window) > 70:
        head = in_window[:30]
        tail = in_window[-30:]
        in_window = head + [None] + tail  # type: ignore[list-item]
    lines: List[str] = []
    for seg in in_window:
        if seg is None:
            lines.append("[…]")
            continue
        lines.append(f"[{format_timestamp(seg.start)}] {seg.text}")
    return "\n".join(lines)


def _candidates_for_topic(
    topic_start: float,
    topic_end: float,
    frames: List[Tuple[Path, float]],
    max_per_topic: int = 5,
    overlap_pre: float = 5.0,
) -> List[Tuple[Path, float]]:
    """Frames whose timestamp falls in [start - overlap, end). Downsample to max_per_topic."""
    in_window = [
        (p, t) for p, t in frames if (topic_start - overlap_pre) <= t < topic_end
    ]
    if len(in_window) <= max_per_topic:
        return in_window
    # Even spacing across the window
    step = len(in_window) / max_per_topic
    return [in_window[int(i * step)] for i in range(max_per_topic)]


def _encode_frame_for_vision(path: Path, max_long_edge: int = 1024) -> str:
    """Resize a frame to <= max_long_edge on the long side and return base64-encoded JPEG bytes."""
    import base64
    import io
    from PIL import Image

    with Image.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        scale = max_long_edge / max(w, h)
        if scale < 1:
            im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=82)
    return base64.standard_b64encode(buf.getvalue()).decode("ascii")


def vision_pick_frames(
    digest,
    frames: List[Tuple[Path, float]],
    video_duration: float,
    model: str,
    segments: Optional[List["TranscriptSegment"]] = None,
    backend=None,
):
    """Use Claude's vision to pick the best frame per topic from in-window candidates.

    Returns a dict {topic_index -> chosen_frame_path}, plus the API usage object.
    Topics with no in-window candidates are omitted (caller falls back to timestamp-based pick).

    If `segments` is provided, the per-topic transcript slice is included so vision
    can ground picks on what the narrator is saying at each candidate's timestamp
    (e.g. "speaker says 'as you can see in this diagram' at 04:23 → frame at 04:23").

    Raises VisionUnsupported when the active backend can't process images
    (e.g. Claude Code with vision opt-out). Caller should fall back.
    """
    from pydantic import BaseModel
    from typing import List as TList

    class TopicChoice(BaseModel):
        topic_index: int
        candidate_index: int
        rationale: str

    class FrameChoices(BaseModel):
        choices: TList[TopicChoice]

    topics = digest.topics

    # Build per-topic candidate lists and a flat list of (topic_idx, cand_idx, path, ts)
    per_topic: List[List[Tuple[Path, float]]] = []
    per_topic_transcript: List[str] = []
    for i, topic in enumerate(topics):
        end = topics[i + 1].start_time if i + 1 < len(topics) else video_duration
        per_topic.append(_candidates_for_topic(topic.start_time, end, frames))
        if segments:
            per_topic_transcript.append(
                _transcript_slice(segments, topic.start_time, end)
            )
        else:
            per_topic_transcript.append("")

    # Build the message: text intro -> for each topic, label + summary + numbered candidate images
    content: list = []
    intro = (
        "For each topic below, pick the candidate frame that best illustrates what the "
        "narrator is discussing. Prefer frames showing the most informative visual content "
        "(diagrams, code, distinctive UI) over generic framing or talking-head shots.\n\n"
        "When multiple candidates show the same scene at different stages of an animation "
        "or progressive reveal — bullets appearing one at a time, diagram elements being "
        "added, code typed line by line — prefer the LATEST candidate in the sequence. "
        "The final state shows the most complete information; partial/early states omit "
        "content the narrator goes on to add.\n\n"
        "Use the per-topic transcript to ground your pick: when the narrator says things "
        "like \"as you can see here\" or refers to a specific element at a specific "
        "moment, prefer the candidate whose timestamp is closest to that mention.\n\n"
        "Return one choice per topic that has candidates.\n\n"
        f"Total topics: {len(topics)}\n"
    )
    content.append({"type": "text", "text": intro})

    for ti, topic in enumerate(topics):
        cands = per_topic[ti]
        if not cands:
            content.append({
                "type": "text",
                "text": f"\n--- Topic {ti} (no candidates available — skip) ---\n"
                        f"Title: {topic.title}\n",
            })
            continue
        header_parts = [
            f"\n--- Topic {ti} ---",
            f"Title: {topic.title}",
            f"Summary: {topic.summary}",
        ]
        if per_topic_transcript[ti]:
            header_parts.append("Transcript:\n" + per_topic_transcript[ti])
        header_parts.append(f"Candidates ({len(cands)} frames):\n")
        content.append({"type": "text", "text": "\n".join(header_parts)})
        for ci, (path, ts) in enumerate(cands):
            content.append({
                "type": "text",
                "text": f"Candidate {ci} (at {format_timestamp(ts)}):",
            })
            content.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/jpeg",
                    "data": _encode_frame_for_vision(path),
                },
            })

    system = (
        "You select illustrative frames for a video digest. You will be shown a list of "
        "topics, each with the topic's title, summary, the transcript spoken during that "
        "topic (timestamped), and a small set of candidate frames (also timestamped). "
        "For each topic with candidates, return the (topic_index, candidate_index) pair "
        "that best illustrates the topic, with a one-sentence rationale. Use the "
        "transcript to ground your choice on what the narrator is saying when each "
        "candidate frame was captured. Skip topics that say 'no candidates available'."
    )

    if backend is None:
        backend = select_backend()
    if not getattr(backend, "vision_supported", False):
        raise VisionUnsupported(
            f"Backend {backend.name!r} does not support vision in the current "
            "configuration."
        )
    parsed, usage = backend.vision_parse(
        system=system, content_blocks=content,
        model=model, max_tokens=4000, schema=FrameChoices,
    )

    chosen: dict = {}
    for choice in parsed.choices:
        if 0 <= choice.topic_index < len(topics):
            cands = per_topic[choice.topic_index]
            if 0 <= choice.candidate_index < len(cands):
                chosen[choice.topic_index] = cands[choice.candidate_index][0]
    return chosen, usage


def _pick_topic_frame(
    topic_start: float,
    topic_end: float,
    candidates: List[Tuple[Path, float]],
    used: set,
) -> Optional[Tuple[Path, float]]:
    """Pick the best frame for a topic: prefer a frame inside [start, end), else closest."""
    in_window = [(p, t) for p, t in candidates if topic_start <= t < topic_end and p not in used]
    if in_window:
        midpoint = (topic_start + min(topic_end, in_window[-1][1] + 1)) / 2
        return min(in_window, key=lambda x: abs(x[1] - midpoint))
    available = [(p, t) for p, t in candidates if p not in used]
    if not available:
        return None
    return min(available, key=lambda x: abs(x[1] - topic_start))


def write_markdown_digest(
    digest,
    candidate_frames: List[Tuple[Path, float]],
    video_duration: float,
    output_md: Path,
    images_dir: Path,
    vision_picks: Optional[dict] = None,
    video_title: Optional[str] = None,
    video_url: Optional[str] = None,
) -> None:
    """Render the digest as Markdown with <img> tags. Copies frames into images_dir.

    If vision_picks is provided, prefer those mappings; fall back to timestamp-based picks
    for any topic not covered.

    video_title (when provided) overrides the LLM-generated title so the digest's
    heading matches the original YouTube title — readers can map it back to the
    source. video_url renders a "Watch on YouTube" link directly under the title.
    """
    images_dir.mkdir(parents=True, exist_ok=True)

    used: set = set()
    topic_images: List[Optional[Path]] = []
    topics = digest.topics
    for i, topic in enumerate(topics):
        # Vision pick takes priority
        if vision_picks and i in vision_picks:
            src_path = vision_picks[i]
            used.add(src_path)
        else:
            end = topics[i + 1].start_time if i + 1 < len(topics) else video_duration
            pick = _pick_topic_frame(topic.start_time, end, candidate_frames, used)
            if pick is None:
                topic_images.append(None)
                continue
            src_path, _ = pick
            used.add(src_path)
        dest = images_dir / f"topic_{i + 1:02d}.jpg"
        shutil.copy(src_path, dest)
        topic_images.append(dest)

    rel_dir = images_dir.name  # render as "images/topic_NN.jpg"
    lines: List[str] = []
    heading = video_title if video_title else digest.title
    lines.append(f"# {heading}")
    lines.append("")
    if video_url:
        lines.append(f"**Watch on YouTube:** <{video_url}>")
        lines.append("")
    lines.append(digest.overview)
    lines.append("")
    for i, topic in enumerate(topics):
        ts = format_timestamp(topic.start_time)
        # Deep-link the timestamp when we know the source URL: clicking
        # jumps the viewer straight to that moment in the original video.
        # Without a URL (local-file digests), keep the bare timestamp.
        if video_url:
            sec = int(topic.start_time)
            sep = "&" if "?" in video_url else "?"
            ts_link = f"[{ts}]({video_url}{sep}t={sec}s)"
            lines.append(f"## {topic.title}  <sub>*{ts_link}*</sub>")
        else:
            lines.append(f"## {topic.title}  <sub>*{ts}*</sub>")
        lines.append("")
        if topic_images[i] is not None:
            # Use the topic title as alt text — accessible to screen readers and
            # readable as a fallback if the image fails to load.
            alt = topic.title.replace('"', "'")
            lines.append(f'<img src="{rel_dir}/{topic_images[i].name}" alt="{alt}" width="800">')
            lines.append("")
        lines.append(topic.summary)
        lines.append("")
        for kp in topic.key_points:
            lines.append(f"- {kp}")
        if topic.key_points:
            lines.append("")

    output_md.write_text("\n".join(lines).rstrip() + "\n")


def render_takeaway_markdown(
    takeaway_text: str,
    *,
    video_url: Optional[str] = None,
) -> str:
    """Post-process the LLM's takeaway prose for writing to takeaway.md.
    Converts the LLM's bracketed [M:SS] / [H:MM:SS] markers into clickable
    markdown links to the source video when video_url is known. Returns the
    body text alone — the file's heading is rendered by the viewer chrome,
    not embedded in the markdown.
    """
    body = takeaway_text.strip()

    if not video_url:
        return body + "\n"

    sep = "&" if "?" in video_url else "?"

    def _ts_to_seconds(ts: str) -> int:
        parts = [int(p) for p in ts.split(":")]
        if len(parts) == 2:
            return parts[0] * 60 + parts[1]
        if len(parts) == 3:
            return parts[0] * 3600 + parts[1] * 60 + parts[2]
        return 0

    # Match bare bracketed timestamps like [3:15] or [1:02:48], but not ones
    # that are already a markdown link's display text — skip patterns
    # immediately followed by '(' which indicates the LLM already wrote a
    # full [text](url) link.
    def _link(m):
        ts = m.group(1)
        sec = _ts_to_seconds(ts)
        return f"[{ts}]({video_url}{sep}t={sec}s)"

    return re.sub(r"\[(\d+:\d+(?::\d+)?)\](?!\()", _link, body) + "\n"


# ---------- Subcommands: watch / meta / serve ----------

LATEST_LIMIT = 10
MAX_NEW_PER_RUN = 3


def channels_file() -> Path:
    return get_data_dir() / "channels.txt"


def state_file() -> Path:
    return get_data_dir() / "state.json"


def read_channels() -> List[str]:
    p = channels_file()
    if not p.exists():
        return []
    return [
        line.strip()
        for line in p.read_text().splitlines()
        if line.strip() and not line.lstrip().startswith("#")
    ]


def write_channels(channels: List[str]) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    body = "# YouTube channels to watch. One URL per line. Lines starting with # are ignored.\n"
    body += "\n".join(channels) + ("\n" if channels else "")
    channels_file().write_text(body)


def load_state() -> dict:
    p = state_file()
    if not p.exists():
        return {"channels": {}}
    return json.loads(p.read_text())


def save_state(state: dict) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    state_file().write_text(json.dumps(state, indent=2, sort_keys=True) + "\n")


# ---- watch subcommands ----

def normalize_channel_url(url: str) -> str:
    """Light normalization for YouTube channel URLs.

    Accepts: '@handle', 'youtube.com/@handle', 'https://www.youtube.com/@handle/videos'.
    Always returns a fully-qualified URL. Adds '/videos' to bare @handle URLs so
    yt-dlp targets the videos tab specifically.
    """
    url = url.strip()
    if url.startswith("@"):
        url = f"https://www.youtube.com/{url}"
    if not url.startswith(("http://", "https://")):
        url = f"https://{url}"
    m = re.match(r"^(https?://(?:www\.)?youtube\.com/@[^/]+)/?$", url)
    if m:
        url = m.group(1) + "/videos"
    return url


def cmd_watch_add(args) -> int:
    url = normalize_channel_url(args.url)
    if not is_url(url):
        sys.exit(f"Not a URL: {url}")
    channels = read_channels()
    if url in channels:
        print(f"Already watching: {url}")
        return 0
    channels.append(url)
    write_channels(channels)
    print(f"Added: {url}")
    return 0


def cmd_watch_list(args) -> int:
    channels = read_channels()
    if not channels:
        print("No channels configured. Add one with: yt2md watch add <URL>")
        return 0
    print(f"Watching {len(channels)} channel(s):")
    for ch in channels:
        print(f"  {ch}")
    print(f"\nConfig: {channels_file()}")
    print(f"State:  {state_file()}")
    print(f"Data:   {get_data_dir()}")
    return 0


def cmd_watch_remove(args) -> int:
    url = args.url.strip()
    channels = read_channels()
    if url not in channels:
        sys.exit(f"Not in list: {url}")
    channels = [c for c in channels if c != url]
    write_channels(channels)
    print(f"Removed: {url}")
    return 0


def _list_channel_videos(url: str, limit: int = LATEST_LIMIT) -> List[str]:
    out = subprocess.check_output(
        ["yt-dlp", "--flat-playlist", "--playlist-end", str(limit), "--print", "%(id)s", url],
        text=True,
    )
    return [line.strip() for line in out.splitlines() if line.strip()]


def _digest_video(video_id: str, output_dir: Path) -> Tuple[int, str]:
    """Run yt2md on a video. Returns (exit_code, combined_stdout_stderr).

    Streams output to the parent's stdout in real time (so poll.log captures
    it as it happens) AND collects it into a buffer the caller can scan for
    permanent-failure patterns.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    digest_path = output_dir / "digest.md"
    yt2md = shutil.which("yt2md") or sys.argv[0]
    proc = subprocess.Popen(
        [yt2md, f"https://youtu.be/{video_id}", "-o", str(digest_path)],
        cwd=output_dir,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
    )
    chunks: list = []
    assert proc.stdout is not None
    for line in proc.stdout:
        sys.stdout.write(line)
        sys.stdout.flush()
        chunks.append(line)
    proc.wait()
    return proc.returncode, "".join(chunks)


# Substrings (case-insensitive) that indicate a video can never be digested by
# this account — gating it forever, deleted, etc. Matches mark the video
# as seen so polling stops re-trying it. Network/transient errors don't match
# and continue to retry next cycle.
_PERMANENT_FAILURE_PATTERNS = (
    "members-only",
    "join this channel",
    "private video",
    "video unavailable",
    "this video is no longer available",
    "removed by the uploader",
    "removed for violating",
    "sign in to confirm your age",
    "video has been removed",
)


def _is_permanent_failure(output: str) -> bool:
    low = output.lower()
    return any(p in low for p in _PERMANENT_FAILURE_PATTERNS)


def cmd_watch_run(args) -> int:
    ensure_api_key()
    channels = read_channels()
    if not channels:
        print("No channels configured. Add one with: yt2md watch add <URL>")
        return 0

    data_dir = get_data_dir()
    digests_dir = data_dir / "digests"
    state = load_state()
    any_failures = False

    for channel_url in channels:
        print(f"--- {channel_url}")
        seen = set(state["channels"].get(channel_url, {}).get("seen", []))
        latest_ids = _list_channel_videos(channel_url)

        if not seen:
            print(f"  first run, seeding state with {len(latest_ids)} videos (no backfill)")
            state["channels"][channel_url] = {"seen": sorted(latest_ids)}
            continue

        new_ids = [vid for vid in latest_ids if vid not in seen][:MAX_NEW_PER_RUN]
        if not new_ids:
            print("  no new videos")
            continue

        print(f"  {len(new_ids)} new: {new_ids}")
        for vid in reversed(new_ids):
            print(f"  processing {vid}...")
            rc, output = _digest_video(vid, digests_dir / vid)
            if rc == 0:
                seen.add(vid)
                state["channels"][channel_url] = {"seen": sorted(seen)}
                save_state(state)
            elif _is_permanent_failure(output):
                # Permanent: mark seen so polling stops cycling on it.
                # Wipe the partial dir (mp4 download, empty digest, etc.) to
                # keep digests/ tidy.
                print(f"  PERMANENTLY UNAVAILABLE: {vid} — marking seen and wiping partial dir")
                shutil.rmtree(digests_dir / vid, ignore_errors=True)
                seen.add(vid)
                state["channels"][channel_url] = {"seen": sorted(seen)}
                save_state(state)
            else:
                print(f"  FAILED on {vid} (transient — will retry next poll)", file=sys.stderr)
                any_failures = True

    save_state(state)
    return 1 if any_failures else 0


# ---- in-process scheduler ----
#
# Cadenced background runner for `yt2md watch run` (subscription poll). Lives
# inside `yt2md serve`: a daemon thread ticks every ~30s, fires due jobs as
# detached subprocesses, tracks pid + exit code in schedule_state.json for
# the /schedule UI.
#
# Tradeoff: scheduling pauses while serve is down — for an interactively-used
# reader this is fine; missed slots fire on next start (catch-up semantics).

DEFAULT_SCHEDULE_CONFIG = {
    "poll_interval_hours": 6,
}

_SCHED_TICK_SECS = 30
_scheduler_jobs: dict = {}
_scheduler_thread = None
_scheduler_lock_obj = None


def _schedule_lock():
    global _scheduler_lock_obj
    if _scheduler_lock_obj is None:
        import threading
        _scheduler_lock_obj = threading.Lock()
    return _scheduler_lock_obj


def _schedule_config_file() -> Path:
    return get_data_dir() / "schedule.json"


def _schedule_state_file() -> Path:
    return get_data_dir() / "schedule_state.json"


def load_schedule_config() -> dict:
    p = _schedule_config_file()
    if not p.exists():
        return dict(DEFAULT_SCHEDULE_CONFIG)
    try:
        cfg = json.loads(p.read_text())
        merged = dict(DEFAULT_SCHEDULE_CONFIG)
        merged.update({k: v for k, v in cfg.items() if k in DEFAULT_SCHEDULE_CONFIG})
        return merged
    except Exception:
        return dict(DEFAULT_SCHEDULE_CONFIG)


def save_schedule_config(cfg: dict) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    _schedule_config_file().write_text(json.dumps(cfg, indent=2) + "\n")


def _load_schedule_state() -> dict:
    p = _schedule_state_file()
    default = {"poll": {}}
    if not p.exists():
        return default
    try:
        s = json.loads(p.read_text())
        return {"poll": s.get("poll") or {}}
    except Exception:
        return default


def _save_schedule_state(state: dict) -> None:
    p = _schedule_state_file()
    p.parent.mkdir(parents=True, exist_ok=True)
    tmp = p.with_suffix(".json.tmp")
    tmp.write_text(json.dumps(state, indent=2) + "\n")
    tmp.replace(p)


# ---- user-tunable settings (model + tooling choices) ----
#
# Lives at ~/yt2md/settings.json. Editable via the /settings page; flows into
# every subprocess spawn (one-off, scheduled poll) as YT2MD_* env vars, so the
# digest CLI's argparse defaults pick them up. The /digests/<id>/discuss route
# reads settings directly (it runs in-process).

DEFAULT_SETTINGS = {
    "digest_model": "claude-sonnet-4-6",
    "panel_model": "claude-opus-4-7",
    "whisper_model": DEFAULT_WHISPER_MODEL,
    "cookies_from_browser": "",
    # "auto" = write the digest in the same language as the transcript;
    # "en" = always translate to English. Applies to both the per-video digest
    # and the panel discussion.
    "digest_language": "auto",
    # Which auth path to use for LLM calls.
    #   "auto":        prefer ANTHROPIC_API_KEY when set; else use the bundled
    #                  Claude Code sandbox if installed + logged in.
    #   "api":         force direct Anthropic API (requires ANTHROPIC_API_KEY).
    #   "claude-code": force the bundled Claude Code subprocess backend.
    "llm_backend": "auto",
    # Vision frame-picking is automatic for the API backend (cheap & native)
    # but disabled by default for Claude Code (no -p image flag; we'd have to
    # base64-embed → token-heavy). Toggle on if you want vision via Claude
    # Code despite the cost.
    "claude_code_vision": False,
    # When generating slides, use a vision-LLM (typically Haiku) to filter
    # raw extracted frames down to actual deck slides. Trades a small LLM
    # cost (~$0.005 per video via 3×3 grid batching with Haiku) for a much
    # cleaner deck. Set False to use pure pHash dedup only.
    "slide_classification": True,
    "slide_classifier_model": "claude-haiku-4-5-20251001",
}


def _settings_file() -> Path:
    return get_data_dir() / "settings.json"


def load_settings() -> dict:
    p = _settings_file()
    if not p.exists():
        return dict(DEFAULT_SETTINGS)
    try:
        s = json.loads(p.read_text())
        merged = dict(DEFAULT_SETTINGS)
        merged.update({k: v for k, v in s.items() if k in DEFAULT_SETTINGS})
        return merged
    except Exception:
        return dict(DEFAULT_SETTINGS)


def save_settings(s: dict) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    tmp = _settings_file().with_suffix(".json.tmp")
    tmp.write_text(json.dumps(s, indent=2) + "\n")
    tmp.replace(_settings_file())


def _settings_to_env(settings: dict) -> dict:
    """Map settings to YT2MD_* env vars for subprocess invocation. Empty values
    are dropped so the subprocess sees only meaningful overrides."""
    out: dict = {}
    if settings.get("digest_model"):
        out["YT2MD_DIGEST_MODEL"] = settings["digest_model"]
    if settings.get("whisper_model"):
        out["YT2MD_WHISPER_MODEL"] = settings["whisper_model"]
    if settings.get("cookies_from_browser"):
        out["YT2MD_COOKIES_FROM_BROWSER"] = settings["cookies_from_browser"]
    if settings.get("panel_model"):
        out["YT2MD_PANEL_MODEL"] = settings["panel_model"]
    if settings.get("digest_language"):
        out["YT2MD_DIGEST_LANGUAGE"] = settings["digest_language"]
    # Spawned subprocesses re-resolve llm_backend / claude_code_vision via
    # load_settings() in the child — no env-var passthrough needed for those.
    # CLAUDE_CONFIG_DIR is required so the bundled `claude` binary in the
    # child finds our sandboxed credentials/settings instead of any system
    # install's defaults.
    out["CLAUDE_CONFIG_DIR"] = str(claude_config_dir())
    return out


def _format_schedule_summary(cfg: dict) -> str:
    """Human-readable description of the schedule config."""
    poll = cfg["poll_interval_hours"]
    poll_str = f"every {poll} hour{'s' if poll != 1 else ''}"
    return f"polling {poll_str}"


def _compute_next_poll(cfg: dict, last_started_at: Optional[float]) -> float:
    """Timestamp of the next scheduled poll. First-run convention: fire ASAP."""
    interval = max(60.0, float(cfg.get("poll_interval_hours", 6)) * 3600.0)
    if last_started_at is None:
        import time as _t
        return _t.time()
    return last_started_at + interval


def _fire_scheduled_job(kind: str) -> Optional[subprocess.Popen]:
    """Spawn yt2md as a subprocess for the given kind. Returns the running
    Popen or None if already running / yt2md not on PATH. Caller holds lock."""
    if kind != "poll":
        return None
    existing = _scheduler_jobs.get(kind)
    if existing is not None and existing.poll() is None:
        return existing
    yt2md_path = shutil.which("yt2md")
    if not yt2md_path:
        print(f"[scheduler] yt2md not on PATH; cannot fire {kind}", file=sys.stderr)
        return None
    args = [yt2md_path, "watch", "run"]
    log_path = get_data_dir() / "logs" / f"{kind}.log"
    log_path.parent.mkdir(parents=True, exist_ok=True)
    import time as _t
    log_fd = open(log_path, "a")
    log_fd.write(f"\n===== {_t.strftime('%Y-%m-%d %H:%M:%S')} {kind} run =====\n")
    log_fd.flush()
    proc = subprocess.Popen(
        args,
        stdout=log_fd,
        stderr=subprocess.STDOUT,
        env={**os.environ, **_settings_to_env(load_settings()), "PYTHONUNBUFFERED": "1"},
        start_new_session=True,
    )
    log_fd.close()
    _scheduler_jobs[kind] = proc

    state = _load_schedule_state()
    state[kind] = {
        **(state.get(kind) or {}),
        "last_started_at": _t.time(),
        "last_pid": proc.pid,
        "last_exit_code": None,
        "last_finished_at": None,
    }
    _save_schedule_state(state)
    return proc


def _reap_scheduled_job(kind: str) -> None:
    """Capture exit code if the kind's subprocess has exited. Caller holds lock."""
    proc = _scheduler_jobs.get(kind)
    if proc is None:
        return
    rc = proc.poll()
    if rc is None:
        return
    import time as _t
    state = _load_schedule_state()
    state[kind] = {
        **(state.get(kind) or {}),
        "last_finished_at": _t.time(),
        "last_exit_code": rc,
    }
    _save_schedule_state(state)
    del _scheduler_jobs[kind]


def _scheduler_tick() -> None:
    with _schedule_lock():
        _reap_scheduled_job("poll")
        cfg = load_schedule_config()
        state = _load_schedule_state()
        import time as _t
        now = _t.time()
        if "poll" not in _scheduler_jobs:
            if now >= _compute_next_poll(cfg, (state.get("poll") or {}).get("last_started_at")):
                _fire_scheduled_job("poll")


def _scheduler_loop() -> None:
    import time as _t
    while True:
        try:
            _scheduler_tick()
        except Exception as e:
            print(f"[scheduler] tick error: {e}", file=sys.stderr)
        _t.sleep(_SCHED_TICK_SECS)


def start_scheduler() -> None:
    """Start the daemon scheduler thread. Idempotent."""
    global _scheduler_thread
    if _scheduler_thread is not None and _scheduler_thread.is_alive():
        return
    import threading
    _scheduler_thread = threading.Thread(target=_scheduler_loop, daemon=True)
    _scheduler_thread.start()
    print("[scheduler] started (in-process; ticks every "
          f"{_SCHED_TICK_SECS}s)")


def _cleanup_legacy_launchd() -> None:
    """Best-effort one-time removal of the old launchctl plists from
    ~/Library/LaunchAgents — so the user doesn't end up with duplicate
    scheduling after this migration. Silent if nothing is present."""
    launchd_dir = Path.home() / "Library" / "LaunchAgents"
    removed = []
    for label in ("com.youtube-to-markdown.poll", "com.youtube-to-markdown.meta"):
        plist = launchd_dir / f"{label}.plist"
        if plist.exists():
            subprocess.run(
                ["launchctl", "bootout", f"gui/{os.getuid()}", str(plist)],
                check=False, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            )
            try:
                plist.unlink()
                removed.append(label)
            except OSError:
                pass
    if removed:
        print(f"[scheduler] removed legacy launchd plists: {', '.join(removed)}")


def _scheduler_status_summary(kind: str, state: dict) -> Tuple[str, str]:
    """English summary + dot class for /schedule and /channels surfaces."""
    s = state.get(kind) or {}
    if kind in _scheduler_jobs:
        return ("Running now.", "dot-on")
    last_started = s.get("last_started_at")
    last_exit = s.get("last_exit_code")
    if last_started is None:
        return ("Set up — first run hasn't happened yet.", "dot-warn")
    if last_exit not in (0, None):
        return (f"Last run failed (exit code {last_exit}). Check the log.", "dot-warn")
    import datetime as dt
    age = dt.datetime.now() - dt.datetime.fromtimestamp(last_started)
    age_secs = int(age.total_seconds())
    if age_secs < 60:
        age_str = f"{age_secs}s ago"
    elif age_secs < 3600:
        age_str = f"{age_secs // 60}m ago"
    elif age_secs < 86400:
        age_str = f"{age_secs // 3600}h ago"
    else:
        age_str = f"{age_secs // 86400}d ago"
    return (f"Healthy — last run {age_str} (clean).", "dot-on")


def _format_next_run(ts: float) -> str:
    """Render an upcoming-run timestamp as 'in 2h 15m' / 'in 3 days' / 'overdue'."""
    import time as _t
    delta = ts - _t.time()
    if delta < 0:
        return "due now"
    if delta < 60:
        return f"in {int(delta)}s"
    if delta < 3600:
        return f"in {int(delta // 60)}m"
    if delta < 86400:
        h = int(delta // 3600)
        m = int((delta % 3600) // 60)
        return f"in {h}h {m}m"
    return f"in {int(delta // 86400)}d"


def _tail_log(path: Path, n: int = 20) -> str:
    if not path.exists():
        return "(no log yet)"
    try:
        text = path.read_text(errors="replace")
    except Exception as e:
        return f"(error reading log: {e})"
    lines = text.splitlines()
    return "\n".join(lines[-n:]) if lines else "(empty)"




# ---- one-off digest jobs (in-memory tracking; cheap) ----

# Module-level dict: PID -> {"video_id": str, "started": float, "url": str, "proc": Popen}.
# Lost on server restart, which is fine — the digest still completes (detached
# subprocess) and shows up in the sidebar when done.
_oneoff_jobs: dict = {}

# Recent failures, most-recent-first, capped at _ONEOFF_FAILURE_CAP.
# Lost on server restart (matches _oneoff_jobs lifecycle).
_oneoff_failures: list = []
_ONEOFF_FAILURE_CAP = 20


_VIDEO_ID_RE = re.compile(r"(?:v=|youtu\.be/|/shorts/|/embed/)([A-Za-z0-9_-]{11})")


def extract_video_id(url: str) -> str:
    """Pull a YouTube video ID from common URL forms. Returns '' if not found."""
    m = _VIDEO_ID_RE.search(url)
    return m.group(1) if m else ""


def _extract_last_error(log_path: Path, video_id: str) -> str:
    """Pull the most relevant error line from oneoff.log for a given video_id.

    The log uses '===== {ts} starting {video_id} ({url}) =====' as section
    markers. We bound the section, then prefer 'RuntimeError: ...' / similar
    summary lines over the bare 'Traceback' header.
    """
    try:
        text = log_path.read_text(errors="replace")
    except OSError:
        return ""
    marker = f"starting {video_id}"
    idx = text.rfind(marker)
    if idx < 0:
        return ""
    next_idx = text.find("\n===== ", idx + len(marker))
    section = text[idx:next_idx if next_idx > 0 else len(text)]
    candidates = [
        ln.strip() for ln in section.splitlines()
        if ln.strip()
        and not ln.lstrip().startswith("[download")
        and ("Error" in ln or "Traceback" in ln)
    ]
    for ln in reversed(candidates):
        if ":" in ln and not ln.startswith("Traceback"):
            return ln
    return candidates[-1] if candidates else ""


# Pipeline stage markers, in pipeline order. The status endpoint scans the log
# section for the LAST matching substring to determine the current stage.
# Reordering here changes which stage is reported when two markers happen to
# match — keep this list in pipeline order (later entries override earlier).
_ONEOFF_STAGE_MARKERS: list = [
    ("starting",            "starting "),  # section header is always present
    ("downloading",         "[0/5] Fetching YouTube video"),
    ("downloading",         "[download]"),
    ("loading whisper",     "loading whisper model"),
    ("transcribing",        "transcribing audio with whisper"),
    ("extracting frames",   "[1/5] Extracting frames"),
    ("deduping frames",     "[2/5] Deduping"),
    ("parsing transcript",  "[3/5] Parsing SRT"),
    ("aligning",            "[4/5] Aligning"),
    ("building deck",       "[5/5] Building deck"),
    ("digesting",           "[+] Generating digest"),
    ("vision pass",         "[+] Vision-picking"),
    ("writing digest",      "Digest written"),
]


def _describe_job_stage(log_text: str, video_id: str) -> str:
    """Return a short human-readable label for the latest stage of a job.

    Scans the log section bounded by the start marker for video_id (or EOF /
    next start marker) and returns the most pipeline-advanced stage whose
    substring marker appears in that section.
    """
    marker = f"starting {video_id}"
    idx = log_text.rfind(marker)
    if idx < 0:
        return "starting"
    next_idx = log_text.find("\n===== ", idx + len(marker))
    section = log_text[idx:next_idx if next_idx > 0 else len(log_text)]
    current = "starting"
    for label, needle in _ONEOFF_STAGE_MARKERS:
        if needle in section:
            current = label
    return current


def _extract_run_summary(log_text: str, video_id: str) -> Optional[dict]:
    """Pull the last `[summary] {...}` JSON line emitted by the pipeline for a job.

    The pipeline prints one line of the form `[summary] {...}` on successful
    completion. Returns the parsed dict, or None if absent / malformed.
    """
    import json as _json
    marker = f"starting {video_id}"
    idx = log_text.rfind(marker)
    if idx < 0:
        return None
    next_idx = log_text.find("\n===== ", idx + len(marker))
    section = log_text[idx:next_idx if next_idx > 0 else len(log_text)]
    last = None
    for ln in section.splitlines():
        s = ln.lstrip()
        if s.startswith("[summary] "):
            last = s[len("[summary] "):]
    if not last:
        return None
    try:
        return _json.loads(last)
    except Exception:
        return None


def _runs_jsonl_path() -> Path:
    return get_data_dir() / "logs" / "runs.jsonl"


def _record_run(row: dict) -> None:
    """Persist a run completion to library.db AND append a JSONL line.

    Two stores by design: SQLite for the activity UI's queries; JSONL for
    `tail -f` debugging and downstream scripts. Both reflect the same data.
    """
    import json as _json
    cols = (
        "video_id", "url", "source", "started_at", "ended_at", "duration_secs",
        "exit_code", "success", "stage_reached", "error", "source_lang",
        "used_whisper", "whisper_model",
        "download_secs", "whisper_secs", "frames_secs", "digest_secs", "vision_secs",
        "digest_input_tokens", "digest_output_tokens",
        "digest_cache_read_tokens", "digest_cache_creation_tokens",
        "digest_path",
    )
    placeholders = ", ".join("?" for _ in cols)
    values = tuple(row.get(c) for c in cols)
    try:
        with _library_connect() as conn:
            conn.execute(
                f"INSERT INTO runs ({', '.join(cols)}) VALUES ({placeholders})",
                values,
            )
    except Exception as e:
        # DB failure shouldn't bring down the reaper. JSONL still gets written.
        print(f"[runs] db insert failed: {e}", file=sys.stderr)

    jsonl = _runs_jsonl_path()
    try:
        jsonl.parent.mkdir(parents=True, exist_ok=True)
        with jsonl.open("a") as fh:
            fh.write(_json.dumps(row) + "\n")
    except OSError as e:
        print(f"[runs] jsonl append failed: {e}", file=sys.stderr)


def _recent_runs(limit: int = 100) -> list:
    """Read the last `limit` runs from SQLite, newest first."""
    try:
        with _library_connect() as conn:
            conn.row_factory = __import__("sqlite3").Row
            rows = conn.execute(
                "SELECT * FROM runs ORDER BY started_at DESC LIMIT ?", (limit,)
            ).fetchall()
            return [dict(r) for r in rows]
    except Exception:
        return []


def _build_run_row(info: dict, exit_code: int, log_text: str) -> dict:
    """Translate a finished job + its log into a row dict suitable for runs table."""
    import time as _t
    video_id = info["video_id"]
    started = info["started"]
    ended = _t.time()
    summary = _extract_run_summary(log_text, video_id)
    success = exit_code == 0 and summary is not None
    error = "" if success else _extract_last_error(_oneoff_log_path(), video_id)
    stage = _describe_job_stage(log_text, video_id)
    timings = (summary or {}).get("timings") or {}
    tokens = (summary or {}).get("tokens") or {}
    return {
        "video_id": video_id,
        "url": info.get("url"),
        "source": "oneoff",
        "started_at": started,
        "ended_at": ended,
        "duration_secs": ended - started,
        "exit_code": exit_code,
        "success": 1 if success else 0,
        "stage_reached": stage,
        "error": error or None,
        "source_lang": (summary or {}).get("source_lang"),
        "used_whisper": 1 if (summary or {}).get("used_whisper") else 0,
        "whisper_model": (summary or {}).get("whisper_model"),
        "download_secs": timings.get("download"),
        "whisper_secs": timings.get("whisper"),
        "frames_secs": timings.get("frames"),
        "digest_secs": timings.get("digest"),
        "vision_secs": timings.get("vision"),
        "digest_input_tokens": tokens.get("input"),
        "digest_output_tokens": tokens.get("output"),
        "digest_cache_read_tokens": tokens.get("cache_read"),
        "digest_cache_creation_tokens": tokens.get("cache_creation"),
        "digest_path": (summary or {}).get("digest_path"),
    }


def _oneoff_log_path() -> Path:
    return get_data_dir() / "logs" / "oneoff.log"


def _record_oneoff_failure(info: dict, exit_code: int) -> None:
    import time as _t
    log_path = _oneoff_log_path()
    last_error = _extract_last_error(log_path, info["video_id"]) if log_path.exists() else ""
    _oneoff_failures.insert(0, {
        "video_id": info["video_id"],
        "url": info["url"],
        "exit_code": exit_code,
        "started": info["started"],
        "ended": _t.time(),
        "error": last_error,
    })
    del _oneoff_failures[_ONEOFF_FAILURE_CAP:]


def _list_active_oneoff_jobs() -> list:
    """Return one-off jobs whose subprocesses are still alive.

    Side effect: jobs that have exited are removed from _oneoff_jobs;
    non-zero exits are appended to _oneoff_failures.
    """
    active = []
    for pid in list(_oneoff_jobs.keys()):
        info = _oneoff_jobs[pid]
        proc = info.get("proc")
        if proc is None:
            # legacy entries with no Popen handle — fall back to kill probe
            try:
                os.kill(pid, 0)
                active.append({"pid": pid, **{k: v for k, v in info.items() if k != "proc"}})
            except (ProcessLookupError, PermissionError):
                del _oneoff_jobs[pid]
            continue
        rc = proc.poll()
        if rc is None:
            active.append({"pid": pid, **{k: v for k, v in info.items() if k != "proc"}})
            continue
        del _oneoff_jobs[pid]
        # Treat any non-zero exit, or zero-exit-with-no-digest, as a failure.
        digest_path = get_data_dir() / "digests" / info["video_id"] / "digest.md"
        try:
            log_text = _oneoff_log_path().read_text(errors="replace")
        except OSError:
            log_text = ""
        if rc != 0 or not digest_path.exists():
            _record_oneoff_failure(info, rc)
        _record_run(_build_run_row(info, rc, log_text))
    return active


def _list_recent_oneoff_failures() -> list:
    """Return a copy of recent failures (most recent first)."""
    return list(_oneoff_failures)


# ---- read-state library (SQLite-backed) ----

def _library_path() -> Path:
    return get_data_dir() / "library.db"


def _library_connect():
    """Open (and lazily migrate) the read-state SQLite database."""
    import sqlite3
    get_data_dir().mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(_library_path())
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS digest_reads (
            digest_id TEXT PRIMARY KEY,
            opened_at INTEGER NOT NULL
        );
        CREATE TABLE IF NOT EXISTS runs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            video_id TEXT NOT NULL,
            url TEXT,
            source TEXT NOT NULL,           -- 'oneoff' | 'poll' | 'meta'
            started_at REAL NOT NULL,
            ended_at REAL NOT NULL,
            duration_secs REAL NOT NULL,
            exit_code INTEGER NOT NULL,
            success INTEGER NOT NULL,
            stage_reached TEXT,
            error TEXT,
            source_lang TEXT,
            used_whisper INTEGER NOT NULL DEFAULT 0,
            whisper_model TEXT,
            download_secs REAL,
            whisper_secs REAL,
            frames_secs REAL,
            digest_secs REAL,
            vision_secs REAL,
            digest_input_tokens INTEGER,
            digest_output_tokens INTEGER,
            digest_cache_read_tokens INTEGER,
            digest_cache_creation_tokens INTEGER,
            digest_path TEXT
        );
        CREATE INDEX IF NOT EXISTS idx_runs_started ON runs(started_at DESC);
    """)
    return conn


def _mark_digest_read(digest_id: str) -> None:
    import time
    with _library_connect() as conn:
        conn.execute(
            "INSERT OR REPLACE INTO digest_reads(digest_id, opened_at) VALUES (?, ?)",
            (digest_id, int(time.time())),
        )


def _read_digest_ids() -> set:
    with _library_connect() as conn:
        rows = conn.execute("SELECT digest_id FROM digest_reads").fetchall()
        return {r[0] for r in rows}


# ---- serve subcommand (local reader UI) ----

# Shared between digest + panel + takeaway viewers. The button has
# data-copy-target pointing at a hidden <textarea> elsewhere on the page; on
# click we read its raw value and copy to the clipboard. Falls back to a
# select+execCommand path for browsers / contexts (eg http on a non-localhost)
# where the modern clipboard API is unavailable.
#
# Optional attribute data-then-open="<url>" navigates to the URL in a new tab
# after a successful copy — used by the "Continue in chat" handoff to land
# the user in claude.ai with the context already in their clipboard.
_COPY_BUTTON_JS = """
<script>
(function () {
  function flash(btn, text) {
    const orig = btn.getAttribute('data-orig-text') || btn.textContent;
    btn.setAttribute('data-orig-text', orig);
    btn.textContent = text;
    setTimeout(() => { btn.textContent = orig; }, 1500);
  }
  function maybeOpen(btn) {
    const url = btn.getAttribute('data-then-open');
    if (!url) return;
    // Brief delay so the user sees the "Copied!" feedback before the new tab.
    setTimeout(() => window.open(url, '_blank', 'noopener'), 600);
  }
  document.addEventListener('click', async (ev) => {
    const btn = ev.target.closest('[data-copy-target]');
    if (!btn) return;
    const id = btn.getAttribute('data-copy-target');
    const src = document.getElementById(id);
    if (!src) return;
    const text = src.value;
    try {
      if (navigator.clipboard && navigator.clipboard.writeText) {
        await navigator.clipboard.writeText(text);
        flash(btn, 'Copied!');
        maybeOpen(btn);
        return;
      }
    } catch (e) { /* fall through to legacy path */ }
    try {
      src.removeAttribute('hidden');
      src.style.position = 'absolute';
      src.style.left = '-9999px';
      src.select();
      document.execCommand('copy');
      src.setAttribute('hidden', '');
      flash(btn, 'Copied!');
      maybeOpen(btn);
    } catch (e) {
      flash(btn, 'Copy failed');
    }
  });
})();
</script>
"""


# ---- shared async job tracker for in-process work surfaced via the web UI ----
#
# In-process daemon threads keyed by "<video_id>:<kind>" — currently used for
# slides generation; designed to also serve panel + takeaway retrofits later.
# Threads die with the Flask process; the work is idempotent so a lost job
# just means the user clicks Generate again. State lives in a module dict
# (single-process Flask, single-writer assumption holds).

_local_jobs: dict = {}


def start_local_job(key: str, fn, *args, **kwargs) -> bool:
    """Start a daemon-thread job under `key` if one isn't already running.
    Returns True if a new job was started, False if one was already in flight.
    Idempotent: a second click while running is a no-op (not an error).
    """
    import threading
    import time as _t

    existing = _local_jobs.get(key)
    if existing and existing["thread"].is_alive():
        return False

    job = {
        "started": _t.time(),
        "kind": key.split(":", 1)[-1],
        "error": None,
        "thread": None,
    }

    def _wrapper():
        try:
            fn(*args, **kwargs)
        except Exception as e:
            job["error"] = f"{type(e).__name__}: {e}"

    t = threading.Thread(target=_wrapper, daemon=True)
    job["thread"] = t
    _local_jobs[key] = job
    t.start()
    return True


def local_job_status(key: str) -> dict:
    """Snapshot of a job's state. Includes a simple 'phase' string the UI can
    render directly: 'idle' (no job ever ran or completed long ago) /
    'running' / 'done' / 'error'.
    """
    import time as _t
    job = _local_jobs.get(key)
    if not job:
        return {"phase": "idle"}
    running = job["thread"].is_alive()
    elapsed = int(_t.time() - job["started"])
    if running:
        return {"phase": "running", "elapsed": elapsed,
                "started": job["started"]}
    if job.get("error"):
        return {"phase": "error", "elapsed": elapsed,
                "error": job["error"]}
    return {"phase": "done", "elapsed": elapsed}


# Job-status polling JS shared by any toolbar element with
# data-poll-url=<url> — polls every 2s, updates a `.elapsed` child with
# "{n}s", and reloads the page when the artifact appears (on success) or
# surfaces the error inline (on failure).
_JOB_POLL_JS = """
<script>
(function () {
  const el = document.querySelector('[data-poll-url]');
  if (!el) return;
  const url = el.getAttribute('data-poll-url');
  const elapsedSpan = el.querySelector('.elapsed');
  let timer = null;
  async function tick() {
    try {
      const res = await fetch(url);
      if (!res.ok) return;
      const s = await res.json();
      if (s.phase === 'running') {
        if (elapsedSpan) elapsedSpan.textContent = s.elapsed + 's';
        return;
      }
      if (timer) { clearInterval(timer); timer = null; }
      if (s.phase === 'done' && s.artifact_exists) {
        // Reload so the toolbar swaps to "Download slides".
        window.location.href = window.location.pathname + '?msg=Slides+ready.';
        return;
      }
      if (s.phase === 'error') {
        el.innerHTML = '<span style="color: #c00;">Generation failed: ' +
          (s.error || '(unknown error)') + '</span>';
        return;
      }
      // Done but artifact missing — likely a write race. Reload anyway.
      window.location.reload();
    } catch (e) { /* network blip — try next tick */ }
  }
  tick();
  timer = setInterval(tick, 2000);
})();
</script>
"""


def build_chat_handoff_prompt(video_id: str, digests_dir: Path) -> str:
    """Assemble a prompt the user can paste into claude.ai (or any chat) to
    continue thinking about a video. Includes whichever artifacts exist on
    disk — digest, panel, takeaway — so the chat has full context without
    a manual paste of each.
    """
    digest_md = digests_dir / video_id / "digest.md"
    panel_md = digests_dir / video_id / "panel.md"
    takeaway_md = digests_dir / video_id / "takeaway.md"

    have = []
    if digest_md.exists():
        have.append("digest")
    if panel_md.exists():
        have.append("panel-of-experts critique")
    if takeaway_md.exists():
        have.append("bottom-line takeaway")
    if not have:
        artifacts_phrase = "summary"
    elif len(have) == 1:
        artifacts_phrase = have[0]
    elif len(have) == 2:
        artifacts_phrase = f"{have[0]} and {have[1]}"
    else:
        artifacts_phrase = ", ".join(have[:-1]) + ", and " + have[-1]

    parts: List[str] = [
        f"I just read a distilled summary of a YouTube video. Below is the "
        f"{artifacts_phrase}. I have a follow-up question.",
        "",
        f"Source: https://www.youtube.com/watch?v={video_id}",
        "",
    ]
    if digest_md.exists():
        parts += ["# Digest", "", digest_md.read_text().rstrip(), ""]
    if panel_md.exists():
        parts += ["# Panel discussion", "", panel_md.read_text().rstrip(), ""]
    if takeaway_md.exists():
        parts += ["# Takeaway", "", takeaway_md.read_text().rstrip(), ""]
    parts += ["---", "", "My question: "]
    return "\n".join(parts)


SERVE_PAGE_TEMPLATE = """<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{{ title }} — yt2md</title>
{% if base_href %}<base href="{{ base_href }}">{% endif %}
<script>
(function() {
  const stored = localStorage.getItem('yt2md-theme');
  if (stored && stored !== 'auto') document.documentElement.setAttribute('data-theme', stored);
})();
</script>
<style>
:root {
  --bg: #fafaf7;
  --fg: #1a1a1a;
  --muted: #6b6b6b;
  --accent: #b65a2c;
  --unread: #2563eb;
  --border: #e5e3dc;
  --sidebar-bg: #f0eee6;
  --code-bg: #ececea;
}
@media (prefers-color-scheme: dark) {
  :root:not([data-theme="light"]) {
    --bg: #1a1a1a;
    --fg: #e8e8e8;
    --muted: #999;
    --accent: #d97a4d;
    --unread: #60a5fa;
    --border: #2e2e2e;
    --sidebar-bg: #141414;
    --code-bg: #232323;
  }
}
:root[data-theme="dark"] {
  --bg: #1a1a1a;
  --fg: #e8e8e8;
  --muted: #999;
  --accent: #d97a4d;
  --unread: #60a5fa;
  --border: #2e2e2e;
  --sidebar-bg: #141414;
  --code-bg: #232323;
}
* { box-sizing: border-box; }
html, body { margin: 0; padding: 0; height: 100%; }
body {
  display: flex;
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Helvetica Neue", sans-serif;
  background: var(--bg);
  color: var(--fg);
  line-height: 1.6;
}
aside {
  width: 320px;
  flex-shrink: 0;
  height: 100vh;
  overflow-y: auto;
  padding: 20px 16px;
  background: var(--sidebar-bg);
  border-right: 1px solid var(--border);
}
aside h1 { margin: 0 0 16px; font-size: 18px; }
aside h1 a { color: var(--fg); text-decoration: none; }
aside h2 {
  font-size: 11px;
  text-transform: uppercase;
  letter-spacing: 0.08em;
  color: var(--muted);
  margin: 20px 0 8px;
  font-weight: 600;
}
aside ul { list-style: none; padding: 0; margin: 0; }
aside li { margin: 0 0 3px; }
aside li a {
  display: -webkit-box;
  -webkit-line-clamp: 3;
  -webkit-box-orient: vertical;
  overflow: hidden;
  padding: 6px 8px;
  border-radius: 4px;
  color: var(--fg);
  text-decoration: none;
  font-size: 13px;
  line-height: 1.35;
}
aside li a:hover { background: rgba(0,0,0,0.05); }
@media (prefers-color-scheme: dark) {
  aside li a:hover { background: rgba(255,255,255,0.05); }
}
aside li.active a { background: var(--accent); color: white; }
aside li.unread a { font-weight: 600; }
aside .empty { color: var(--muted); font-size: 13px; padding: 6px 8px; }
aside .unread-count {
  display: inline-block; padding: 2px 7px; background: var(--unread); color: white;
  border-radius: 10px; font-size: 11px; font-weight: 600;
  text-transform: none; letter-spacing: 0; margin-left: 4px;
}
aside .unread-dot {
  display: inline-block; width: 6px; height: 6px; border-radius: 50%;
  background: var(--unread); margin-right: 6px; vertical-align: middle;
}
aside .meta-card.unread { border-color: var(--unread); }
aside .meta-card.unread .week { font-weight: 700; }
/* When an item is the currently-viewed one, suppress the unread signals
   to avoid two competing color cues. */
aside li.active .unread-dot,
aside .meta-card.active .unread-dot { display: none; }
aside .meta-card.active.unread { border-color: var(--accent); }
aside .meta-card {
  display: block; padding: 10px 12px; border-radius: 4px;
  border: 1px solid var(--border); margin-bottom: 6px;
  text-decoration: none; color: var(--fg);
}
aside .meta-card:hover { border-color: var(--accent); }
aside .meta-card.active { background: var(--accent); color: white; border-color: var(--accent); }
aside .meta-card .week { font-weight: 600; font-size: 13px; line-height: 1.2; }
aside .meta-card .count {
  color: var(--muted); font-size: 11px; margin-top: 2px;
  text-transform: uppercase; letter-spacing: 0.04em;
}
aside .meta-card.active .count { color: white; opacity: 0.85; }
main {
  flex: 1;
  height: 100vh;
  overflow-y: auto;
  padding: 32px 48px 80px;
}
main .reader {
  max-width: 720px;
  margin: 0 auto;
}
main h1 { font-size: 28px; line-height: 1.25; margin-top: 0; }
main h2 { font-size: 22px; line-height: 1.3; margin-top: 32px; border-bottom: 1px solid var(--border); padding-bottom: 4px; }
main h3 { font-size: 17px; }
main img { max-width: 100%; height: auto; border: 1px solid var(--border); border-radius: 4px; }
main a { color: var(--accent); }
main code { background: var(--code-bg); padding: 2px 5px; border-radius: 3px; font-size: 0.9em; }
main pre { background: var(--code-bg); padding: 12px; border-radius: 4px; overflow-x: auto; }
main pre code { background: none; padding: 0; }
main blockquote { border-left: 3px solid var(--border); padding-left: 16px; color: var(--muted); margin-left: 0; }
main ul, main ol { padding-left: 24px; }
main hr { border: none; border-top: 1px solid var(--border); margin: 24px 0; }
main sub { color: var(--muted); font-size: 0.85em; }
.empty-state { color: var(--muted); margin-top: 80px; text-align: center; }
.meta-info { color: var(--muted); font-size: 13px; margin-top: -12px; margin-bottom: 24px; }
.featured-eyebrow {
  color: var(--muted); font-size: 12px; text-transform: uppercase;
  letter-spacing: 0.06em; margin-bottom: 24px;
}
.featured-eyebrow a { color: var(--accent); text-decoration: none; font-weight: 600; }
.cta {
  display: inline-block; background: var(--accent); color: white;
  padding: 12px 20px; border-radius: 4px; text-decoration: none;
  font-weight: 500; margin-top: 12px;
}
.cta:hover { opacity: 0.9; }
.recent-list { list-style: none; padding: 0; margin: 12px 0; }
.recent-list li { padding: 6px 0; border-bottom: 1px solid var(--border); }
.recent-list a { color: var(--fg); text-decoration: none; }
.recent-list a:hover { color: var(--accent); }
.add-form { display: flex; gap: 8px; margin: 24px 0; }
.add-form input[type="text"] {
  flex: 1; padding: 10px 12px; font-size: 14px;
  border: 1px solid var(--border); border-radius: 4px;
  background: var(--bg); color: var(--fg);
}
.add-form button, .channel-list button {
  padding: 10px 16px; font-size: 14px;
  border: 1px solid var(--accent); border-radius: 4px;
  background: var(--accent); color: white; cursor: pointer;
}
.channel-list { list-style: none; padding: 0; }
.channel-list li {
  display: flex; align-items: center; justify-content: space-between;
  padding: 10px 12px; border: 1px solid var(--border); border-radius: 4px;
  margin-bottom: 8px; gap: 12px;
}
.channel-list .url { flex: 1; word-break: break-all; font-size: 14px; }
.channel-list button {
  background: transparent; color: var(--muted);
  border-color: var(--border); padding: 6px 12px; font-size: 12px;
}
.channel-list button:hover { color: var(--accent); border-color: var(--accent); }
.flash { padding: 10px 14px; border-radius: 4px; margin: 16px 0;
  background: var(--code-bg); border-left: 3px solid var(--accent); }
.next-step {
  padding: 14px 18px; border-radius: 6px; margin: 16px 0 24px;
  background: var(--sidebar-bg); border: 1px solid var(--border);
  border-left: 3px solid var(--accent); font-size: 14px; line-height: 1.5;
}
.next-step strong { color: var(--accent); }
.next-step a { color: var(--accent); }
.schedule-form {
  background: var(--sidebar-bg); border: 1px solid var(--border);
  border-radius: 6px; padding: 16px 20px; margin: 16px 0;
}
.schedule-fields {
  display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
  gap: 12px 20px; margin-bottom: 16px;
}
.schedule-fields label {
  display: flex; flex-direction: column; gap: 4px;
  font-size: 12px; color: var(--muted); text-transform: uppercase;
  letter-spacing: 0.04em;
}
.schedule-fields input, .schedule-fields select {
  padding: 8px 10px; font-size: 14px; border: 1px solid var(--border);
  border-radius: 4px; background: var(--bg); color: var(--fg);
  font-family: inherit; text-transform: none; letter-spacing: normal;
}
.schedule-fields .suffix {
  position: absolute; right: 10px; top: 50%; transform: translateY(-50%);
  color: var(--muted); font-size: 12px; pointer-events: none;
}
.status-table { width: 100%; border-collapse: collapse; font-size: 13px;
  margin: 12px 0 16px; }
.status-table td { padding: 6px 12px; border-bottom: 1px solid var(--border); }
.status-table td:first-child { color: var(--muted); width: 140px; }
.status-table td:last-child { font-family: ui-monospace, "SF Mono", Menlo, monospace; }
.activity-table { width: 100%; border-collapse: collapse; font-size: 13px; margin-top: 12px; }
.activity-table th, .activity-table td {
  text-align: left; padding: 10px 8px; border-bottom: 1px solid var(--border);
  vertical-align: top;
}
.activity-table th { color: var(--muted); font-weight: 600; font-size: 11px;
  text-transform: uppercase; letter-spacing: 0.04em; }
.activity-table td a { color: var(--accent); text-decoration: none; }
.activity-table td a:hover { text-decoration: underline; }
.activity-table .ok { color: #4a9f56; font-weight: 600; }
.activity-table .fail { color: #d04545; font-weight: 600; }
.activity-meta { color: var(--muted); font-size: 12px; margin-top: 3px; }
.activity-error { font-family: ui-monospace, "SF Mono", Menlo, monospace; word-break: break-word; }
.activity-stages, .activity-tokens { font-family: ui-monospace, "SF Mono", Menlo, monospace;
  font-size: 12px; white-space: nowrap; }
.filter-row { display: flex; gap: 8px; margin: 12px 0 4px; flex-wrap: wrap; }
.filter-chip {
  padding: 5px 12px; border: 1px solid var(--border); border-radius: 999px;
  font-size: 12px; color: var(--fg); text-decoration: none; background: var(--bg);
}
.filter-chip:hover { border-color: var(--accent); }
.filter-chip.active { background: var(--accent); color: white; border-color: var(--accent); }
.filter-chip-count { opacity: 0.7; }
.delete-btn {
  padding: 6px 14px; font-size: 13px; cursor: pointer;
  background: transparent; color: #b13030;
  border: 1px solid #b13030; border-radius: 4px;
}
.delete-btn:hover { background: #b13030; color: white; }
.digest-actions { display: flex; gap: 12px; align-items: center; flex-wrap: wrap; }
.digest-toolbar { margin-top: -8px; margin-bottom: 8px; }
.digest-toolbar .delete-btn { margin-left: auto; }  /* push delete to the far right */
.discuss-btn {
  padding: 8px 16px; font-size: 13px; cursor: pointer;
  background: var(--accent); color: white;
  border: 1px solid var(--accent); border-radius: 4px;
  text-decoration: none; display: inline-block;
}
.discuss-btn:hover { opacity: 0.9; }
.discuss-btn-secondary {
  padding: 6px 12px; font-size: 12px; cursor: pointer;
  background: transparent; color: var(--fg);
  border: 1px solid var(--border); border-radius: 4px;
}
.discuss-btn-secondary:hover { border-color: var(--accent); }
.job-block { border: 1px solid var(--border); border-radius: 4px; padding: 16px 20px;
  margin: 16px 0; background: var(--sidebar-bg); }
.job-block h3 { margin: 0 0 8px; font-size: 15px; }
.job-actions { display: flex; gap: 8px; margin: 8px 0 16px; flex-wrap: wrap; }
.job-actions button {
  padding: 8px 14px; font-size: 13px; cursor: pointer;
  border: 1px solid var(--border); border-radius: 4px;
  background: var(--bg); color: var(--fg);
}
.job-actions button.primary { background: var(--accent); color: white; border-color: var(--accent); }
.job-actions button:hover { border-color: var(--accent); }
.job-summary { font-size: 15px; margin: 0 0 12px; }
details summary {
  cursor: pointer; color: var(--muted); font-size: 12px;
  text-transform: uppercase; letter-spacing: 0.04em;
  padding: 4px 0; user-select: none;
}
details summary:hover { color: var(--accent); }
details[open] summary { margin-bottom: 8px; }
.log-block {
  background: var(--code-bg); padding: 12px; border-radius: 4px;
  font-family: ui-monospace, "SF Mono", Menlo, monospace;
  font-size: 12px; line-height: 1.4; max-height: 240px;
  overflow: auto; white-space: pre-wrap; word-break: break-all;
}
.dot { display: inline-block; width: 8px; height: 8px; border-radius: 50%;
  margin-right: 6px; vertical-align: middle; }
.dot-on { background: #4caf50; }
.dot-off { background: #999; }
.dot-warn { background: #d97a4d; }

/* Accessibility: skip to content for keyboard users */
.skip-link {
  position: absolute; left: -1000px; top: 0; padding: 8px 12px;
  background: var(--accent); color: white; text-decoration: none;
  border-radius: 0 0 4px 0; z-index: 100;
}
.skip-link:focus { left: 0; }

/* Form labels (visually hidden, exposed to screen readers) */
.sr-only {
  position: absolute; width: 1px; height: 1px; padding: 0; margin: -1px;
  overflow: hidden; clip: rect(0,0,0,0); white-space: nowrap; border: 0;
}

/* Sidebar header (yt2md title + theme toggle inline) */
.sidebar-header {
  display: flex; justify-content: space-between; align-items: center;
  margin-bottom: 16px;
}
.sidebar-header h1 { margin: 0; }
.theme-toggle {
  background: transparent; border: 1px solid var(--border); border-radius: 4px;
  padding: 4px 8px; font-size: 14px; cursor: pointer; color: var(--fg);
  font-family: inherit; line-height: 1;
}
.theme-toggle:hover { border-color: var(--accent); }

/* Mobile: stack sidebar above main, slimmer padding */
@media (max-width: 720px) {
  body { flex-direction: column; }
  aside { width: 100%; height: auto; max-height: 40vh; border-right: none; border-bottom: 1px solid var(--border); }
  main { height: auto; padding: 24px 20px 60px; }
  main .reader { max-width: 100%; }
}
</style>
</head>
<body>
<a class="skip-link" href="#main-content">Skip to main content</a>
<aside>
  <div class="sidebar-header">
    <h1><a href="/">yt2md</a></h1>
    <button class="theme-toggle" type="button" onclick="cycleTheme()" aria-label="Cycle theme: auto / light / dark">🌓</button>
  </div>

  <nav aria-label="Per-video digests">
  <h2>Digests {% if unread_digest_count %}<span class="unread-count">{{ unread_digest_count }} new</span>{% else %}({{ digests|length }}){% endif %}</h2>
  <ul>
    {% for d in digests %}
    <li{% if current == 'digest:' + d.id %} class="active"{% endif %}{% if d.unread %} class="unread"{% endif %}>
      <a href="/digests/{{ d.id }}/">{% if d.unread %}<span class="unread-dot" aria-label="unread"></span>{% endif %}{{ d.title }}</a>
    </li>
    {% else %}
    <li class="empty">none yet</li>
    {% endfor %}
  </ul>
  </nav>

  <nav aria-label="Manage">
  <h2>Manage</h2>
  <ul>
    <li{% if current == 'channels' %} class="active"{% endif %}><a href="/channels">Subscriptions ({{ channel_count }})</a></li>
    <li{% if current == 'one-off' %} class="active"{% endif %}><a href="/one-off">One-off digest</a></li>
    <li{% if current == 'schedule' %} class="active"{% endif %}><a href="/schedule">Schedule</a></li>
    <li{% if current == 'activity' %} class="active"{% endif %}><a href="/activity">Activity</a></li>
    <li{% if current == 'settings' %} class="active"{% endif %}><a href="/settings">Settings</a></li>
  </ul>
  </nav>
</aside>
<main id="main-content" tabindex="-1">
  <div class="reader">
    {{ body|safe }}
  </div>
</main>
<script>
function applyTheme() {
  const stored = localStorage.getItem('yt2md-theme') || 'auto';
  const root = document.documentElement;
  if (stored === 'auto') root.removeAttribute('data-theme');
  else root.setAttribute('data-theme', stored);
  const btn = document.querySelector('.theme-toggle');
  if (btn) {
    const icons = {auto: '🌓', light: '☀️', dark: '🌙'};
    btn.textContent = icons[stored];
    btn.title = 'Theme: ' + stored + ' (click to cycle)';
  }
}
function cycleTheme() {
  const cur = localStorage.getItem('yt2md-theme') || 'auto';
  const next = {auto: 'light', light: 'dark', dark: 'auto'}[cur];
  localStorage.setItem('yt2md-theme', next);
  applyTheme();
}
applyTheme();
</script>
</body>
</html>
"""


def _list_digests(digests_dir: Path) -> List[dict]:
    if not digests_dir.exists():
        return []
    results = []
    for d in sorted(digests_dir.iterdir(), key=lambda p: p.stat().st_mtime, reverse=True):
        digest_md = d / "digest.md"
        if not d.is_dir() or not digest_md.exists():
            continue
        title = d.name
        try:
            for line in digest_md.read_text().splitlines():
                if line.startswith("# "):
                    title = line[2:].strip()
                    break
        except Exception:
            pass
        results.append({"id": d.name, "title": title, "mtime": d.stat().st_mtime})
    return results


def _render_markdown(text: str) -> str:
    import markdown as md_lib
    html = md_lib.markdown(text, extensions=["fenced_code", "tables", "sane_lists"])
    # Rewrite cross-references to other digests (e.g. ../digests/X/digest.md) into view URLs.
    html = re.sub(r'href="[^"]*digests/([^"/]+)/digest\.md"', r'href="/digests/\1/"', html)
    # Open external links in a new tab so the reader doesn't lose their place
    # when clicking a YouTube timestamp / "Watch on YouTube" link. rel=noopener
    # blocks the new tab from manipulating window.opener (web-security best
    # practice). Skip anchors that already have a target= attribute.
    html = re.sub(
        r'<a (href="https?://[^"]+")(?![^>]*\btarget=)',
        r'<a \1 target="_blank" rel="noopener"',
        html,
    )
    return html


def cmd_serve(args) -> int:
    try:
        from flask import Flask, render_template_string, abort, send_from_directory
    except ImportError:
        sys.exit(
            "Flask is required for the reader. Reinstall with:\n"
            "  uv tool install --reinstall git+https://github.com/jyouturner/youtube-to-markdown"
        )

    data_dir = get_data_dir()
    digests_dir = data_dir / "digests"

    app = Flask(__name__)
    # Disable Flask's default request logging — keep stdout clean.
    import logging
    logging.getLogger("werkzeug").setLevel(logging.WARNING)

    def page(body: str, *, title: str, current: str, base_href: str = None):
        # Annotate listing items with read state so the sidebar can show "new" markers.
        digests = _list_digests(digests_dir)
        try:
            read_digests = _read_digest_ids()
        except Exception:
            read_digests = set()
        for d in digests:
            d["unread"] = d["id"] not in read_digests
        unread_digest_count = sum(1 for d in digests if d["unread"])

        # Persistent banner above every page when neither auth path is
        # configured. Reading cached digests still works; only generation does,
        # so we warn rather than gate. The Setup page is the one place this is
        # hidden (the page itself IS the configuration UI).
        has_api_key = bool(os.environ.get("ANTHROPIC_API_KEY"))
        has_claude_code = (
            claude_code_installed()
            and _claude_code_session_state.get("logged_in", False)
        )
        if current != "setup" and not (has_api_key or has_claude_code):
            banner = (
                '<div class="flash" style="border-left-color: #c00;">'
                '<strong>LLM auth not configured.</strong> '
                'Generating digests and panel discussions requires either an '
                'Anthropic API key or a Claude Code subscription login. '
                'Reading existing digests still works. '
                '<a href="/setup">Set it up →</a>'
                '</div>'
            )
            body = banner + body

        return render_template_string(
            SERVE_PAGE_TEMPLATE,
            body=body,
            title=title,
            current=current,
            base_href=base_href,
            digests=digests,
            channel_count=len(read_channels()),
            unread_digest_count=unread_digest_count,
        )

    def _require_llm_or_redirect():
        """Helper for action endpoints: returns a redirect Response if neither
        an API key nor a logged-in Claude Code sandbox is configured. Use as:
        r = _require_llm_or_redirect(); if r is not None: return r."""
        from flask import redirect
        if os.environ.get("ANTHROPIC_API_KEY"):
            return None
        if claude_code_installed() and _claude_code_session_state.get("logged_in"):
            return None
        return redirect("/setup?msg=Auth+required+to+run+this+action.")

    @app.route("/")
    def home():
        from flask import redirect
        digests = _list_digests(digests_dir)
        channels = read_channels()

        # True first-run (no auth at all + nothing on disk yet): land directly
        # on the setup page so the user isn't asked to subscribe before they
        # can generate anything. Skip when there are existing digests — those
        # should still be readable even without a key.
        has_api_key = bool(os.environ.get("ANTHROPIC_API_KEY"))
        has_claude_code = (
            claude_code_installed()
            and _claude_code_session_state.get("logged_in", False)
        )
        if not digests and not channels and not (has_api_key or has_claude_code):
            return redirect("/setup")

        # Empty states first — show a real CTA, not a list of zero items.
        if not digests:
            if not channels:
                body = (
                    "<h1>Welcome to yt2md</h1>"
                    "<p>You haven't subscribed to any channels yet.</p>"
                    '<p><a class="cta" href="/channels">Add your first channel →</a></p>'
                )
            else:
                body = (
                    "<h1>Polling is set up</h1>"
                    f"<p>You're watching {len(channels)} channel(s). Your first digest "
                    "will appear after the next polling run (every few hours, or "
                    'fire one now from the <a href="/schedule">Schedule</a> page).</p>'
                )
            return page(body, title="yt2md", current="home")

        # Featured content: the latest digest's body, with an eyebrow link.
        featured = digests[0]
        featured_md = (digests_dir / featured["id"] / "digest.md").read_text()
        body = (
            f'<p class="featured-eyebrow">Latest digest · '
            f'<a href="/digests/{featured["id"]}/">{featured["title"]}</a></p>'
        )
        body += _render_markdown(featured_md)
        base_href = f"/digests/{featured['id']}/"

        # No "More digests" footer here — sidebar is the navigation surface;
        # showing the same list twice is just noise.

        return page(body, title="Home", current="home", base_href=base_href)

    @app.route("/channels", methods=["GET"])
    def channels_page():
        from flask import request
        channels = read_channels()
        digests = _list_digests(digests_dir)
        sched_state = _load_schedule_state()
        poll_has_run = bool((sched_state.get("poll") or {}).get("last_started_at"))
        flash = request.args.get("msg", "")

        body = "<h1>Subscriptions</h1>"
        if flash:
            body += f'<div class="flash">{flash}</div>'

        # "What's next" guidance — adapts to current state. Only shown until the
        # user has at least one digest, then disappears.
        next_step = None
        if not channels:
            next_step = (
                "Paste a YouTube channel URL below to get started. "
                "Each new video on this channel will be auto-digested."
            )
        elif not poll_has_run:
            next_step = (
                "You\'re subscribed but the scheduler hasn\'t fired its first poll yet. "
                'It runs every few hours — or trigger one now from the '
                '<a href="/schedule">Schedule page</a>.'
            )
        elif not digests:
            next_step = (
                "Polling fires every 6 hours. Your first digest will land after the next run "
                '— or fire one now from the <a href="/schedule">Schedule page</a>.'
            )
        if next_step:
            body += f'<div class="next-step"><strong>Next step:</strong> {next_step}</div>'

        body += (
            '<form method="post" action="/channels" class="add-form">'
            '<label for="channel-url" class="sr-only">YouTube channel URL</label>'
            '<input id="channel-url" type="text" name="url" '
            'placeholder="https://www.youtube.com/@channel/videos  (or @handle)" '
            'autofocus required>'
            '<button type="submit">Add</button>'
            '</form>'
        )
        if channels:
            body += '<ul class="channel-list">'
            for ch in channels:
                body += (
                    '<li>'
                    f'<span class="url">{ch}</span>'
                    '<form method="post" action="/channels/remove" style="margin:0;">'
                    f'<input type="hidden" name="url" value="{ch}">'
                    '<button type="submit">Remove</button>'
                    '</form>'
                    '</li>'
                )
            body += '</ul>'
        else:
            body += "<p class='empty-state'>No subscriptions yet. Paste a YouTube channel URL above.</p>"
        body += (
            "<p class='meta-info' style='margin-top:32px'>"
            "Stored in <code>~/yt2md/channels.txt</code>."
            "</p>"
        )
        return page(body, title="Subscriptions", current="channels")

    @app.route("/channels", methods=["POST"])
    def channels_add():
        from flask import request, redirect
        url = normalize_channel_url(request.form.get("url", ""))
        if not is_url(url):
            return redirect("/channels?msg=Not+a+valid+URL")
        channels = read_channels()
        if url in channels:
            return redirect(f"/channels?msg=Already+watching+{url}")
        channels.append(url)
        write_channels(channels)
        return redirect(f"/channels?msg=Added+{url}")

    @app.route("/channels/remove", methods=["POST"])
    def channels_remove():
        from flask import request, redirect
        url = request.form.get("url", "").strip()
        channels = [c for c in read_channels() if c != url]
        write_channels(channels)
        return redirect(f"/channels?msg=Removed+{url}")

    @app.route("/schedule")
    def schedule_page():
        from flask import request
        from html import escape as h
        import time as _t
        flash = request.args.get("msg", "")
        cfg = load_schedule_config()
        sched_state = _load_schedule_state()

        body = "<h1>Schedule</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'

        body += '<form method="post" action="/schedule/save" class="schedule-form">'
        body += '<div class="schedule-fields">'
        body += '<label>Polling interval'
        body += f'  <input type="number" name="poll_hours" value="{cfg["poll_interval_hours"]}" min="0.1" step="0.1" required>'
        body += '  <span class="suffix">hours</span>'
        body += '</label>'
        body += '</div>'  # schedule-fields

        body += '<button type="submit" class="primary">Save</button>'
        body += '</form>'

        body += (
            f'<p class="meta-info">Current schedule: {h(_format_schedule_summary(cfg))}. '
            'Scheduling runs inside this server — pauses while it\'s down, '
            'catches up on missed slots when you start it again.</p>'
        )

        for kind, friendly, desc in [
            ("poll", "Polling", "fires <code>yt2md watch run</code>"),
        ]:
            sentence, dot_class = _scheduler_status_summary(kind, sched_state)
            s = sched_state.get(kind) or {}
            next_at = _compute_next_poll(cfg, s.get("last_started_at"))

            body += f'<div class="job-block"><h3><span class="dot {dot_class}"></span>{friendly}</h3>'
            body += f'<p class="meta-info" style="margin: 0 0 12px;">{desc}</p>'
            body += f'<p class="job-summary">{sentence}</p>'
            body += (f'<p class="meta-info">Next run: {h(_format_next_run(next_at))} '
                     f'({h(_t.strftime("%Y-%m-%d %H:%M", _t.localtime(next_at)))}).</p>')
            running = kind in _scheduler_jobs
            disabled = " disabled" if running else ""
            running_label = " (already running)" if running else ""
            body += (
                f'<div class="job-actions">'
                f'<form method="post" action="/schedule/run/{kind}" style="margin:0;">'
                f'<button type="submit"{disabled}>Run now{running_label}</button>'
                f'</form></div>'
            )

            body += '<details><summary>Diagnostics</summary>'
            body += '<table class="status-table">'
            for k in ("last_started_at", "last_finished_at", "last_exit_code", "last_pid"):
                v = s.get(k)
                if v is None:
                    continue
                if k.endswith("_at"):
                    v = _t.strftime("%Y-%m-%d %H:%M:%S", _t.localtime(v))
                body += f'<tr><td>{k.replace("_", " ")}</td><td>{h(str(v))}</td></tr>'
            body += '</table></details>'

            log_path = data_dir / "logs" / f"{kind}.log"
            body += '<details style="margin-top: 8px;"><summary>Recent log (last 20 lines)</summary>'
            body += f'<div class="log-block">{h(_tail_log(log_path, 20))}</div>'
            body += '</details>'
            body += '</div>'

        body += '<p class="meta-info">Refresh the page to see updated status after a run.</p>'
        return page(body, title="Schedule", current="schedule")

    @app.route("/schedule/save", methods=["POST"])
    def schedule_save():
        from flask import redirect, request
        cfg = load_schedule_config()
        try:
            if request.form.get("poll_hours"):
                cfg["poll_interval_hours"] = float(request.form["poll_hours"])
        except Exception as e:
            return redirect(f"/schedule?msg=Invalid+input:+{e}")
        save_schedule_config(cfg)
        return redirect("/schedule?msg=Saved+(scheduler+picks+up+within+30s)")

    @app.route("/schedule/run/<job>", methods=["POST"])
    def schedule_run(job):
        from flask import redirect
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        if job != "poll":
            abort(404)
        with _schedule_lock():
            existing = _scheduler_jobs.get(job)
            if existing is not None and existing.poll() is None:
                return redirect(f"/schedule?msg={job}+is+already+running")
            proc = _fire_scheduled_job(job)
        if proc is None:
            return redirect(f"/schedule?msg=Failed+to+fire+{job}+(yt2md+not+on+PATH)")
        return redirect(
            f"/schedule?msg=Started+{job}+(pid+{proc.pid}).+Refresh+for+status."
        )

    @app.route("/settings", methods=["GET"])
    def settings_page():
        from flask import request
        from html import escape as h
        flash = request.args.get("msg", "")
        s = load_settings()
        # Show the EFFECTIVE value (settings.json with .env fallback) so the
        # form matches what the system is actually using. On Save we persist
        # whatever the user submits, which becomes the new canonical value.
        for key, env_name in (
            ("digest_model", "YT2MD_DIGEST_MODEL"),
            ("panel_model", "YT2MD_PANEL_MODEL"),
            ("whisper_model", "YT2MD_WHISPER_MODEL"),
            ("cookies_from_browser", "YT2MD_COOKIES_FROM_BROWSER"),
            ("digest_language", "YT2MD_DIGEST_LANGUAGE"),
        ):
            if not s.get(key) and os.environ.get(env_name):
                s[key] = os.environ[env_name]

        whisper_choices = ("tiny", "base", "small", "medium", "large-v2", "large-v3")
        cookie_choices = ("", "chrome", "firefox", "safari", "brave", "edge",
                          "chromium", "opera", "vivaldi")

        body = "<h1>Settings</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'
        body += (
            '<p class="meta-info">Stored in <code>~/yt2md/settings.json</code> '
            '(API key in <code>~/yt2md/.env</code>). '
            'New values take effect on the next one-off submit / scheduled poll / '
            '"Discuss with experts" click — no restart required.</p>'
        )

        body += '<form method="post" action="/settings/save" class="schedule-form">'
        body += '<div class="schedule-fields" style="grid-template-columns: 1fr;">'

        cur_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if cur_key:
            cur_state = (
                f'<span style="color: var(--accent);">set</span> '
                f'(<code>{h(cur_key[:7])}…{h(cur_key[-4:])}</code>)'
            )
        else:
            cur_state = '<span style="color: #c00;">not set</span>'
        body += (
            '<label>Anthropic API key'
            '  <input type="password" name="anthropic_api_key" '
            '    placeholder="sk-ant-... (leave blank to keep current)" '
            '    autocomplete="off">'
            '  <span class="suffix" style="display:block;">'
            f'    Current: {cur_state}. '
            f'    {API_KEY_COST_NOTE} '
            '    Get a key at '
            '    <a href="https://console.anthropic.com/settings/keys" target="_blank" '
            '       rel="noopener">console.anthropic.com/settings/keys</a>. '
            '    The key is validated on save with a 1-token test call.'
            '  </span>'
            '</label>'
        )

        body += (
            '<label>Digest model'
            f'  <input type="text" name="digest_model" value="{h(s["digest_model"])}" required>'
            '  <span class="suffix" style="display:block;">'
            'Anthropic model ID for the per-video digest. e.g. <code>claude-sonnet-4-6</code>, '
            '<code>claude-opus-4-7</code>, <code>claude-haiku-4-5-20251001</code>.'
            '</span>'
            '</label>'
        )

        body += (
            '<label>Panel-discussion model'
            f'  <input type="text" name="panel_model" value="{h(s["panel_model"])}" required>'
            '  <span class="suffix" style="display:block;">'
            'Used when you click "Discuss with experts" on a digest. Multi-perspective '
            'synthesis benefits from a stronger reasoning model — Opus is the default.'
            '</span>'
            '</label>'
        )

        body += '<label>Whisper model'
        body += '  <select name="whisper_model">'
        for w in whisper_choices:
            sel = ' selected' if s["whisper_model"] == w else ''
            body += f'    <option value="{w}"{sel}>{w}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Local STT fallback when YouTube has no captions. Larger = better quality, '
            'slower, bigger first-run download. <code>medium</code> is a good default.'
            '</span>'
            '</label>'
        )

        body += '<label>Digest language'
        body += '  <select name="digest_language">'
        for code, label in (
            ("auto", "auto — match the transcript's language"),
            ("en", "en — always English"),
        ):
            sel = ' selected' if s.get("digest_language", "auto") == code else ''
            body += f'    <option value="{code}"{sel}>{label}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Applies to both the per-video digest and the panel discussion. '
            '<code>auto</code> writes in the source language (e.g. Chinese for a Chinese-language video). '
            '<code>en</code> forces English regardless.'
            '</span>'
            '</label>'
        )

        body += '<label>Cookies from browser'
        body += '  <select name="cookies_from_browser">'
        for c in cookie_choices:
            sel = ' selected' if s.get("cookies_from_browser", "") == c else ''
            label = "(none)" if c == "" else c
            body += f'    <option value="{c}"{sel}>{label}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'YouTube increasingly requires logged-in cookies. Pick the browser you\'re '
            'signed into YouTube on; yt-dlp extracts cookies on each run. Leave as '
            '"(none)" if you only digest publicly-accessible videos.'
            '</span>'
            '</label>'
        )

        # LLM backend selector
        backend_choices = (
            ("auto", "auto — pick API when ANTHROPIC_API_KEY is set, else Claude Code"),
            ("api", "api — direct Anthropic API (requires ANTHROPIC_API_KEY)"),
            ("claude-code", "claude-code — bundled Claude Code subprocess (requires Claude.ai login)"),
        )
        body += '<label>LLM backend'
        body += '  <select name="llm_backend">'
        for code, label in backend_choices:
            sel = ' selected' if s.get("llm_backend", "auto") == code else ''
            body += f'    <option value="{code}"{sel}>{h(label)}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Which auth path to use for digest / panel calls. The "auto" mode '
            'picks the cheapest available path. Switch from <a href="/setup">/setup</a>.'
            '</span>'
            '</label>'
        )

        # Claude Code vision toggle
        cc_vision_checked = ' checked' if s.get("claude_code_vision") else ''
        body += (
            '<label>'
            f'  <input type="checkbox" name="claude_code_vision" value="1"{cc_vision_checked}> '
            'Enable vision frame-picking under Claude Code backend'
            '  <span class="suffix" style="display:block;">'
            'Off by default. The Claude Code CLI has no native image flag, so '
            'enabling this base64-embeds frames into prompts (token-heavy). '
            'Has no effect when the API backend is in use (which always uses '
            'native vision).'
            '  </span>'
            '</label>'
        )

        body += '</div>'  # schedule-fields
        body += '<button type="submit" class="primary">Save</button>'
        body += '</form>'
        return page(body, title="Settings", current="settings")

    @app.route("/settings/save", methods=["POST"])
    def settings_save():
        from flask import redirect, request
        from urllib.parse import quote_plus
        s = load_settings()
        for key in ("digest_model", "panel_model", "whisper_model",
                    "cookies_from_browser", "digest_language", "llm_backend"):
            v = request.form.get(key)
            if v is not None:
                s[key] = v.strip()
        # Checkboxes are absent from request.form when unchecked.
        s["claude_code_vision"] = bool(request.form.get("claude_code_vision"))
        save_settings(s)

        # API key is stored separately (.env, not settings.json) so non-secret
        # config can be checked into a shared settings file without leaking it.
        new_key = (request.form.get("anthropic_api_key") or "").strip()
        if new_key:
            err = validate_api_key(new_key)
            if err:
                return redirect(
                    f"/settings?msg=Settings+saved+but+API+key+rejected:+{quote_plus(err)}"
                )
            set_env_var("ANTHROPIC_API_KEY", new_key)
            return redirect("/settings?msg=Saved+(API+key+validated).")
        return redirect("/settings?msg=Saved.")

    @app.route("/setup", methods=["GET"])
    def setup_page():
        from flask import request
        from html import escape as h
        flash = request.args.get("msg", "")

        body = "<h1>Connect Claude</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'
        body += (
            '<p class="meta-info">Pick one of the two paths below. yt2md needs '
            'access to a Claude model to generate digests and panel discussions; '
            'either an Anthropic API key (per-call billing) or a Claude.ai '
            'subscription via the bundled Claude Code (no extra billing).</p>'
        )

        # --- Side-by-side dual-auth grid ---
        body += (
            '<div class="auth-grid" style="display: grid; '
            'grid-template-columns: repeat(auto-fit, minmax(360px, 1fr)); '
            'gap: 20px; margin-top: 20px;">'
        )

        # API key panel
        cur_key = os.environ.get("ANTHROPIC_API_KEY", "")
        api_already_set = bool(cur_key)
        api_panel = '<div class="schedule-form">'
        api_panel += '<h2 style="margin-top: 0;">Anthropic API key</h2>'
        if api_already_set:
            api_panel += (
                '<div class="flash" style="border-left-color: var(--accent); '
                'margin-top: 0;">'
                f'<strong>Configured</strong> '
                f'(<code>{h(cur_key[:7])}…{h(cur_key[-4:])}</code>). '
                '<a href="/">Go to library →</a>'
                '</div>'
            )
        api_panel += (
            '<p class="meta-info" style="margin-top: 0;">'
            'For developers / users who already have an Anthropic API account. '
            'Pay-per-call. Native vision support, prompt caching.'
            '</p>'
            '<ol style="margin: 8px 0 16px 20px; padding: 0; font-size: 14px;">'
            '<li>Get a key at <a href="https://console.anthropic.com/settings/keys" '
            'target="_blank" rel="noopener">console.anthropic.com</a>.</li>'
            '<li>Add a payment method (subscription does NOT cover API usage).</li>'
            '<li>Paste below.</li>'
            '</ol>'
            f'<p class="meta-info" style="font-size: 13px;">{API_KEY_COST_NOTE}</p>'
            '<form method="post" action="/setup/save-api-key">'
            '<label>API key'
            '  <input type="password" name="anthropic_api_key" '
            '    placeholder="sk-ant-..." autocomplete="off" '
            'style="width: 100%; box-sizing: border-box;">'
            '</label>'
            '<button type="submit" class="primary" style="margin-top: 12px;">'
            'Save and validate</button>'
            '</form>'
            '</div>'
        )

        # Claude Code panel — installation + login flow
        snap = claude_setup_snapshot()
        cc_panel = '<div class="schedule-form">'
        cc_panel += '<h2 style="margin-top: 0;">Claude.ai subscription</h2>'
        cc_panel += (
            '<p class="meta-info" style="margin-top: 0;">'
            'Uses your Pro/Max plan via a bundled, sandboxed copy of Claude Code. '
            'No extra billing. Vision off by default (toggle in Settings). '
            'Prompt caching disappears, so very long videos may use more tokens.'
            '</p>'
        )

        # Status & action area, populated by JS poll, with a server-rendered
        # initial state for users who hit the page without JS.
        cc_panel += '<div id="cc-status" style="margin-bottom: 12px;">'
        if snap["logged_in"]:
            cc_panel += (
                '<div class="flash" style="border-left-color: var(--accent); margin: 0;">'
                '<strong>Signed in.</strong> '
                '<a href="/">Go to library →</a>'
                '</div>'
            )
        elif not snap["node_ok"]:
            cc_panel += (
                '<div class="flash" style="border-left-color: #c00; margin: 0;">'
                f'<strong>Node.js {MIN_NODE_MAJOR}+ required.</strong> '
                'Install with <code>brew install node</code> (macOS) or your '
                'package manager, then refresh this page.'
                '</div>'
            )
        elif not snap["installed"]:
            cc_panel += (
                '<p class="meta-info" style="margin: 0;">'
                f'Step 1: install Claude Code into <code>~/yt2md/claude-code/</code> '
                f'(~200&nbsp;MB; isolated from any system install).'
                '</p>'
            )
        else:
            cc_panel += (
                '<p class="meta-info" style="margin: 0;">'
                'Step 2: sign in. A new browser tab will open for Claude.ai OAuth. '
                'Complete sign-in there; this page will update automatically.'
                '</p>'
            )
        cc_panel += '</div>'

        # Action buttons (the JS toggles these based on status).
        install_disabled = (
            ' disabled' if not snap["node_ok"] or snap["install_running"]
            or snap["installed"] else ''
        )
        login_visible = snap["installed"] and not snap["logged_in"]
        cc_panel += '<div id="cc-actions" style="display: flex; gap: 8px; flex-wrap: wrap;">'
        cc_panel += (
            f'<form method="post" action="/setup/install-claude" style="display:inline;">'
            f'<button type="submit" class="primary"{install_disabled} '
            'id="cc-install-btn">Install Claude Code</button>'
            '</form>'
        )
        if login_visible:
            cc_panel += (
                '<form method="post" action="/setup/login-claude" style="display:inline;">'
                '<button type="submit" class="primary" id="cc-login-btn">'
                'Sign in with Claude.ai</button>'
                '</form>'
            )
        if snap["logged_in"]:
            cc_panel += (
                '<form method="post" action="/setup/logout-claude" style="display:inline;" '
                'onsubmit="return confirm(\'Sign out of Claude Code?\');">'
                '<button type="submit">Sign out</button>'
                '</form>'
            )
        cc_panel += '</div>'

        # Live log tail (hidden when both logs are empty).
        cc_panel += (
            '<details id="cc-log-details" style="margin-top: 12px;'
            + ('' if (snap["install_log_tail"] or snap["login_log_tail"]) else ' display: none;')
            + '">'
            '<summary>Recent log</summary>'
            '<pre id="cc-log-tail" class="log-block" style="max-height: 200px; '
            'overflow: auto; font-size: 11px;">'
            + h((snap["login_log_tail"] or snap["install_log_tail"] or "").strip())
            + '</pre>'
            '</details>'
        )
        cc_panel += '</div>'

        body += api_panel + cc_panel + "</div>"  # close auth-grid

        # Polling JS: refreshes the Claude Code panel every 2s while a job is
        # running, so the user sees install / login progress without refreshing.
        body += """
<script>
(function () {
  const statusEl = document.getElementById('cc-status');
  const actionsEl = document.getElementById('cc-actions');
  const logDetailsEl = document.getElementById('cc-log-details');
  const logTailEl = document.getElementById('cc-log-tail');
  if (!statusEl) return;
  let pollTimer = null;

  function escapeHtml(s) {
    const d = document.createElement('div'); d.textContent = String(s ?? ''); return d.innerHTML;
  }

  function renderStatus(s) {
    if (s.logged_in) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: var(--accent); margin: 0;">' +
        '<strong>Signed in.</strong> <a href="/">Go to library →</a></div>';
    } else if (!s.node_ok) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: #c00; margin: 0;">' +
        '<strong>Node.js ' + 18 + '+ required.</strong> ' +
        'Install with <code>brew install node</code> (macOS) or your package manager, then refresh this page.' +
        '</div>';
    } else if (s.install_running) {
      statusEl.innerHTML =
        '<p class="meta-info" style="margin: 0;">Installing Claude Code… this takes ~30s.</p>';
    } else if (s.login_running) {
      statusEl.innerHTML =
        '<p class="meta-info" style="margin: 0;">Waiting for Claude.ai OAuth to complete in the browser tab that just opened…</p>';
    } else if (s.install_error) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: #c00; margin: 0;"><strong>Install failed.</strong> ' +
        escapeHtml(s.install_error) + '</div>';
    } else if (s.login_error) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: #c00; margin: 0;"><strong>Login failed.</strong> ' +
        escapeHtml(s.login_error) + '</div>';
    } else if (!s.installed) {
      statusEl.innerHTML = '<p class="meta-info" style="margin: 0;">Step 1: install Claude Code into <code>~/yt2md/claude-code/</code>.</p>';
    } else {
      statusEl.innerHTML = '<p class="meta-info" style="margin: 0;">Step 2: sign in. A new browser tab will open for Claude.ai OAuth.</p>';
    }
  }

  function renderActions(s) {
    let html = '';
    const installDisabled = (!s.node_ok || s.install_running || s.installed) ? ' disabled' : '';
    html += '<form method="post" action="/setup/install-claude" style="display:inline;">' +
            '<button type="submit" class="primary"' + installDisabled + ' id="cc-install-btn">' +
            (s.install_running ? 'Installing…' : 'Install Claude Code') + '</button></form>';
    if (s.installed && !s.logged_in) {
      html += '<form method="post" action="/setup/login-claude" style="display:inline;">' +
              '<button type="submit" class="primary"' + (s.login_running ? ' disabled' : '') + ' id="cc-login-btn">' +
              (s.login_running ? 'Waiting for OAuth…' : 'Sign in with Claude.ai') + '</button></form>';
    }
    if (s.logged_in) {
      html += '<form method="post" action="/setup/logout-claude" style="display:inline;" ' +
              'onsubmit="return confirm(\\'Sign out of Claude Code?\\');">' +
              '<button type="submit">Sign out</button></form>';
    }
    actionsEl.innerHTML = html;
  }

  function renderLog(s) {
    const tail = (s.login_log_tail || s.install_log_tail || '').trim();
    if (!tail) {
      logDetailsEl.style.display = 'none';
      return;
    }
    logDetailsEl.style.display = '';
    logTailEl.textContent = tail;
  }

  async function poll() {
    try {
      const res = await fetch('/setup/claude-status');
      if (!res.ok) return;
      const s = await res.json();
      renderStatus(s);
      renderActions(s);
      renderLog(s);
      // Auto-redirect home once login lands.
      if (s.logged_in) {
        clearInterval(pollTimer);
        setTimeout(() => { window.location.href = '/?msg=Signed+in+via+Claude.ai+subscription.'; }, 1500);
        return;
      }
      // Stop polling if nothing is in flight (saves cycles).
      if (!s.install_running && !s.login_running) {
        clearInterval(pollTimer);
        pollTimer = null;
      }
    } catch (e) { /* network blip — try again next tick */ }
  }

  // Always do an immediate poll so transient state from a just-submitted
  // form is reflected without the 2s delay.
  poll();
  pollTimer = setInterval(poll, 2000);
})();
</script>
"""

        return page(body, title="Set up", current="setup")

    @app.route("/setup/save-api-key", methods=["POST"])
    def setup_save_api_key():
        from flask import redirect, request
        from urllib.parse import quote_plus
        new_key = (request.form.get("anthropic_api_key") or "").strip()
        if not new_key:
            return redirect("/setup?msg=Paste+a+key+first.")
        err = validate_api_key(new_key)
        if err:
            return redirect(f"/setup?msg=Key+rejected:+{quote_plus(err)}")
        set_env_var("ANTHROPIC_API_KEY", new_key)
        return redirect("/?msg=API+key+saved.+You+can+now+generate+digests.")

    # Legacy route name retained for back-compat with bookmarks/redirects.
    @app.route("/setup/save", methods=["POST"])
    def setup_save_legacy():
        return setup_save_api_key()

    @app.route("/setup/install-claude", methods=["POST"])
    def setup_install_claude():
        from flask import redirect
        err = start_install_job()
        if err:
            from urllib.parse import quote_plus
            return redirect(f"/setup?msg={quote_plus(err)}")
        return redirect("/setup")

    @app.route("/setup/login-claude", methods=["POST"])
    def setup_login_claude():
        from flask import redirect
        err = start_login_job()
        if err:
            from urllib.parse import quote_plus
            return redirect(f"/setup?msg={quote_plus(err)}")
        return redirect("/setup")

    @app.route("/setup/logout-claude", methods=["POST"])
    def setup_logout_claude():
        from flask import redirect
        rc, out = claude_logout()
        if rc != 0:
            from urllib.parse import quote_plus
            return redirect(f"/setup?msg=Logout+failed:+{quote_plus(out[:200])}")
        return redirect("/setup?msg=Signed+out+of+Claude+Code.")

    @app.route("/setup/claude-status")
    def setup_claude_status():
        from flask import jsonify
        return jsonify(claude_setup_snapshot())

    @app.route("/one-off", methods=["GET"])
    def one_off_page():
        from flask import request
        from html import escape as h
        flash = request.args.get("msg", "")
        active = _list_active_oneoff_jobs()  # also reaps exited subprocesses
        failures = _list_recent_oneoff_failures()

        body = "<h1>One-off digest</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'
        body += (
            '<p class="meta-info">Paste a YouTube video URL. The digest runs in the '
            'background and lands in your library when complete (1–25 min depending on '
            'video length and the vision pass). You can close this tab — it keeps running.</p>'
        )
        body += (
            '<form method="post" action="/one-off" class="add-form">'
            '<label for="video-url" class="sr-only">YouTube video URL</label>'
            '<input id="video-url" type="text" name="url" '
            'placeholder="https://youtu.be/... (or full watch URL)" '
            'autofocus required>'
            '<button type="submit">Digest</button>'
            '</form>'
        )

        import time as _t
        # Read the log once so per-job stage lookups don't hit disk repeatedly.
        log_path = data_dir / "logs" / "oneoff.log"
        try:
            log_text = log_path.read_text(errors="replace")
        except OSError:
            log_text = ""

        # Always render the section containers so the polling JS can target them
        # (display:none hides them when empty).
        active_hidden = "" if active else " style='display:none'"
        body += f"<section id='oneoff-active-section'{active_hidden}>"
        body += "<h2>In progress</h2>"
        body += "<ul id='oneoff-active-list' class='channel-list'>"
        for j in active:
            elapsed = int(_t.time() - j["started"])
            stage = _describe_job_stage(log_text, j["video_id"])
            body += (
                '<li>'
                f'<span class="url"><strong>{h(j["video_id"])}</strong> · '
                f'{h(j["url"])}</span>'
                f'<span style="color: var(--muted); font-size: 12px;">'
                f'{elapsed//60}m {elapsed%60}s · {h(stage)}</span>'
                '</li>'
            )
        body += "</ul></section>"

        failures_hidden = "" if failures else " style='display:none'"
        body += f"<section id='oneoff-failures-section'{failures_hidden}>"
        body += "<h2>Recent failures</h2>"
        body += "<ul id='oneoff-failures-list' class='channel-list'>"
        for f in failures:
            ago = int(_t.time() - f["ended"])
            if ago < 60:
                when = f"{ago}s ago"
            elif ago < 3600:
                when = f"{ago // 60}m ago"
            else:
                when = f"{ago // 3600}h ago"
            err = f["error"] or f"exit code {f['exit_code']}"
            body += (
                '<li>'
                f'<span class="url"><strong>{h(f["video_id"])}</strong> · '
                f'{h(f["url"])}</span>'
                f'<span style="color: var(--muted); font-size: 12px;">{when}</span>'
                f'<div style="color: var(--muted); font-size: 13px; margin-top: 4px;">{h(err)}</div>'
                '</li>'
            )
        body += "</ul>"
        body += (
            "<p class='meta-info'>"
            f"Full output: <code>{h(str(log_path))}</code>"
            "</p>"
        )
        body += "</section>"

        body += (
            "<p class='meta-info' style='margin-top: 32px;'>"
            "One-off digests share the same library as subscription mode. "
            "They appear in the sidebar's <strong>Digests</strong> section once ready."
            "</p>"
        )

        # Polling: refresh in-progress + failures every 2s without reloading the page.
        body += """
<script>
(function () {
  const activeSection  = document.getElementById('oneoff-active-section');
  const activeList     = document.getElementById('oneoff-active-list');
  const failSection    = document.getElementById('oneoff-failures-section');
  const failList       = document.getElementById('oneoff-failures-list');
  if (!activeSection || !failSection) return;

  const fmtElapsed = s => `${Math.floor(s/60)}m ${s%60}s`;
  const fmtAgo = s => s < 60 ? `${s}s ago` : (s < 3600 ? `${Math.floor(s/60)}m ago` : `${Math.floor(s/3600)}h ago`);
  const esc = s => { const d = document.createElement('div'); d.textContent = String(s ?? ''); return d.innerHTML; };

  function render(data) {
    const active = data.active || [];
    activeSection.style.display = active.length ? '' : 'none';
    activeList.innerHTML = active.map(j => `
      <li>
        <span class="url"><strong>${esc(j.video_id)}</strong> &middot; ${esc(j.url)}</span>
        <span style="color: var(--muted); font-size: 12px;">${esc(fmtElapsed(j.elapsed_secs))} &middot; ${esc(j.stage)}</span>
      </li>
    `).join('');

    const failures = data.failures || [];
    failSection.style.display = failures.length ? '' : 'none';
    failList.innerHTML = failures.map(f => `
      <li>
        <span class="url"><strong>${esc(f.video_id)}</strong> &middot; ${esc(f.url)}</span>
        <span style="color: var(--muted); font-size: 12px;">${esc(fmtAgo(f.ago_secs))}</span>
        <div style="color: var(--muted); font-size: 13px; margin-top: 4px;">${esc(f.error || ('exit code ' + f.exit_code))}</div>
      </li>
    `).join('');
  }

  async function poll() {
    try {
      const res = await fetch('/one-off/status', { cache: 'no-store' });
      if (res.ok) render(await res.json());
    } catch (_) { /* transient — try again next tick */ }
  }

  poll();
  setInterval(poll, 2000);
})();
</script>
"""
        return page(body, title="One-off digest", current="one-off")

    @app.route("/one-off/status")
    def one_off_status():
        from flask import jsonify
        import time as _t
        active = _list_active_oneoff_jobs()  # also reaps any just-exited subprocesses
        failures = _list_recent_oneoff_failures()
        log_path = data_dir / "logs" / "oneoff.log"
        try:
            log_text = log_path.read_text(errors="replace")
        except OSError:
            log_text = ""
        now = _t.time()
        return jsonify({
            "active": [
                {
                    "video_id": j["video_id"],
                    "url": j["url"],
                    "started": j["started"],
                    "elapsed_secs": int(now - j["started"]),
                    "stage": _describe_job_stage(log_text, j["video_id"]),
                }
                for j in active
            ],
            "failures": [
                {
                    "video_id": f["video_id"],
                    "url": f["url"],
                    "started": f["started"],
                    "ended": f["ended"],
                    "ago_secs": int(now - f["ended"]),
                    "exit_code": f["exit_code"],
                    "error": f["error"],
                }
                for f in failures
            ],
        })

    @app.route("/one-off", methods=["POST"])
    def one_off_submit():
        from flask import redirect, request
        import time as _t
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        url = request.form.get("url", "").strip()
        if not url:
            return redirect("/one-off?msg=URL+is+required")

        video_id = extract_video_id(url)
        if not video_id:
            return redirect(
                f"/one-off?msg=Couldn%27t+extract+a+YouTube+video+ID+from:+{url}"
            )

        # Already in library? Send them straight to the existing digest.
        existing = digests_dir / video_id / "digest.md"
        if existing.exists():
            return redirect(f"/digests/{video_id}/")

        # Already in progress? Show the page with a message rather than re-firing.
        for active in _list_active_oneoff_jobs():
            if active["video_id"] == video_id:
                return redirect(f"/one-off?msg=Already+digesting+{video_id}")

        # Fire and forget. start_new_session=True detaches the child so it survives
        # if the web server is killed. stdout/stderr go to oneoff.log.
        digest_path = digests_dir / video_id / "digest.md"
        digest_path.parent.mkdir(parents=True, exist_ok=True)
        log_path = data_dir / "logs" / "oneoff.log"
        log_path.parent.mkdir(parents=True, exist_ok=True)
        log_fd = open(log_path, "a")
        log_fd.write(f"\n===== {_t.strftime('%Y-%m-%d %H:%M:%S')} starting {video_id} ({url}) =====\n")
        log_fd.flush()

        yt2md_path = shutil.which("yt2md")
        if not yt2md_path:
            log_fd.close()
            return redirect("/one-off?msg=yt2md+not+on+PATH")

        try:
            proc = subprocess.Popen(
                [yt2md_path, url, "-o", str(digest_path)],
                cwd=digest_path.parent,
                stdout=log_fd,
                stderr=subprocess.STDOUT,
                env={**os.environ, **_settings_to_env(load_settings())},
                start_new_session=True,
            )
        finally:
            # Subprocess holds its own copy of the fd; safe for us to close.
            log_fd.close()

        _oneoff_jobs[proc.pid] = {
            "video_id": video_id,
            "started": _t.time(),
            "url": url,
            "proc": proc,
        }
        return redirect(
            f"/one-off?msg=Started+digesting+{video_id}+(check+sidebar+in+a+few+minutes)"
        )

    @app.route("/activity")
    def activity_page():
        from flask import request
        from html import escape as h
        import time as _t
        flt = request.args.get("status", "all")  # all | success | failed
        rows = _recent_runs(limit=200)
        if flt == "success":
            rows = [r for r in rows if r.get("success")]
        elif flt == "failed":
            rows = [r for r in rows if not r.get("success")]

        body = "<h1>Activity</h1>"
        body += '<p class="meta-info">Every completed one-off digest, success or failure. Persists across server restarts.</p>'

        # Cost summary across windows. Reads the central LLM usage log once
        # and bucket-totals — cheap, runs in a few ms even with thousands
        # of entries.
        usage_entries = read_llm_usage_log()
        now_ts = _t.time()
        WINDOWS = [
            ("today", 24 * 3600),
            ("7d", 7 * 24 * 3600),
            ("30d", 30 * 24 * 3600),
            ("all", None),
        ]
        window_totals: dict = {label: 0.0 for label, _ in WINDOWS}
        window_counts: dict = {label: 0 for label, _ in WINDOWS}
        backend_counts: dict = {}
        for e in usage_entries:
            cost = float(e.get("cost_usd", 0.0) or 0.0)
            ts = float(e.get("ts", 0) or 0)
            backend_counts[e.get("backend", "?")] = (
                backend_counts.get(e.get("backend", "?"), 0) + 1
            )
            for label, window in WINDOWS:
                if window is None or (now_ts - ts) <= window:
                    window_totals[label] += cost
                    window_counts[label] += 1
        backend_summary = ", ".join(
            f"{n} {b}" for b, n in sorted(backend_counts.items(), key=lambda kv: -kv[1])
        ) if backend_counts else ""
        body += (
            "<div class='schedule-form' style='display: grid; "
            "grid-template-columns: repeat(auto-fit, minmax(160px, 1fr)); "
            "gap: 12px 20px;'>"
        )
        for label, _ in WINDOWS:
            body += (
                f"<div><div style='color: var(--muted); font-size: 12px;'>"
                f"{h(label)}</div>"
                f"<div style='font-size: 22px; font-weight: 600;'>"
                f"${window_totals[label]:.2f}</div>"
                f"<div style='color: var(--muted); font-size: 11px;'>"
                f"{window_counts[label]} call(s)</div></div>"
            )
        body += "</div>"
        if backend_summary:
            body += (
                f"<p class='meta-info'>By backend: {h(backend_summary)}. "
                "Subscription (Claude Code) calls report $0 — billed via the "
                "user's plan. Pricing is estimated; treat as a guide.</p>"
            )

        # In-progress section — live one-off jobs the reaper hasn't logged yet.
        # Hidden by default; populated by the polling JS at the bottom of the page.
        active_now = _list_active_oneoff_jobs()
        log_path = data_dir / "logs" / "oneoff.log"
        try:
            log_text = log_path.read_text(errors="replace")
        except OSError:
            log_text = ""
        active_hidden = "" if active_now else " style='display:none'"
        body += f"<section id='activity-active-section'{active_hidden}>"
        body += "<h2>In progress</h2>"
        body += "<ul id='activity-active-list' class='channel-list'>"
        for j in active_now:
            elapsed = int(_t.time() - j["started"])
            stage = _describe_job_stage(log_text, j["video_id"])
            body += (
                '<li>'
                f'<span class="url"><strong>{h(j["video_id"])}</strong> · '
                f'{h(j["url"])}</span>'
                f'<span style="color: var(--muted); font-size: 12px;">'
                f'{elapsed//60}m {elapsed%60}s · {h(stage)}</span>'
                '</li>'
            )
        body += "</ul></section>"

        # Polling JS for the in-progress section — emitted here so both the
        # "no runs yet" early return and the full table path include it.
        active_poll_js = """
<script>
(function () {
  const section = document.getElementById('activity-active-section');
  const list = document.getElementById('activity-active-list');
  if (!section || !list) return;
  let prevCount = list.children.length;
  const fmtElapsed = s => `${Math.floor(s/60)}m ${s%60}s`;
  const esc = s => { const d = document.createElement('div'); d.textContent = String(s ?? ''); return d.innerHTML; };
  async function poll() {
    try {
      const res = await fetch('/one-off/status', { cache: 'no-store' });
      if (!res.ok) return;
      const data = await res.json();
      const active = data.active || [];
      section.style.display = active.length ? '' : 'none';
      list.innerHTML = active.map(j => `
        <li>
          <span class="url"><strong>${esc(j.video_id)}</strong> &middot; ${esc(j.url)}</span>
          <span style="color: var(--muted); font-size: 12px;">${esc(fmtElapsed(j.elapsed_secs))} &middot; ${esc(j.stage)}</span>
        </li>
      `).join('');
      if (active.length < prevCount) {
        window.location.reload();
        return;
      }
      prevCount = active.length;
    } catch (_) { /* try again */ }
  }
  poll();
  setInterval(poll, 2000);
})();
</script>
"""

        # Filter chips
        chip = lambda v, label, count: (
            f'<a href="/activity?status={v}" class="filter-chip'
            + (' active' if flt == v else '')
            + f'">{h(label)} <span class="filter-chip-count">({count})</span></a>'
        )
        all_runs = _recent_runs(limit=200)
        n_all = len(all_runs)
        n_ok = sum(1 for r in all_runs if r.get("success"))
        n_fail = n_all - n_ok
        body += "<div class='filter-row'>"
        body += chip("all", "All", n_all)
        body += chip("success", "Success", n_ok)
        body += chip("failed", "Failed", n_fail)
        body += "</div>"

        if not rows:
            body += "<p class='meta-info'>No runs recorded yet. Submit a one-off digest from the <a href='/one-off'>One-off digest</a> page.</p>"
            body += active_poll_js
            return page(body, title="Activity", current="activity")

        def _fmt_ago(secs: float) -> str:
            secs = int(secs)
            if secs < 60: return f"{secs}s ago"
            if secs < 3600: return f"{secs//60}m ago"
            if secs < 86400: return f"{secs//3600}h ago"
            return f"{secs//86400}d ago"

        def _fmt_dur(secs) -> str:
            if secs is None: return "—"
            secs = float(secs)
            if secs < 60: return f"{secs:.1f}s"
            return f"{int(secs)//60}m {int(secs)%60}s"

        def _fmt_int(n) -> str:
            if n is None: return "—"
            n = int(n)
            return f"{n:,}"

        now = _t.time()
        # Index usage by video_id for O(1) per-row cost lookup. Each row's
        # cost = sum of entries whose ts is within the row's [started_at,
        # ended_at] window AND video_id matches.
        usage_by_vid: dict = {}
        for e in usage_entries:
            usage_by_vid.setdefault(e.get("video_id", ""), []).append(e)

        body += "<table class='activity-table'>"
        body += (
            "<thead><tr>"
            "<th>When</th><th>Video</th><th>Outcome</th>"
            "<th>Duration</th><th>Stages</th><th>Tokens</th><th>Cost</th>"
            "</tr></thead><tbody>"
        )
        for r in rows:
            video_id = r.get("video_id") or ""
            url = r.get("url") or ""
            ago = _fmt_ago(now - (r.get("started_at") or now))
            dur = _fmt_dur(r.get("duration_secs"))
            success = r.get("success")
            if success:
                outcome = "<span class='ok'>✓ done</span>"
                if r.get("digest_path"):
                    title = f'<a href="/digests/{h(video_id)}/">{h(video_id)}</a>'
                else:
                    title = h(video_id)
            else:
                stage = r.get("stage_reached") or "?"
                outcome = f"<span class='fail'>✗ failed at {h(stage)}</span>"
                title = h(video_id)
            # Stage breakdown
            parts = []
            for label, key in [
                ("dl", "download_secs"),
                ("whisper", "whisper_secs"),
                ("frames", "frames_secs"),
                ("digest", "digest_secs"),
                ("vision", "vision_secs"),
            ]:
                v = r.get(key)
                if v is not None and v > 0.05:
                    parts.append(f"{label} {v:.1f}s")
            stages_cell = h(", ".join(parts)) if parts else "—"
            # Tokens (digest only)
            tin = r.get("digest_input_tokens")
            tout = r.get("digest_output_tokens")
            cache = r.get("digest_cache_read_tokens") or 0
            if tin or tout:
                tokens_cell = f"in {_fmt_int(tin)} · out {_fmt_int(tout)}"
                if cache:
                    tokens_cell += f" · cache {_fmt_int(cache)}"
                tokens_cell = h(tokens_cell)
            else:
                tokens_cell = "—"

            # Per-row cost: every LLM call recorded for this video_id whose
            # ts falls in the run's [started_at, ended_at] window.
            row_cost = 0.0
            row_backend = None
            started_at = r.get("started_at") or 0
            ended_at = r.get("ended_at") or now
            for e in usage_by_vid.get(video_id, []):
                ets = float(e.get("ts", 0) or 0)
                if started_at <= ets <= ended_at:
                    row_cost += float(e.get("cost_usd", 0.0) or 0.0)
                    row_backend = e.get("backend") or row_backend
            if row_cost > 0:
                cost_cell = f"${row_cost:.4f}"
            elif row_backend == "claude-code":
                cost_cell = "<span style='color: var(--muted);'>subscription</span>"
            else:
                cost_cell = "—"

            body += "<tr>"
            body += f"<td title='{h(url)}'>{h(ago)}</td>"
            body += f"<td>{title}"
            extras = []
            if r.get("source_lang"): extras.append(f"lang: {h(r['source_lang'])}")
            if r.get("used_whisper"):
                wm = r.get("whisper_model") or "?"
                extras.append(f"whisper: {h(wm)}")
            if extras:
                body += f"<div class='activity-meta'>{' · '.join(extras)}</div>"
            body += "</td>"
            body += f"<td>{outcome}"
            err = r.get("error")
            if err and not success:
                body += f"<div class='activity-meta activity-error'>{h(err)}</div>"
            body += "</td>"
            body += f"<td>{h(dur)}</td>"
            body += f"<td class='activity-stages'>{stages_cell}</td>"
            body += f"<td class='activity-tokens'>{tokens_cell}</td>"
            body += f"<td class='activity-cost'>{cost_cell}</td>"
            body += "</tr>"
        body += "</tbody></table>"
        body += (
            "<p class='meta-info' style='margin-top: 24px;'>"
            f"Raw log: <code>{h(str(_oneoff_log_path()))}</code><br>"
            f"JSONL: <code>{h(str(_runs_jsonl_path()))}</code>"
            "</p>"
        )
        body += active_poll_js
        return page(body, title="Activity", current="activity")

    @app.route("/digests/<video_id>/")
    def view_digest(video_id):
        from html import escape as h
        digest_md = digests_dir / video_id / "digest.md"
        if not digest_md.exists():
            abort(404)
        try:
            _mark_digest_read(video_id)
        except Exception:
            pass  # never block reading on a DB error
        md_source = digest_md.read_text()
        rendered = _render_markdown(md_source)

        # Build the action toolbar. Lives at the TOP of the page so the user
        # doesn't have to scroll past the whole digest to find these. Reading
        # order in the toolbar mirrors the typical "fast food" path: takeaway
        # first (the synthesized bottom line), panel second (the deeper
        # critique), then utility actions.
        panel_md = digests_dir / video_id / "panel.md"
        panel_exists = panel_md.exists()
        takeaway_md = digests_dir / video_id / "takeaway.md"
        takeaway_exists = takeaway_md.exists()
        gen_submit_js_panel = (
            "this.querySelector('button').disabled=true;"
            "this.querySelector('button').textContent='Generating panel… (~60–120s)';"
        )
        gen_submit_js_takeaway = (
            "this.querySelector('button').disabled=true;"
            "this.querySelector('button').textContent='Generating takeaway… (~30s)';"
        )

        toolbar = "<div class='digest-actions digest-toolbar'>"
        # Takeaway: primary CTA when present (it's what most readers want
        # next). When missing, offer to generate it.
        if takeaway_exists:
            toolbar += (
                f"<a class='discuss-btn' href='/digests/{h(video_id)}/takeaway/'>"
                "View takeaway</a>"
            )
        else:
            toolbar += (
                f"<form method='post' action='/digests/{h(video_id)}/takeaway' "
                f"style='display:inline;' onsubmit=\"{gen_submit_js_takeaway}\">"
                "<button type='submit' class='discuss-btn' "
                "title='Writes a 1–3 paragraph takeaway and saves it to takeaway.md'>"
                "Generate takeaway</button>"
                "</form>"
            )
        # Panel: secondary nav when present. Regenerate button removed —
        # without prompt editability there's no useful reason to spend
        # another Opus call. Legacy digests without a panel can still
        # generate one on demand (auto-pipeline didn't run for them).
        if panel_exists:
            toolbar += (
                f"<a class='discuss-btn-secondary' style='text-decoration:none;' "
                f"href='/digests/{h(video_id)}/panel/'>View panel discussion</a>"
            )
        else:
            toolbar += (
                f"<form method='post' action='/digests/{h(video_id)}/discuss' "
                f"style='display:inline;' onsubmit=\"{gen_submit_js_panel}\">"
                "<button type='submit' class='discuss-btn-secondary' "
                "title='Generates a panel-of-experts discussion (~60–120s, one Opus call)'>"
                "Generate panel</button>"
                "</form>"
            )
        # Slides — sibling artifact alongside takeaway/panel. The visual
        # layer (intelligently-picked frames + transcript snippets) is the
        # tool's main differentiator vs. a pure-LLM workflow. Generation
        # is async (daemon thread); the toolbar renders three states.
        slides_path = digests_dir / video_id / "slides.pptx"
        slides_job = local_job_status(f"{video_id}:slides")
        slides_running = slides_job.get("phase") == "running"
        if slides_path.exists():
            toolbar += (
                f"<a class='discuss-btn-secondary' style='text-decoration:none;' "
                f"href='/digests/{h(video_id)}/slides.pptx' "
                "title='Download the auto-generated PowerPoint deck "
                "(one slide per topic, with intelligently-picked frames).'>"
                "Download slides</a>"
            )
        elif slides_running:
            elapsed = slides_job.get("elapsed", 0)
            toolbar += (
                f"<span class='discuss-btn-secondary' style='cursor: default;' "
                f"data-poll-url='/digests/{h(video_id)}/job-status?kind=slides'>"
                f"Generating slides… <span class='elapsed'>{elapsed}s</span>"
                "</span>"
            )
        else:
            toolbar += (
                f"<form method='post' action='/digests/{h(video_id)}/slides' "
                f"style='display:inline;'>"
                "<button type='submit' class='discuss-btn-secondary' "
                "title='Builds slides.pptx from the cached video — local frame "
                "extraction + alignment + PowerPoint assembly. No LLM call. "
                "Runs in the background; this page will refresh when ready.'>"
                "Generate slides</button>"
                "</form>"
            )
        # Copy markdown — lets the reader paste into Notion, Obsidian, an
        # email, another LLM, etc. Markdown is the most-portable form.
        toolbar += (
            "<button type='button' class='discuss-btn-secondary' "
            "data-copy-target='digest-md-source' "
            "title='Copy the markdown source — paste into your notes app, email, or another LLM'>"
            "Copy markdown</button>"
        )
        toolbar += (
            f"<form method='post' action='/digests/{h(video_id)}/delete' style='display:inline;' "
            "onsubmit=\"return confirm('Delete this digest? "
            "The rendered output, frames, and cached video will be wiped.');\">"
            "<button type='submit' class='delete-btn' "
            "title='Wipes digest.md, frames, and cached video. Re-submit via One-off digest to regenerate.'>"
            "Delete digest</button>"
            "</form>"
        )
        toolbar += "</div>"

        # Hidden textarea holds the raw markdown for clipboard copy. <textarea>
        # treats content as raw text up to </textarea>, so escaping & and < is
        # sufficient (html.escape handles both).
        hidden_md = (
            f'<textarea id="digest-md-source" hidden aria-hidden="true">'
            f'{h(md_source)}</textarea>'
        )

        # Inject the polling JS only when there's actually a job in flight,
        # so the steady state has zero JS overhead beyond the existing copy
        # handler.
        poll_js = _JOB_POLL_JS if slides_running else ""
        body = (
            toolbar
            + hidden_md
            + _COPY_BUTTON_JS
            + poll_js
            + "<hr style='margin: 16px 0 32px; border: none; border-top: 1px solid var(--border);'>"
            + rendered
        )
        return page(body, title=video_id, current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/")

    @app.route("/digests/<video_id>/discuss", methods=["POST"])
    def generate_panel(video_id):
        from flask import redirect
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        digest_md = digests_dir / video_id / "digest.md"
        if not digest_md.exists():
            abort(404)
        # Find the cached SRT to feed into the panel prompt. The download dir
        # holds whichever language track was picked (en / zh-Hans / etc.).
        srt_dir = digests_dir / video_id / "downloads" / video_id
        srt_files = list(srt_dir.glob("*.srt")) if srt_dir.exists() else []
        if not srt_files:
            return redirect(
                f"/digests/{video_id}/?msg=No+cached+transcript.+Re-digest+the+video+first."
            )
        # Filename pattern: <video_id>.<lang>.srt — pull the lang back out so
        # the panel can be written in the same language as the digest.
        srt_path = srt_files[0]
        lang = srt_path.stem[len(video_id) + 1:] if srt_path.stem.startswith(video_id + ".") else "en"
        s = load_settings()
        try:
            segments = parse_srt(srt_path)
            backend = select_backend()
            panel_model = s.get("panel_model") or DEFAULT_PANEL_MODEL
            text, p_usage = generate_panel_discussion(
                digest_md.read_text(),
                segments,
                model=panel_model,
                source_lang=lang,
                output_language=s.get("digest_language") or "auto",
                backend=backend,
            )
            record_llm_usage(
                video_id=video_id, kind="panel", model=panel_model,
                backend_name=backend.name, usage=p_usage,
            )
        except Exception as e:
            return redirect(f"/digests/{video_id}/?msg=Panel+generation+failed:+{e}")
        (digests_dir / video_id / "panel.md").write_text(text)
        return redirect(f"/digests/{video_id}/panel/")

    @app.route("/digests/<video_id>/takeaway", methods=["POST"])
    def generate_takeaway_route(video_id):
        from flask import redirect
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        digest_md = digests_dir / video_id / "digest.md"
        if not digest_md.exists():
            abort(404)
        # Look up the cached SRT (same convention as the panel route).
        srt_dir = digests_dir / video_id / "downloads" / video_id
        srt_files = list(srt_dir.glob("*.srt")) if srt_dir.exists() else []
        if not srt_files:
            return redirect(
                f"/digests/{video_id}/?msg=No+cached+transcript.+Re-digest+the+video+first."
            )
        srt_path = srt_files[0]
        lang = (
            srt_path.stem[len(video_id) + 1:]
            if srt_path.stem.startswith(video_id + ".") else "en"
        )
        # Pull the panel text if it exists so the takeaway can weave its
        # critique into the synthesis; otherwise pass None and the prompt
        # copes.
        panel_md_path = digests_dir / video_id / "panel.md"
        panel_text = panel_md_path.read_text() if panel_md_path.exists() else None
        # We don't track the publish_date in saved metadata; pass None and
        # the takeaway will simply omit the "as of <date>" anchor.
        s = load_settings()
        try:
            segments = parse_srt(srt_path)
            backend = select_backend()
            takeaway_model = (os.environ.get("YT2MD_TAKEAWAY_MODEL")
                              or DEFAULT_TAKEAWAY_MODEL)
            takeaway_text, t_usage = generate_takeaway(
                digest_md.read_text(), panel_text, segments,
                model=takeaway_model,
                publish_date=None,
                source_lang=lang,
                output_language=s.get("digest_language") or "auto",
                backend=backend,
            )
            record_llm_usage(
                video_id=video_id, kind="takeaway", model=takeaway_model,
                backend_name=backend.name, usage=t_usage,
            )
            body = render_takeaway_markdown(
                takeaway_text,
                video_url=f"https://www.youtube.com/watch?v={video_id}",
            )
            (digests_dir / video_id / "takeaway.md").write_text(body)
        except Exception as e:
            return redirect(
                f"/digests/{video_id}/?msg=Takeaway+failed:+{type(e).__name__}:+{e}"
            )
        # Drop the user straight into the takeaway — that's what they
        # clicked Generate to read.
        return redirect(f"/digests/{video_id}/takeaway/")

    @app.route("/digests/<video_id>/takeaway/")
    def view_takeaway(video_id):
        from html import escape as h
        takeaway_md = digests_dir / video_id / "takeaway.md"
        if not takeaway_md.exists():
            abort(404)
        md_source = takeaway_md.read_text()
        rendered = _render_markdown(md_source)
        # Pre-assemble the chat handoff prompt server-side so the button can
        # just copy a hidden textarea (same pattern as Copy markdown). Loads
        # whichever artifacts exist (digest + panel + takeaway).
        chat_prompt = build_chat_handoff_prompt(video_id, digests_dir)
        nav = (
            '<div class="digest-actions digest-toolbar" style="margin-top:0;">'
            f'<a href="/digests/{h(video_id)}/" class="discuss-btn-secondary" '
            'style="text-decoration:none; padding: 6px 12px;">← Back to digest</a>'
            "<button type='button' class='discuss-btn' "
            "data-copy-target='chat-handoff-source' "
            "data-then-open='https://claude.ai/new' "
            "title='Copies digest + panel + takeaway to your clipboard and "
            "opens claude.ai in a new tab — paste to continue the discussion "
            "with full context.'>Continue in chat</button>"
            "<button type='button' class='discuss-btn-secondary' "
            "data-copy-target='takeaway-md-source' "
            "title='Copy the takeaway markdown'>Copy markdown</button>"
            "</div>"
        )
        hidden_md = (
            f'<textarea id="takeaway-md-source" hidden aria-hidden="true">'
            f'{h(md_source)}</textarea>'
            f'<textarea id="chat-handoff-source" hidden aria-hidden="true">'
            f'{h(chat_prompt)}</textarea>'
        )
        return page(nav + hidden_md + _COPY_BUTTON_JS + rendered,
                    title=f"Takeaway · {video_id}",
                    current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/takeaway/")

    @app.route("/digests/<video_id>/panel/")
    def view_panel(video_id):
        from html import escape as h
        panel_md = digests_dir / video_id / "panel.md"
        if not panel_md.exists():
            abort(404)
        md_source = panel_md.read_text()
        rendered = _render_markdown(md_source)
        nav = (
            f'<div class="digest-actions digest-toolbar" style="margin-top:0;">'
            f'<a href="/digests/{h(video_id)}/" class="discuss-btn-secondary" '
            f'style="text-decoration:none; padding: 6px 12px;">← Back to digest</a>'
            "<button type='button' class='discuss-btn-secondary' "
            "data-copy-target='panel-md-source' "
            "title='Copy the panel discussion markdown'>Copy markdown</button>"
            "</div>"
        )
        hidden_md = (
            f'<textarea id="panel-md-source" hidden aria-hidden="true">'
            f'{h(md_source)}</textarea>'
        )
        return page(nav + hidden_md + _COPY_BUTTON_JS + rendered,
                    title=f"Panel · {video_id}",
                    current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/panel/")

    @app.route("/digests/<video_id>/delete", methods=["POST"])
    def delete_digest(video_id):
        from flask import redirect
        target = digests_dir / video_id
        if not target.exists():
            abort(404)
        # Wipe artifacts + cached source media. shutil handles missing children.
        shutil.rmtree(target, ignore_errors=True)
        # Clear read-state row.
        try:
            with _library_connect() as conn:
                conn.execute("DELETE FROM digest_reads WHERE digest_id = ?", (video_id,))
        except Exception:
            pass
        return redirect("/?msg=Deleted+" + video_id)

    @app.route("/digests/<video_id>/digest_images/<path:filename>")
    def digest_image(video_id, filename):
        return send_from_directory(digests_dir / video_id / "digest_images", filename)

    @app.route("/digests/<video_id>/slides.pptx")
    def download_slides(video_id):
        slides_path = digests_dir / video_id / "slides.pptx"
        if not slides_path.exists():
            abort(404)
        return send_from_directory(
            digests_dir / video_id, "slides.pptx",
            as_attachment=True, download_name=f"{video_id}.pptx",
        )

    def _build_slides_for_digest(video_id: str) -> None:
        """Synchronous slides-only pipeline against cached MP4 + SRT.
        Raises on any failure; the caller (sync or async) is expected to
        catch and surface the error. Idempotent: writes slides.pptx atomically
        via a temp suffix swap so a partial write doesn't leave a corrupt file.
        """
        from concurrent.futures import ThreadPoolExecutor
        cache_dir = digests_dir / video_id / "downloads" / video_id
        if not cache_dir.exists():
            raise RuntimeError("No cached video. Re-digest the video first.")
        mp4_path = cache_dir / f"{video_id}.mp4"
        srt_files = list(cache_dir.glob("*.srt"))
        if not mp4_path.exists() or not srt_files:
            raise RuntimeError(
                "Cached video or SRT missing. Re-digest the video first."
            )
        srt_path = srt_files[0]
        workdir = Path(tempfile.mkdtemp(prefix="v2d_slides_"))
        scene_dir = workdir / "scene"
        interval_dir = workdir / "interval"
        try:
            duration = get_video_duration(mp4_path)
            with ThreadPoolExecutor(max_workers=2) as pool:
                scene_fut = pool.submit(
                    extract_scene_frames, mp4_path, scene_dir, 0.2,
                )
                interval_fut = pool.submit(
                    extract_interval_frames, mp4_path, interval_dir, 20.0, duration,
                )
                scene_frames = scene_fut.result()
                interval_frames = interval_fut.result()
            frames = merge_frames(scene_frames, interval_frames)
            frames = dedupe_frames(frames, 4)
            # Slide-aware filtering: global pHash cluster (talk-deck pattern)
            # then optional vision-LLM classifier via 3×3 grids. Same logic
            # as the CLI auto-pipeline path.
            deck_frames = global_phash_cluster(frames, distance=6)
            settings = load_settings()
            if (settings.get("slide_classification", True)
                    and len(deck_frames) > _GRID_CELLS):
                try:
                    backend = select_backend()
                    if getattr(backend, "vision_supported", False):
                        deck_frames = classify_slides_via_grids(
                            deck_frames, backend=backend,
                            model=settings.get("slide_classifier_model")
                                  or "claude-haiku-4-5-20251001",
                            workdir=workdir,
                            log_video_id=video_id,
                        )
                except Exception:
                    # Backend unavailable / call failed → use pHash-only set.
                    pass
            segments = parse_srt(srt_path)
            slides_data = assign_transcript_to_frames(deck_frames, segments, duration)
            # Title for the deck's title slide. Prefer the digest's H1 (the
            # original YouTube title) — falls back to the video_id if absent.
            digest_md_path = digests_dir / video_id / "digest.md"
            title = video_id
            if digest_md_path.exists():
                for line in digest_md_path.read_text().splitlines():
                    if line.startswith("# "):
                        title = line[2:].strip()
                        break
            # Write through a temp path then rename so the polling UI
            # never sees a half-written .pptx as "ready".
            tmp_out = digests_dir / video_id / "slides.pptx.tmp"
            build_deck(slides_data, tmp_out, title)
            tmp_out.replace(digests_dir / video_id / "slides.pptx")
        finally:
            shutil.rmtree(workdir, ignore_errors=True)

    @app.route("/digests/<video_id>/slides", methods=["POST"])
    def generate_slides_route(video_id):
        from flask import redirect
        # Slides already on disk — no-op, just bounce back.
        if (digests_dir / video_id / "slides.pptx").exists():
            return redirect(f"/digests/{video_id}/?msg=Slides+already+exist.")
        # Spawn the job and return immediately. The UI polls
        # /job-status?kind=slides for progress and reloads when done.
        start_local_job(f"{video_id}:slides", _build_slides_for_digest, video_id)
        return redirect(f"/digests/{video_id}/")

    @app.route("/digests/<video_id>/job-status")
    def job_status_route(video_id):
        from flask import jsonify, request
        kind = request.args.get("kind", "slides")
        snap = local_job_status(f"{video_id}:{kind}")
        # Surface artifact presence so the polling UI can decide whether to
        # reload to the success state or show "missing" after a clean exit.
        if kind == "slides":
            snap["artifact_exists"] = (
                (digests_dir / video_id / "slides.pptx").exists()
            )
        return jsonify(snap)

    @app.errorhandler(404)
    def not_found(e):
        return page(
            "<h1>Not found</h1><p>That digest doesn't exist yet.</p>",
            title="404", current="home",
        ), 404

    url = f"http://127.0.0.1:{args.port}/"
    print(f"yt2md reader: {url}")
    print(f"Data dir: {data_dir}")
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if api_key:
        print(f"API key: set ({api_key[:7]}…{api_key[-4:]})")
    else:
        print(f"API key: (not set)")
    # Cheap sentinel-based probe to populate the session-state cache so the
    # banner / gate know whether Claude Code is logged in across restarts.
    claude_probe_login_state()
    cc_status = (
        "logged in" if _claude_code_session_state.get("logged_in")
        else ("installed (run /setup to log in)" if claude_code_installed()
              else "not installed")
    )
    print(f"Claude Code: {cc_status}")
    if not api_key and not _claude_code_session_state.get("logged_in"):
        print(f"  → first-run setup at {url}setup")
    # Surface the YouTube prerequisites at startup so the user notices missing
    # cookies / JS runtime before a one-off digest fails 30s later.
    cookies_browser = os.environ.get("YT2MD_COOKIES_FROM_BROWSER")
    print(f"Cookies: {cookies_browser or '(none — set YT2MD_COOKIES_FROM_BROWSER for paywalled YouTube)'}")
    js_rt = _ensure_js_runtime_available()
    print(f"JS runtime: {js_rt or '(not found — yt-dlp n-challenge will fail; install deno or node)'}")
    _cleanup_legacy_launchd()
    start_scheduler()
    print("Press Ctrl-C to stop.\n")

    if not args.no_browser:
        import webbrowser
        import threading
        threading.Timer(0.5, lambda: webbrowser.open(url)).start()

    app.run(host="127.0.0.1", port=args.port, debug=False, use_reloader=False)
    return 0


def cmd_doctor(args) -> int:
    """Diagnose prerequisites and config. Prints check/X per item with fix hints.

    Exits 0 if everything looks usable, 1 if any blocking issue was found.
    Designed to be the first thing a new user runs after install — gives a
    concrete punch list instead of failing mid-pipeline 30s into the first
    digest.
    """
    OK, WARN, FAIL = "\033[32m✓\033[0m", "\033[33m!\033[0m", "\033[31m✗\033[0m"
    blocking: list = []
    advisory: list = []

    def ok(msg: str) -> None:
        print(f"  {OK} {msg}")

    def warn(msg: str, hint: Optional[str] = None) -> None:
        print(f"  {WARN} {msg}")
        if hint:
            print(f"      → {hint}")
        advisory.append(msg)

    def fail(msg: str, hint: Optional[str] = None) -> None:
        print(f"  {FAIL} {msg}")
        if hint:
            print(f"      → {hint}")
        blocking.append(msg)

    print("\nyt2md doctor — checking prerequisites and config\n")

    print("System tools:")
    for tool in ("ffmpeg", "ffprobe"):
        if shutil.which(tool):
            ok(f"{tool} on PATH")
        else:
            fail(f"{tool} not found",
                 "macOS: `brew install ffmpeg` · Linux: `apt install ffmpeg`")

    if shutil.which("uv"):
        try:
            v = subprocess.run(["uv", "--version"], capture_output=True, text=True,
                               timeout=5).stdout.strip()
            ok(f"uv ({v})")
        except Exception:
            ok("uv on PATH")
    else:
        warn("uv not on PATH",
             "Recommended: install via `curl -LsSf https://astral.sh/uv/install.sh | sh`")

    rt = _ensure_js_runtime_available()
    if rt:
        rt_path = shutil.which(rt) or "(unknown)"
        try:
            ver = subprocess.run([rt, "--version"], capture_output=True, text=True,
                                 timeout=5).stdout.strip()
        except Exception:
            ver = "?"
        ok(f"JS runtime: {rt} {ver} ({rt_path})")
    else:
        fail("No JS runtime found (needed for yt-dlp's n-challenge solver)",
             "`brew install deno` (simplest) or `nvm install 20`")

    print("\nPython packages:")
    try:
        import yt_dlp  # type: ignore
        ok(f"yt-dlp {yt_dlp.version.__version__}")
    except ImportError:
        fail("yt-dlp not installed", "Run `uv sync` from the project directory")
    try:
        import faster_whisper  # type: ignore  # noqa: F401
        ok("faster-whisper installed")
    except ImportError:
        warn("faster-whisper not installed",
             "Whisper fallback won't work for captionless videos. `uv sync` to install.")
    try:
        import anthropic  # type: ignore  # noqa: F401
        ok("anthropic SDK installed")
    except ImportError:
        fail("anthropic SDK not installed", "Run `uv sync`")

    print("\nAPI / auth:")
    load_env_files()
    key = os.environ.get("ANTHROPIC_API_KEY", "")
    if key:
        ok(f"ANTHROPIC_API_KEY set ({key[:10]}…{key[-4:]})")
    else:
        fail("ANTHROPIC_API_KEY not set",
             f"Get a key at https://console.anthropic.com/settings/keys; first run "
             f"of yt2md will prompt and save it to {get_data_dir() / '.env'}")

    settings = load_settings()
    cookies = settings.get("cookies_from_browser") or os.environ.get("YT2MD_COOKIES_FROM_BROWSER", "")
    if cookies:
        ok(f"YouTube cookies: from browser '{cookies}'")
    else:
        warn("YouTube cookies not configured",
             "Many videos now require login. Set in /settings or as "
             "YT2MD_COOKIES_FROM_BROWSER=firefox in ~/yt2md/.env")

    print("\nConfig:")
    data_dir = get_data_dir()
    print(f"  data dir: {data_dir}")
    if _settings_file().exists():
        ok(f"settings.json exists")
    else:
        warn("settings.json not yet created — defaults in effect",
             f"Open http://localhost:7682/settings to configure (after `yt2md serve`)")
    print(f"  digest model: {settings.get('digest_model')}")
    print(f"  panel model:  {settings.get('panel_model')}")
    print(f"  whisper model: {settings.get('whisper_model')}")
    print(f"  digest language: {settings.get('digest_language')}")

    cfg = load_schedule_config()
    print(f"  schedule: {_format_schedule_summary(cfg)}")

    digests_dir = data_dir / "digests"
    n_digests = len([d for d in digests_dir.iterdir() if (d / "digest.md").exists()]) if digests_dir.exists() else 0
    n_channels = len(read_channels())
    print(f"  library: {n_digests} digest(s), {n_channels} channel(s)")

    print()
    if blocking:
        print(f"{FAIL} {len(blocking)} blocking issue{'s' if len(blocking) != 1 else ''}; "
              "fix the items marked above.")
        if advisory:
            print(f"{WARN} {len(advisory)} advisory item{'s' if len(advisory) != 1 else ''} "
                  "(non-blocking but worth setting up).")
        return 1
    if advisory:
        n = len(advisory)
        print(f"{OK} Core requirements met. {n} advisory item"
              f"{'s' if n != 1 else ''} {'are' if n != 1 else 'is'} optional.")
    else:
        print(f"{OK} All checks passed. Run `yt2md serve` to start the reader.")
    return 0


# ---- subcommand dispatcher ----

def _subcommand_main(argv: List[str]) -> int:
    """Handle yt2md {watch,serve} ..."""
    ap = argparse.ArgumentParser(prog="yt2md", description="yt2md subcommands")
    sub = ap.add_subparsers(dest="cmd", required=True)

    watch = sub.add_parser("watch", help="Manage watched channels and run polling")
    watch_sub = watch.add_subparsers(dest="watch_cmd", required=True)
    p = watch_sub.add_parser("add", help="Add a channel URL"); p.add_argument("url")
    p.set_defaults(func=cmd_watch_add)
    p = watch_sub.add_parser("list", help="List watched channels"); p.set_defaults(func=cmd_watch_list)
    p = watch_sub.add_parser("remove", help="Remove a channel URL"); p.add_argument("url")
    p.set_defaults(func=cmd_watch_remove)
    p = watch_sub.add_parser("run", help="Poll all channels and digest new videos")
    p.set_defaults(func=cmd_watch_run)

    serve = sub.add_parser("serve", help="Start a local web reader (also runs the in-process scheduler)")
    serve.add_argument("--port", type=int, default=7682, help="Port (default: 7682)")
    serve.add_argument("--no-browser", action="store_true",
                       help="Don't auto-open a browser tab on start")
    serve.set_defaults(func=cmd_serve)

    doctor = sub.add_parser("doctor", help="Check prerequisites and config; print a punch list")
    doctor.set_defaults(func=cmd_doctor)

    args = ap.parse_args(argv)
    return args.func(args)


# ---------- Main ----------

def main():
    # Subcommand dispatch — short-circuit the single-video flow when the user
    # invokes yt2md watch / serve / doctor.
    if len(sys.argv) > 1 and sys.argv[1] in ("watch", "serve", "doctor"):
        load_env_files()
        sys.exit(_subcommand_main(sys.argv[1:]))

    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("video", help="Input MP4 file path OR a YouTube URL (auto-downloads mp4 + SRT)")
    ap.add_argument("srt", nargs="?", default=None,
                    help="SRT transcript file (required if 'video' is a local path; "
                         "ignored for URLs — fetched automatically)")
    ap.add_argument("-o", "--output", type=Path, default=None, metavar="PATH",
                    help="Digest output path (default: <video-name>_digest.md)")
    ap.add_argument("--no-slides", action="store_true",
                    help="Skip the PowerPoint deck. By default a slides.pptx file is "
                         "written alongside the digest — the visual layer (intelligently-"
                         "selected frames + transcript snippets) is the tool's main "
                         "differentiator and worth shipping for every digest.")
    ap.add_argument("--no-slide-classification", action="store_true",
                    help="When building slides, skip the vision-LLM classifier that filters "
                         "raw frames down to actual deck slides. Falls back to pHash dedup "
                         "only — cheaper but typically produces 2-3x more (mostly redundant) "
                         "slides for talks where the speaker has a deck.")
    ap.add_argument("--slide-classifier-model",
                    default=os.environ.get("YT2MD_SLIDE_CLASSIFIER_MODEL"),
                    help="Vision model for the slide-classification step (default: "
                         "claude-haiku-4-5-20251001 — cheap and accurate enough for the task).")
    ap.add_argument("--deck", nargs="?", const="__default__", default=None, metavar="PATH",
                    help="Override the slides output path (default: <digest_dir>/slides.pptx).")
    ap.add_argument("--deck-only", action="store_true",
                    help="Skip the digest entirely — only build the deck. No API key needed.")
    ap.add_argument("--no-vision", action="store_true",
                    help="Disable vision-based frame picking for the digest. Cheaper but the "
                         "frames may be less illustrative.")
    ap.add_argument("--no-panel", action="store_true",
                    help="Skip the panel-of-experts discussion (saves ~1 Opus call per video). "
                         "Distillation will still run but without panel-informed confidence tags.")
    ap.add_argument("--no-takeaway", action="store_true",
                    help="Skip the takeaway step (synthesis prose) appended to digest.md.")
    ap.add_argument("--digest-model",
                    default=os.environ.get("YT2MD_DIGEST_MODEL") or "claude-sonnet-4-6",
                    help="Claude model for the digest (default: claude-sonnet-4-6). "
                         "Use claude-opus-4-7 for the highest-quality summarization.")
    ap.add_argument("--panel-model",
                    default=os.environ.get("YT2MD_PANEL_MODEL") or DEFAULT_PANEL_MODEL,
                    help=f"Claude model for the panel discussion (default: {DEFAULT_PANEL_MODEL}). "
                         "Multi-perspective synthesis benefits from a stronger model.")
    ap.add_argument("--takeaway-model",
                    default=os.environ.get("YT2MD_TAKEAWAY_MODEL") or DEFAULT_TAKEAWAY_MODEL,
                    help=f"Claude model for the takeaway step (default: "
                         f"{DEFAULT_TAKEAWAY_MODEL}).")
    ap.add_argument("--scene-threshold", type=float, default=0.2,
                    help="Scene-detection sensitivity, 0.1=lots of frames, 0.5=only major changes (default: 0.2)")
    ap.add_argument("--interval", type=float, default=20.0,
                    help="Also sample one frame every N seconds (0 to disable). Useful for "
                         "screen recordings where gradual changes don't trip scene detection. (default: 20)")
    ap.add_argument("--hash-distance", type=int, default=4,
                    help="Perceptual hash dedup threshold; lower = stricter. Compared only against "
                         "the previous kept frame, so recurring views are preserved. (default: 4)")
    ap.add_argument("--keep-frames", action="store_true",
                    help="Keep extracted frames in ./frames_<videoname>/ instead of cleaning up")
    ap.add_argument("--downloads-dir", type=Path, default=Path("downloads"),
                    help="Where to cache YouTube downloads (default: ./downloads)")
    ap.add_argument("--source-lang", default=None, metavar="CODE",
                    help="BCP-47 language code of a local SRT (e.g. 'zh-Hans'). "
                         "Drives the digest's output language when --digest-language=auto. "
                         "Ignored for URLs — yt-dlp / Whisper picks the track and the "
                         "lang is read from the file.")
    ap.add_argument("--digest-language",
                    default=os.environ.get("YT2MD_DIGEST_LANGUAGE") or "auto",
                    choices=["auto", "en"],
                    help="Output language for the digest + panel discussion. "
                         "'auto' (default) writes in the transcript's language. "
                         "'en' forces English regardless of source.")
    ap.add_argument("--whisper-model",
                    default=os.environ.get("YT2MD_WHISPER_MODEL") or DEFAULT_WHISPER_MODEL,
                    choices=["tiny", "base", "small", "medium", "large-v2", "large-v3"],
                    help=f"faster-whisper model used as fallback when a YouTube "
                         f"video has no captions (default: {DEFAULT_WHISPER_MODEL}).")
    ap.add_argument("--no-whisper", action="store_true",
                    help="Disable Whisper fallback; fail when a video has no captions.")
    ap.add_argument("--cookies-from-browser", default=os.environ.get("YT2MD_COOKIES_FROM_BROWSER"),
                    choices=["chrome", "firefox", "safari", "brave", "edge",
                             "chromium", "opera", "vivaldi"],
                    metavar="BROWSER",
                    help="Pass cookies from this browser to yt-dlp. Required when "
                         "YouTube returns 'Sign in to confirm you're not a bot'. "
                         "Defaults to $YT2MD_COOKIES_FROM_BROWSER if set.")
    args = ap.parse_args()

    load_env_files()

    for tool in ("ffmpeg", "ffprobe"):
        if not shutil.which(tool):
            sys.exit(f"{tool} not found on PATH. Install with `brew install ffmpeg` or your package manager.")

    do_digest = not args.deck_only
    if do_digest:
        ensure_api_key()

    import time as _time
    timings: dict = {}
    fetch_meta: dict = {
        "used_whisper": False, "whisper_model": None,
    }
    video_title: Optional[str] = None
    video_url: Optional[str] = None
    upload_date: Optional[str] = None
    # Run-start timestamp lets the summary block aggregate cost from the
    # usage log without needing to thread a list through every call site
    # (slide_classifier in particular records from inside its function).
    _run_start_ts = _time.time()

    if is_url(args.video):
        print(f"[0/5] Fetching YouTube video: {args.video}")
        result = fetch_youtube(
            args.video,
            args.downloads_dir,
            whisper_model=args.whisper_model,
            allow_whisper=not args.no_whisper,
            cookies_from_browser=args.cookies_from_browser,
        )
        video_path = result["mp4"]
        srt_path = result["srt"]
        source_lang = result["lang"]
        video_title = result.get("title")
        video_url = result.get("webpage_url")
        upload_date = result.get("upload_date")
        timings["download"] = round(result["download_secs"], 3)
        timings["whisper"] = round(result["whisper_secs"], 3)
        fetch_meta["used_whisper"] = result["used_whisper"]
        fetch_meta["whisper_model"] = result["whisper_model"]
        print(f"      mp4: {video_path}")
        print(f"      srt: {srt_path} (lang: {source_lang})")
    else:
        video_path = Path(args.video)
        if not video_path.exists():
            sys.exit(f"Video not found: {video_path}")
        if args.srt is None:
            sys.exit("SRT path is required when 'video' is a local file.")
        srt_path = Path(args.srt)
        if not srt_path.exists():
            sys.exit(f"SRT not found: {srt_path}")
        # Local-file path: caller didn't tell us the language; assume English.
        # Override with --source-lang if you're digesting a non-English local SRT.
        source_lang = args.source_lang or "en"

    base = video_path.stem
    digest_path = args.output if args.output is not None else Path(f"{base}_digest.md")
    # Slides default-on: write `slides.pptx` next to digest.md so the
    # web reader (and any KB ingester) finds it at a predictable path.
    # --no-slides opts out; --deck path overrides; --deck-only forces it
    # on even when --no-slides was passed (the user explicitly asked
    # for the deck).
    if args.deck and args.deck != "__default__":
        deck_path: Optional[Path] = Path(args.deck)
    elif args.deck_only or not args.no_slides:
        deck_path = digest_path.parent / "slides.pptx"
    else:
        deck_path = None

    workdir = Path(tempfile.mkdtemp(prefix="v2d_"))
    scene_dir = workdir / "scene"
    interval_dir = workdir / "interval"

    try:
        duration = get_video_duration(video_path)

        # Run scene detection and interval sampling concurrently — each is
        # an independent ffmpeg pass that decodes the full video, so doing
        # them in parallel cuts wall time roughly in half on slide-heavy
        # talks where scene detection is the slow leg. Each subprocess gets
        # its own ffmpeg process; CPU contention is real but small.
        from concurrent.futures import ThreadPoolExecutor
        print(f"[1/5] Extracting frames "
              f"(scene threshold={args.scene_threshold}, interval={args.interval}s, "
              f"in parallel)...")
        _frames_t0 = _time.monotonic()
        with ThreadPoolExecutor(max_workers=2) as pool:
            scene_fut = pool.submit(
                extract_scene_frames, video_path, scene_dir, args.scene_threshold,
            )
            interval_fut = pool.submit(
                extract_interval_frames, video_path, interval_dir, args.interval, duration,
            )
            _scene_t0 = _time.monotonic()
            scene_frames = scene_fut.result()
            scene_secs = _time.monotonic() - _scene_t0
            interval_frames = interval_fut.result()
        timings["frames_extract"] = round(_time.monotonic() - _frames_t0, 3)
        frames = merge_frames(scene_frames, interval_frames)
        cap_note = (
            f" (capped from >={SCENE_FRAME_HARD_CAP})"
            if len(scene_frames) == SCENE_FRAME_HARD_CAP else ""
        )
        print(f"      {len(scene_frames)} scene{cap_note} + "
              f"{len(interval_frames)} interval = {len(frames)} candidate frames "
              f"({timings['frames_extract']}s)")

        print(f"[2/5] Deduping consecutive near-identical frames "
              f"(hash distance <= {args.hash_distance})...")
        _dedupe_t0 = _time.monotonic()
        frames = dedupe_frames(frames, args.hash_distance)
        timings["frames_dedupe"] = round(_time.monotonic() - _dedupe_t0, 3)
        print(f"      {len(frames)} unique frames ({timings['frames_dedupe']}s)")
        timings["frames"] = round(_time.monotonic() - _frames_t0, 3)

        print(f"[3/5] Parsing SRT: {srt_path.name}")
        segments = parse_srt(srt_path)
        print(f"      {len(segments)} transcript segments")

        print("[4/5] Aligning transcript to frames...")

        if deck_path is not None:
            # Build a slides-specific frame set: pHash cluster globally to
            # collapse "speaker → slide → speaker → same slide" into one
            # representative, then (optionally) ask a vision LLM to filter
            # to actual deck slides only. The digest's vision-pick step
            # still uses the richer `frames` pool — the two consumers want
            # different shapes of the same data.
            deck_frames = global_phash_cluster(frames, distance=6)
            print(f"      slides: {len(frames)} candidates → "
                  f"{len(deck_frames)} after global pHash dedup")
            settings = load_settings()
            slide_classify_enabled = (
                settings.get("slide_classification", True)
                and not args.no_slide_classification
            )
            if slide_classify_enabled and len(deck_frames) > _GRID_CELLS:
                try:
                    classifier_backend = select_backend()
                    if getattr(classifier_backend, "vision_supported", False):
                        _classify_t0 = _time.monotonic()
                        deck_frames = classify_slides_via_grids(
                            deck_frames,
                            backend=classifier_backend,
                            model=(args.slide_classifier_model
                                   or settings.get("slide_classifier_model")
                                   or "claude-haiku-4-5-20251001"),
                            workdir=workdir,
                            log_video_id=video_path.stem,
                        )
                        timings["slide_classify"] = round(
                            _time.monotonic() - _classify_t0, 3,
                        )
                        print(f"      slides: vision-classified → "
                              f"{len(deck_frames)} kept "
                              f"({timings['slide_classify']}s)")
                except Exception as e:
                    print(f"      slides: classification skipped "
                          f"({type(e).__name__}: {e}); using pHash dedup only.")
            slides_data = assign_transcript_to_frames(
                deck_frames, segments, duration,
            )
            print(f"[5/5] Building slides ({len(slides_data)} slides) -> {deck_path}")
            build_deck(slides_data, deck_path, video_title or video_path.stem)
        else:
            print("[5/5] Slides skipped (--no-slides)")
            # Still need slides_data for downstream code if any consumes it,
            # though currently only build_deck does. Use the rich pool.
            slides_data = assign_transcript_to_frames(frames, segments, duration)

        usage = None
        if do_digest:
            digest_path.parent.mkdir(parents=True, exist_ok=True)
            images_dir = digest_path.parent / f"{digest_path.stem}_images"
            backend = select_backend()
            log_video_id = video_path.stem
            print(f"[+] Generating digest with {args.digest_model} via {backend.name} backend -> {digest_path}")
            _digest_t0 = _time.monotonic()
            digest, usage = generate_digest(
                segments, video_title or video_path.stem, args.digest_model,
                source_lang=source_lang,
                output_language=args.digest_language,
                backend=backend,
            )
            digest_log_entry = record_llm_usage(
                video_id=log_video_id, kind="digest", model=args.digest_model,
                backend_name=backend.name, usage=usage,
            )
            timings["digest"] = round(_time.monotonic() - _digest_t0, 3)
            print(f"      {len(digest.topics)} topics  |  "
                  f"input: {usage.input_tokens} tokens "
                  f"(cache read: {getattr(usage, 'cache_read_input_tokens', 0)}, "
                  f"cache write: {getattr(usage, 'cache_creation_input_tokens', 0)})  |  "
                  f"output: {usage.output_tokens} tokens  |  "
                  f"cost: ${digest_log_entry['cost_usd']:.4f}")

            vision_picks = None
            if not args.no_vision:
                if not getattr(backend, "vision_supported", False):
                    print(f"[+] Vision skipped — {backend.name} backend has vision disabled "
                          "(timestamp-based picks will be used).")
                else:
                    print(f"[+] Vision-picking frames with {args.digest_model}...")
                    _vision_t0 = _time.monotonic()
                    try:
                        vision_picks, v_usage = vision_pick_frames(
                            digest, frames, duration, args.digest_model,
                            segments=segments, backend=backend,
                        )
                        v_log_entry = record_llm_usage(
                            video_id=log_video_id, kind="vision_pick",
                            model=args.digest_model,
                            backend_name=backend.name, usage=v_usage,
                        )
                        timings["vision"] = round(_time.monotonic() - _vision_t0, 3)
                        print(f"      vision-selected {len(vision_picks)}/{len(digest.topics)} topics  |  "
                              f"input: {v_usage.input_tokens} tokens  |  "
                              f"output: {v_usage.output_tokens} tokens  |  "
                              f"cost: ${v_log_entry['cost_usd']:.4f}")
                    except VisionUnsupported as e:
                        print(f"      vision unsupported: {e}; falling back to timestamp picks.")

            write_markdown_digest(
                digest, frames, duration, digest_path, images_dir, vision_picks,
                video_title=video_title, video_url=video_url,
            )
            print(f"      Digest written. Images in {images_dir}/")

            # Render the digest to markdown text once for the panel + takeaway
            # prompts (saves a re-read on each step). For the panel/takeaway we
            # want the source-of-truth digest the user will see, so read back
            # the just-written file rather than reconstructing from `digest`.
            digest_md_text = digest_path.read_text()

            panel_md_text: Optional[str] = None
            panel_path = digest_path.parent / "panel.md"
            if not args.no_panel:
                print(f"[+] Generating panel discussion with {args.panel_model}...")
                _panel_t0 = _time.monotonic()
                try:
                    panel_md_text, p_usage = generate_panel_discussion(
                        digest_md_text, segments, model=args.panel_model,
                        source_lang=source_lang,
                        output_language=args.digest_language,
                        backend=backend,
                    )
                    p_log_entry = record_llm_usage(
                        video_id=log_video_id, kind="panel",
                        model=args.panel_model,
                        backend_name=backend.name, usage=p_usage,
                    )
                    panel_path.write_text(panel_md_text)
                    timings["panel"] = round(_time.monotonic() - _panel_t0, 3)
                    print(f"      Panel written -> {panel_path}  |  "
                          f"input: {p_usage.input_tokens} tokens  |  "
                          f"output: {p_usage.output_tokens} tokens  |  "
                          f"cost: ${p_log_entry['cost_usd']:.4f}")
                except Exception as e:
                    # Takeaway can still run without a panel — just with less
                    # explicit pushback to weave in. Don't take the whole
                    # pipeline down for one downstream step's failure.
                    print(f"      Panel generation failed ({type(e).__name__}: {e}); "
                          "continuing without panel.")

            takeaway_path = digest_path.parent / "takeaway.md"
            if not args.no_takeaway:
                print(f"[+] Generating takeaway with {args.takeaway_model}...")
                _take_t0 = _time.monotonic()
                try:
                    takeaway_text, t_usage = generate_takeaway(
                        digest_md_text, panel_md_text, segments,
                        model=args.takeaway_model,
                        publish_date=upload_date,
                        source_lang=source_lang,
                        output_language=args.digest_language,
                        backend=backend,
                    )
                    t_log_entry = record_llm_usage(
                        video_id=log_video_id, kind="takeaway",
                        model=args.takeaway_model,
                        backend_name=backend.name, usage=t_usage,
                    )
                    body = render_takeaway_markdown(
                        takeaway_text, video_url=video_url,
                    )
                    takeaway_path.write_text(body)
                    timings["takeaway"] = round(_time.monotonic() - _take_t0, 3)
                    print(f"      Takeaway written -> {takeaway_path}  |  "
                          f"input: {t_usage.input_tokens} tokens  |  "
                          f"output: {t_usage.output_tokens} tokens  |  "
                          f"cost: ${t_log_entry['cost_usd']:.4f}")
                except Exception as e:
                    print(f"      Takeaway failed ({type(e).__name__}: {e}); "
                          "digest is still complete without it.")

        if args.keep_frames:
            dest = Path.cwd() / f"frames_{video_path.stem}"
            dest.mkdir(parents=True, exist_ok=True)
            for d in (scene_dir, interval_dir):
                if d.exists():
                    shutil.copytree(d, dest, dirs_exist_ok=True)
            print(f"      frames saved to {dest}")

        outputs = []
        if do_digest:
            outputs.append(str(digest_path))
        if deck_path is not None:
            outputs.append(str(deck_path))
        print(f"\nDone. Wrote: {', '.join(outputs)}")

        # Cost audit: aggregate every LLM call recorded during this run.
        # Reads from the canonical usage log so it picks up records made
        # inside helper functions (e.g. slide_classifier) too.
        try:
            log_video_id = video_path.stem
            run_log_entries = [
                e for e in read_llm_usage_log()
                if e.get("ts", 0) >= _run_start_ts
                and e.get("video_id") == log_video_id
            ]
            costs_by_kind: dict = {}
            for e in run_log_entries:
                costs_by_kind[e["kind"]] = (
                    costs_by_kind.get(e["kind"], 0.0) + float(e.get("cost_usd", 0.0))
                )
            total_cost = round(sum(costs_by_kind.values()), 4)
            costs_by_kind = {k: round(v, 4) for k, v in costs_by_kind.items()}
            backend_used = (
                run_log_entries[0]["backend"] if run_log_entries else None
            )
        except Exception:
            costs_by_kind, total_cost, backend_used = {}, 0.0, None

        # Print a human-readable cost summary above the structured line.
        if costs_by_kind:
            print(f"\nCost summary (backend: {backend_used}):")
            for kind in ("digest", "vision_pick", "panel", "takeaway",
                         "slide_classifier", "validation"):
                if kind in costs_by_kind:
                    print(f"  {kind:18s} ${costs_by_kind[kind]:.4f}")
            note = " (subscription — no per-call billing)" if backend_used == "claude-code" else ""
            print(f"  {'TOTAL':18s} ${total_cost:.4f}{note}")

        # Structured one-line summary parsed by the web reaper. Keep this on
        # one line and as the LAST thing printed on success.
        summary = {
            "source_lang": source_lang,
            "used_whisper": fetch_meta["used_whisper"],
            "whisper_model": fetch_meta["whisper_model"],
            "timings": timings,
            "tokens": {
                "input": getattr(usage, "input_tokens", None) if usage else None,
                "output": getattr(usage, "output_tokens", None) if usage else None,
                "cache_read": getattr(usage, "cache_read_input_tokens", None) if usage else None,
                "cache_creation": getattr(usage, "cache_creation_input_tokens", None) if usage else None,
            } if usage else None,
            "cost": {
                "total_usd": total_cost,
                "by_kind": costs_by_kind,
                "backend": backend_used,
            },
            "digest_path": str(digest_path) if do_digest else None,
        }
        print("[summary] " + json.dumps(summary))
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


if __name__ == "__main__":
    main()
