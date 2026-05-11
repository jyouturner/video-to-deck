#!/usr/bin/env python3
"""Compare a generated slide deck against a hand-curated reference deck.

Used as a manual regression check: after changing anything in the slide
pipeline (frame extraction, dedup, the vision-LLM classifier), regenerate
the deck for a known-good video and run this script to verify nothing
critical was dropped or added.

The check is image-based, not structural. For each slide in the REFERENCE
deck, we compute its perceptual hash and find the nearest slide in the
CANDIDATE deck. A "match" means pHash Hamming distance ≤ threshold.

Outputs:
  - recall:    (matched reference slides) / (total reference slides)
  - precision: (matched candidate slides) / (total candidate slides)
  - missing:   reference slides with no good match in candidate
  - extra:     candidate slides with no good match in reference

This isn't an eval (we only have one curated example, not a dataset) —
it's a regression test for a single high-signal case the maintainer has
already inspected and trusts as ground truth.

Usage:
    uv run python tests/compare_slides.py REFERENCE.pptx CANDIDATE.pptx
    uv run python tests/compare_slides.py tests/fixtures/igO8iyca2_g.pptx ~/yt2md/digests/igO8iyca2_g/slides.pptx

Exits 0 if recall ≥ threshold (default 0.90), non-zero otherwise — so it
can be wired into a CI gate later if useful.
"""
from __future__ import annotations

import argparse
import io
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

# Imports gated to runtime so `python -h` works without the project venv.
def _imports():
    from pptx import Presentation
    import imagehash
    from PIL import Image
    return Presentation, imagehash, Image


def _slide_image_and_label(slide) -> Tuple[Optional[bytes], str]:
    """Pull the first picture's bytes + a human label (from text-box content,
    typically the "MM:SS – MM:SS | Slide N" footer or a slide title)."""
    image_blob = None
    label_parts: List[str] = []
    for shape in slide.shapes:
        if shape.shape_type == 13 and image_blob is None:  # PICTURE
            try:
                image_blob = shape.image.blob
            except Exception:
                pass
        if shape.has_text_frame and shape.text_frame.text.strip():
            label_parts.append(shape.text_frame.text.strip().split("\n")[0])
    label = " | ".join(label_parts)[:80]
    return image_blob, label


def _hash_deck(path: Path):
    """Return [(slide_index, label, phash, raw_image_bytes)] for content
    slides (skips the title slide and any slide without an image)."""
    Presentation, imagehash, Image = _imports()
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        blob, label = _slide_image_and_label(slide)
        if blob is None:
            continue
        try:
            with Image.open(io.BytesIO(blob)) as im:
                im.load()
                h = imagehash.phash(im)
        except Exception:
            continue
        out.append((i, label, h, blob))
    return out


def _short_label(label: str) -> str:
    """Pull the "Slide N" / time-range portion out of a footer label so the
    output table is scannable."""
    m = re.search(r"(\d+:\d+\s*[-–]\s*\d+:\d+(?:\s*\|\s*Slide\s*\d+)?)", label)
    return m.group(1) if m else label[:60]


def compare(reference_path: Path, candidate_path: Path, *,
            match_threshold: int = 10) -> dict:
    ref = _hash_deck(reference_path)
    cand = _hash_deck(candidate_path)

    print(f"reference:  {reference_path.name}  ({len(ref)} content slides)")
    print(f"candidate:  {candidate_path.name}  ({len(cand)} content slides)")
    print(f"match threshold (pHash Hamming distance ≤): {match_threshold}\n")

    if not ref:
        print("Reference deck has no content slides. Nothing to compare.")
        return {"recall": 0.0, "precision": 0.0}
    if not cand:
        print("Candidate deck has no content slides — full miss.")
        return {"recall": 0.0, "precision": 0.0}

    # For each reference slide, find the nearest candidate slide by pHash.
    cand_hashes = [(c[0], c[2]) for c in cand]
    matches = []
    for r_idx, r_label, r_h, _ in ref:
        best_d = 999
        best_c_idx = None
        for c_idx, c_h in cand_hashes:
            d = r_h - c_h
            if d < best_d:
                best_d = d
                best_c_idx = c_idx
        is_match = best_d <= match_threshold
        matches.append((r_idx, _short_label(r_label), best_c_idx, best_d, is_match))

    matched_ref = [m for m in matches if m[4]]
    missing = [m for m in matches if not m[4]]

    # Precision: candidate slides that DON'T match any reference (extras).
    used_cand_indices = {m[2] for m in matched_ref if m[2] is not None}
    extras = []
    for c_idx, c_label, c_h, _ in cand:
        if c_idx in used_cand_indices:
            continue
        # See if this candidate has a match in reference under threshold —
        # if yes, it's a duplicate match (already covered); if no, extra.
        best = min((r_h - c_h for _, _, r_h, _ in ref), default=999)
        if best > match_threshold:
            extras.append((c_idx, _short_label(c_label), best))

    recall = len(matched_ref) / len(ref)
    matched_cand = len(cand) - len(extras)
    precision = matched_cand / len(cand) if cand else 0.0

    # Detailed table — one line per reference slide.
    print(f"{'ref':>3s}  {'label':36s}  {'cand':>4s}  {'dist':>4s}  status")
    print("-" * 70)
    for r_idx, r_label, c_idx, d, is_match in matches:
        c_str = str(c_idx) if c_idx is not None else "-"
        flag = "✓" if is_match else "✗ MISSING"
        print(f"{r_idx:>3d}  {r_label:36s}  {c_str:>4s}  {d:>4d}  {flag}")

    if extras:
        print(f"\nExtras (candidate slides with no reference match within {match_threshold}):")
        for c_idx, c_label, d in extras:
            print(f"  candidate {c_idx}: {c_label}  (nearest ref dist={d})")

    print()
    print(f"Recall:    {len(matched_ref):>2d} / {len(ref)} = {recall:.1%}")
    print(f"Precision: {matched_cand:>2d} / {len(cand)} = {precision:.1%}")
    if missing:
        print(f"Missing:   {len(missing)} reference slide(s) not in candidate")

    return {
        "recall": recall, "precision": precision,
        "matches": matches, "extras": extras,
    }


def main() -> int:
    p = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument("reference", type=Path, help="Path to the curated reference .pptx")
    p.add_argument("candidate", type=Path, help="Path to the pipeline-generated .pptx")
    p.add_argument("--threshold", type=int, default=10,
                   help="pHash Hamming distance for considering two slides a match (default: 10)")
    p.add_argument("--min-recall", type=float, default=0.90,
                   help="Fail the script if recall is below this (default: 0.90)")
    args = p.parse_args()

    if not args.reference.exists():
        sys.exit(f"reference deck not found: {args.reference}")
    if not args.candidate.exists():
        sys.exit(f"candidate deck not found: {args.candidate}")

    result = compare(args.reference, args.candidate, match_threshold=args.threshold)
    if result["recall"] < args.min_recall:
        print(f"\nFAIL: recall {result['recall']:.1%} < required {args.min_recall:.1%}")
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
